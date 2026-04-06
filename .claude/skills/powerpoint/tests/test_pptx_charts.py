# Copyright (c) Microsoft Corporation.
# SPDX-License-Identifier: MIT
"""Tests for pptx_charts module."""

import pytest
from pptx.enum.chart import XL_CHART_TYPE
from pptx_charts import (
    CHART_TYPE_MAP,
    CHART_TYPE_REVERSE,
    add_chart_element,
    extract_chart,
)


def _make_chart(slide, elem, colors=None):
    """Helper to create a chart on a slide."""
    return add_chart_element(slide, elem, colors or {})


class TestAddChartElement:
    """Tests for add_chart_element."""

    def test_category_chart(self, blank_slide):
        elem = {
            "chart_type": "column_clustered",
            "left": 1.0,
            "top": 1.0,
            "width": 8.0,
            "height": 4.5,
            "categories": ["Q1", "Q2", "Q3"],
            "series": [
                {"name": "East", "values": [10, 20, 30]},
            ],
        }
        shape = _make_chart(blank_slide, elem)
        chart = shape.chart
        assert chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED
        plot = chart.plots[0]
        assert list(plot.categories) == ["Q1", "Q2", "Q3"]

    def test_chart_title(self, blank_slide):
        elem = {
            "chart_type": "column_clustered",
            "left": 1.0,
            "top": 1.0,
            "width": 8.0,
            "height": 4.5,
            "title": "Sales Report",
            "categories": ["Q1"],
            "series": [{"name": "S1", "values": [10]}],
        }
        shape = _make_chart(blank_slide, elem)
        assert shape.chart.has_title is True
        assert shape.chart.chart_title.text_frame.text == "Sales Report"

    def test_chart_legend(self, blank_slide):
        elem = {
            "chart_type": "column_clustered",
            "left": 1.0,
            "top": 1.0,
            "width": 8.0,
            "height": 4.5,
            "has_legend": True,
            "categories": ["Q1"],
            "series": [{"name": "S1", "values": [10]}],
        }
        shape = _make_chart(blank_slide, elem)
        assert shape.chart.has_legend is True

    def test_chart_style(self, blank_slide):
        elem = {
            "chart_type": "column_clustered",
            "left": 1.0,
            "top": 1.0,
            "width": 8.0,
            "height": 4.5,
            "chart_style": 10,
            "categories": ["Q1"],
            "series": [{"name": "S1", "values": [10]}],
        }
        shape = _make_chart(blank_slide, elem)
        assert shape.chart.style == 10

    def test_multiple_series(self, blank_slide):
        elem = {
            "chart_type": "column_clustered",
            "left": 1.0,
            "top": 1.0,
            "width": 8.0,
            "height": 4.5,
            "categories": ["Q1", "Q2"],
            "series": [
                {"name": "East", "values": [10, 20]},
                {"name": "West", "values": [15, 25]},
            ],
        }
        shape = _make_chart(blank_slide, elem)
        assert len(shape.chart.plots[0].series) == 2

    def test_series_coloring(self, blank_slide):
        elem = {
            "chart_type": "column_clustered",
            "left": 1.0,
            "top": 1.0,
            "width": 8.0,
            "height": 4.5,
            "categories": ["Q1"],
            "series": [
                {"name": "East", "values": [10], "color": "#0078D4"},
            ],
        }
        _make_chart(blank_slide, elem)
        # Verify no errors during creation

    def test_scatter_chart(self, blank_slide):
        elem = {
            "chart_type": "scatter",
            "left": 1.0,
            "top": 1.0,
            "width": 8.0,
            "height": 4.5,
            "series": [
                {
                    "name": "Series1",
                    "x_values": [1.0, 2.0, 3.0],
                    "y_values": [4.0, 5.0, 6.0],
                },
            ],
        }
        shape = _make_chart(blank_slide, elem)
        assert shape.chart.chart_type == XL_CHART_TYPE.XY_SCATTER

    def test_bubble_chart(self, blank_slide):
        elem = {
            "chart_type": "bubble",
            "left": 1.0,
            "top": 1.0,
            "width": 8.0,
            "height": 4.5,
            "series": [
                {
                    "name": "Bubbles",
                    "x_values": [1.0, 2.0],
                    "y_values": [3.0, 4.0],
                    "sizes": [10, 20],
                },
            ],
        }
        shape = _make_chart(blank_slide, elem)
        assert shape.chart.chart_type == XL_CHART_TYPE.BUBBLE

    def test_shape_name(self, blank_slide):
        elem = {
            "chart_type": "column_clustered",
            "left": 1.0,
            "top": 1.0,
            "width": 8.0,
            "height": 4.5,
            "name": "MyChart",
            "categories": ["Q1"],
            "series": [{"name": "S1", "values": [10]}],
        }
        shape = _make_chart(blank_slide, elem)
        assert shape.name == "MyChart"


class TestExtractChart:
    """Tests for extract_chart."""

    def test_basic_roundtrip(self, blank_slide):
        elem = {
            "chart_type": "column_clustered",
            "left": 1.0,
            "top": 1.0,
            "width": 8.0,
            "height": 4.5,
            "categories": ["Q1", "Q2", "Q3"],
            "series": [
                {"name": "East", "values": [10, 20, 30]},
            ],
        }
        shape = _make_chart(blank_slide, elem)
        result = extract_chart(shape)
        assert result["type"] == "chart"
        assert result["chart_type"] == "column_clustered"
        assert result["categories"] == ["Q1", "Q2", "Q3"]
        assert len(result["series"]) == 1
        assert result["series"][0]["name"] == "East"
        assert result["series"][0]["values"] == pytest.approx([10.0, 20.0, 30.0])

    def test_position_roundtrip(self, blank_slide):
        elem = {
            "chart_type": "column_clustered",
            "left": 2.0,
            "top": 3.0,
            "width": 6.0,
            "height": 4.0,
            "categories": ["A"],
            "series": [{"name": "S", "values": [1]}],
        }
        shape = _make_chart(blank_slide, elem)
        result = extract_chart(shape)
        assert result["left"] == pytest.approx(2.0, abs=0.01)
        assert result["top"] == pytest.approx(3.0, abs=0.01)

    def test_title_roundtrip(self, blank_slide):
        elem = {
            "chart_type": "column_clustered",
            "left": 1.0,
            "top": 1.0,
            "width": 8.0,
            "height": 4.5,
            "title": "Test Title",
            "categories": ["A"],
            "series": [{"name": "S", "values": [1]}],
        }
        shape = _make_chart(blank_slide, elem)
        result = extract_chart(shape)
        assert result["title"] == "Test Title"

    def test_legend_roundtrip(self, blank_slide):
        elem = {
            "chart_type": "column_clustered",
            "left": 1.0,
            "top": 1.0,
            "width": 8.0,
            "height": 4.5,
            "has_legend": True,
            "categories": ["A"],
            "series": [{"name": "S", "values": [1]}],
        }
        shape = _make_chart(blank_slide, elem)
        result = extract_chart(shape)
        assert result["has_legend"] is True

    def test_multiple_series_roundtrip(self, blank_slide):
        elem = {
            "chart_type": "column_clustered",
            "left": 1.0,
            "top": 1.0,
            "width": 8.0,
            "height": 4.5,
            "categories": ["Q1", "Q2"],
            "series": [
                {"name": "East", "values": [10, 20]},
                {"name": "West", "values": [15, 25]},
            ],
        }
        shape = _make_chart(blank_slide, elem)
        result = extract_chart(shape)
        assert len(result["series"]) == 2
        assert result["series"][0]["name"] == "East"
        assert result["series"][1]["name"] == "West"

    def test_name_roundtrip(self, blank_slide):
        elem = {
            "chart_type": "column_clustered",
            "left": 1.0,
            "top": 1.0,
            "width": 8.0,
            "height": 4.5,
            "name": "TestChart",
            "categories": ["A"],
            "series": [{"name": "S", "values": [1]}],
        }
        shape = _make_chart(blank_slide, elem)
        result = extract_chart(shape)
        assert result["name"] == "TestChart"


class TestChartConstants:
    """Tests for chart module constants."""

    def test_chart_type_map_keys(self):
        expected = {
            "column_clustered",
            "column_stacked",
            "bar_clustered",
            "bar_stacked",
            "line",
            "line_markers",
            "pie",
            "doughnut",
            "area",
            "radar",
            "scatter",
            "bubble",
        }
        assert set(CHART_TYPE_MAP.keys()) == expected

    def test_chart_type_reverse_consistency(self):
        for name, val in CHART_TYPE_MAP.items():
            assert CHART_TYPE_REVERSE[val] == name
