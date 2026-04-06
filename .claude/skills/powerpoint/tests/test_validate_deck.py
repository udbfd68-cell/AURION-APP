# Copyright (c) Microsoft Corporation.
# SPDX-License-Identifier: MIT
"""Tests for validate_deck module."""

import json

import pytest
from pptx import Presentation
from pptx.util import Inches
from validate_deck import (
    check_speaker_notes,
    create_parser,
    generate_report,
    main,
    max_severity,
    validate_deck,
)


@pytest.fixture()
def simple_deck(tmp_path):
    """Create a minimal PPTX with 2 slides and save it."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    prs.slides.add_slide(layout)
    prs.slides.add_slide(layout)
    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return path


@pytest.fixture()
def deck_with_notes(tmp_path):
    """PPTX with speaker notes on slide 1, empty on slide 2."""
    prs = Presentation()
    layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(layout)
    slide1.notes_slide.notes_text_frame.text = "Notes here"
    slide2 = prs.slides.add_slide(layout)
    slide2.notes_slide.notes_text_frame.text = ""
    path = tmp_path / "notes.pptx"
    prs.save(str(path))
    return path


class TestCheckSpeakerNotes:
    """Tests for check_speaker_notes."""

    def test_missing_notes(self, blank_slide):
        issues = check_speaker_notes(blank_slide, 1)
        assert len(issues) == 1
        assert issues[0]["severity"] == "warning"

    def test_present_notes(self, blank_slide):
        blank_slide.notes_slide.notes_text_frame.text = "Speaker notes"
        issues = check_speaker_notes(blank_slide, 1)
        assert len(issues) == 0

    def test_empty_notes(self, blank_slide):
        blank_slide.notes_slide.notes_text_frame.text = ""
        issues = check_speaker_notes(blank_slide, 1)
        assert len(issues) == 1
        assert issues[0]["severity"] == "info"
        assert "empty" in issues[0]["description"].lower()


class TestValidateDeck:
    """Tests for validate_deck."""

    def test_basic_validation(self, simple_deck):
        results = validate_deck(simple_deck)
        assert results["source"] == "pptx-properties"
        assert results["slide_count"] == 2
        assert len(results["slides"]) == 2

    def test_slide_filter(self, simple_deck):
        results = validate_deck(simple_deck, slide_filter={1})
        assert len(results["slides"]) == 1
        assert results["slides"][0]["slide_number"] == 1

    def test_content_dir_match(self, simple_deck, tmp_path):
        # Create matching content dirs
        for i in (1, 2):
            (tmp_path / f"slide-{i:03d}").mkdir()
        results = validate_deck(simple_deck, content_dir=tmp_path)
        assert "deck_issues" not in results

    def test_content_dir_mismatch(self, simple_deck, tmp_path):
        # Only 1 content dir for 2 slides — partial content info
        (tmp_path / "slide-001").mkdir()
        results = validate_deck(simple_deck, content_dir=tmp_path)
        assert "deck_issues" in results

    def test_with_notes(self, deck_with_notes):
        results = validate_deck(deck_with_notes)
        # Slide 1 has notes (good), slide 2 has empty notes (info)
        slide1 = results["slides"][0]
        slide2 = results["slides"][1]
        assert slide1["overall_quality"] == "good"
        assert slide2["overall_quality"] == "needs-attention"


class TestGenerateReport:
    """Tests for generate_report."""

    def test_report_contains_header(self, simple_deck):
        results = validate_deck(simple_deck)
        report = generate_report(results)
        assert "# PPTX Property Validation Report" in report

    def test_report_summary_table(self, simple_deck):
        results = validate_deck(simple_deck)
        report = generate_report(results)
        assert "Errors" in report
        assert "Warnings" in report

    def test_report_per_slide(self, simple_deck):
        results = validate_deck(simple_deck)
        report = generate_report(results)
        assert "Slide 1" in report
        assert "Slide 2" in report


class TestMaxSeverity:
    """Tests for max_severity."""

    def test_no_issues(self):
        results = {"slides": [{"issues": []}], "deck_issues": []}
        assert max_severity(results) == "none"

    def test_info_only(self):
        results = {
            "slides": [{"issues": [{"severity": "info"}]}],
        }
        assert max_severity(results) == "info"

    def test_warning(self):
        results = {
            "slides": [{"issues": [{"severity": "warning"}, {"severity": "info"}]}],
        }
        assert max_severity(results) == "warning"

    def test_error_highest(self):
        results = {
            "slides": [{"issues": [{"severity": "error"}, {"severity": "warning"}]}],
        }
        assert max_severity(results) == "error"

    def test_deck_issues_included(self):
        results = {
            "slides": [{"issues": []}],
            "deck_issues": [{"severity": "error"}],
        }
        assert max_severity(results) == "error"


class TestCreateParser:
    """Tests for create_parser."""

    def test_required_args(self):
        parser = create_parser()
        args = parser.parse_args(["--input", "test.pptx"])
        assert str(args.input) == "test.pptx"

    def test_optional_args(self):
        parser = create_parser()
        args = parser.parse_args(
            [
                "--input",
                "test.pptx",
                "--content-dir",
                "content/",
                "--slides",
                "1,3",
                "--output",
                "results.json",
                "--report",
                "report.md",
            ]
        )
        assert str(args.content_dir) == "content"
        assert args.slides == "1,3"


class TestValidateDeckExtended:
    """Extended validate_deck tests for uncovered branches."""

    def test_content_dir_more_than_slides(self, simple_deck, tmp_path):
        # 3 content dirs for 2 slides — mismatch warning
        for i in (1, 2, 3):
            (tmp_path / f"slide-{i:03d}").mkdir()
        results = validate_deck(simple_deck, content_dir=tmp_path)
        assert "deck_issues" in results
        issue = results["deck_issues"][0]
        assert issue["severity"] == "warning"
        assert "mismatch" in issue["description"].lower()

    def test_content_dir_skipped_with_filter(self, simple_deck, tmp_path):
        # With slide_filter, content_dir comparison is skipped
        (tmp_path / "slide-001").mkdir()
        results = validate_deck(simple_deck, content_dir=tmp_path, slide_filter={1})
        assert "deck_issues" not in results

    def test_notes_try_except(self, blank_slide):
        # Slides that don't have notes_slide accessed trigger the
        # "missing" path (no notes at all)
        issues = check_speaker_notes(blank_slide, 1)
        assert any(i["severity"] == "warning" for i in issues)


class TestGenerateReportExtended:
    """Extended generate_report tests for deck-level issues and per-slide."""

    def test_report_with_deck_issues(self, simple_deck, tmp_path):
        # More content dirs than slides → deck_issues
        for i in (1, 2, 3):
            (tmp_path / f"slide-{i:03d}").mkdir()
        results = validate_deck(simple_deck, content_dir=tmp_path)
        report = generate_report(results)
        assert "Deck-Level Findings" in report
        assert "slide_count" in report

    def test_report_slide_with_issues_table(self, deck_with_notes):
        results = validate_deck(deck_with_notes)
        report = generate_report(results)
        # Slide 2 has empty notes → should have issue table
        assert "Severity" in report

    def test_report_slide_no_issues(self):
        results = {
            "source": "pptx-properties",
            "slide_count": 1,
            "slides": [{"slide_number": 1, "issues": [], "overall_quality": "good"}],
        }
        report = generate_report(results)
        assert "No issues found" in report


class TestMain:
    """Tests for main() entry point."""

    def test_file_not_found(self, tmp_path, mocker):
        mocker.patch(
            "sys.argv",
            [
                "validate_deck.py",
                "--input",
                str(tmp_path / "nonexistent.pptx"),
            ],
        )
        result = main()
        assert result != 0

    def test_basic_run(self, simple_deck, tmp_path, mocker):
        out_json = tmp_path / "results.json"
        report_md = tmp_path / "report.md"
        mocker.patch(
            "sys.argv",
            [
                "validate_deck.py",
                "--input",
                str(simple_deck),
                "--output",
                str(out_json),
                "--report",
                str(report_md),
            ],
        )
        result = main()
        # Slides have missing notes (warning severity) → EXIT_FAILURE
        assert result in (0, 1)
        assert out_json.exists()
        assert report_md.exists()
        data = json.loads(out_json.read_text())
        assert data["slide_count"] == 2

    def test_with_slide_filter(self, simple_deck, mocker):
        mocker.patch(
            "sys.argv",
            [
                "validate_deck.py",
                "--input",
                str(simple_deck),
                "--slides",
                "1",
            ],
        )
        result = main()
        # Missing notes → warning → EXIT_FAILURE
        assert result in (0, 1)

    def test_warnings_exit_failure(self, deck_with_notes, tmp_path, mocker):
        mocker.patch(
            "sys.argv",
            [
                "validate_deck.py",
                "--input",
                str(deck_with_notes),
            ],
        )
        result = main()
        # Slide 2 has empty notes (info severity) → exit success
        # unless it has warnings
        assert result in (0, 1)

    def test_per_slide_dir_output(self, simple_deck, tmp_path, mocker):
        per_slide_dir = tmp_path / "per-slide"
        mocker.patch(
            "sys.argv",
            [
                "validate_deck.py",
                "--input",
                str(simple_deck),
                "--per-slide-dir",
                str(per_slide_dir),
            ],
        )
        result = main()
        assert result in (0, 1)
        assert per_slide_dir.exists()
        slide_1 = per_slide_dir / "slide-001-deck-validation.json"
        slide_2 = per_slide_dir / "slide-002-deck-validation.json"
        assert slide_1.exists()
        assert slide_2.exists()
        data = json.loads(slide_1.read_text())
        assert data["slide_number"] == 1
        assert "issues" in data
        assert "overall_quality" in data

    def test_per_slide_dir_with_filter(self, simple_deck, tmp_path, mocker):
        per_slide_dir = tmp_path / "per-slide-filtered"
        mocker.patch(
            "sys.argv",
            [
                "validate_deck.py",
                "--input",
                str(simple_deck),
                "--slides",
                "1",
                "--per-slide-dir",
                str(per_slide_dir),
            ],
        )
        result = main()
        assert result in (0, 1)
        assert (per_slide_dir / "slide-001-deck-validation.json").exists()
        assert not (per_slide_dir / "slide-002-deck-validation.json").exists()


class TestCreateParserPerSlideDir:
    """Tests for --per-slide-dir argument."""

    def test_per_slide_dir_arg(self):
        parser = create_parser()
        args = parser.parse_args(["--input", "test.pptx", "--per-slide-dir", "output/"])
        assert str(args.per_slide_dir) == "output"
