# Copyright (c) Microsoft Corporation.
# SPDX-License-Identifier: MIT
"""Tests for pptx_colors module."""

import pytest
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx_colors import (
    THEME_COLOR_MAP,
    apply_color_spec,
    apply_color_to_fill,
    apply_color_to_font,
    extract_color,
    hex_brightness,
    resolve_color,
    rgb_to_hex,
)


class TestResolveColor:
    """Tests for resolve_color."""

    def test_hex_value(self):
        result = resolve_color("#0078D4")
        assert "rgb" in result
        assert result["rgb"] == RGBColor(0x00, 0x78, 0xD4)

    def test_hex_no_hash(self):
        result = resolve_color("0078D4")
        assert "rgb" in result
        assert result["rgb"] == RGBColor(0x00, 0x78, 0xD4)

    def test_theme_reference(self):
        result = resolve_color("@accent_1")
        assert "theme" in result
        assert result["theme"] == MSO_THEME_COLOR.ACCENT_1

    def test_theme_dark_1(self):
        result = resolve_color("@dark_1")
        assert result["theme"] == MSO_THEME_COLOR.DARK_1

    def test_invalid_theme(self):
        result = resolve_color("@nonexistent")
        assert "rgb" in result
        assert result["rgb"] == RGBColor(0, 0, 0)

    def test_dict_with_theme(self):
        result = resolve_color({"theme": "accent_2"})
        assert result["theme"] == MSO_THEME_COLOR.ACCENT_2

    def test_dict_with_theme_and_brightness(self):
        result = resolve_color({"theme": "accent_1", "brightness": 0.5})
        assert result["theme"] == MSO_THEME_COLOR.ACCENT_1
        assert result["brightness"] == 0.5

    def test_dict_with_invalid_theme_falls_back_to_color_key(self):
        result = resolve_color({"theme": "invalid", "color": "#FF0000"})
        assert "rgb" in result
        assert result["rgb"] == RGBColor(0xFF, 0x00, 0x00)

    def test_dict_with_invalid_theme_no_color(self):
        result = resolve_color({"theme": "invalid"})
        assert "rgb" in result
        assert result["rgb"] == RGBColor(0, 0, 0)

    def test_non_string_input(self):
        result = resolve_color(42)
        assert "rgb" in result
        assert result["rgb"] == RGBColor(0, 0, 0)

    def test_white(self):
        result = resolve_color("#FFFFFF")
        assert result["rgb"] == RGBColor(0xFF, 0xFF, 0xFF)

    def test_black(self):
        result = resolve_color("#000000")
        assert result["rgb"] == RGBColor(0, 0, 0)

    def test_depth_limit_raises_on_nested_dicts(self):
        """Deeply nested dict color values exceed the depth limit."""
        nested = {"color": {"color": {"color": "#FF0000"}}}
        with pytest.raises(ValueError, match="exceeds limit"):
            resolve_color(nested, _depth=0, max_depth=2)

    def test_depth_limit_allows_single_dict_unwrap(self):
        """A single dict unwrap stays within the default depth limit."""
        result = resolve_color({"color": "#0078D4"})
        assert "rgb" in result
        assert result["rgb"] == RGBColor(0x00, 0x78, 0xD4)


class TestApplyColorSpec:
    """Tests for apply_color_spec with real python-pptx objects."""

    def test_apply_rgb(self, sample_shape):
        spec = {"rgb": RGBColor(0xFF, 0x00, 0x00)}
        apply_color_spec(sample_shape.fill.fore_color, spec)
        assert sample_shape.fill.fore_color.rgb == RGBColor(0xFF, 0x00, 0x00)

    def test_apply_theme(self, sample_shape):
        spec = {"theme": MSO_THEME_COLOR.ACCENT_1}
        apply_color_spec(sample_shape.fill.fore_color, spec)
        assert sample_shape.fill.fore_color.theme_color == MSO_THEME_COLOR.ACCENT_1

    def test_apply_theme_with_brightness(self, sample_shape):
        spec = {"theme": MSO_THEME_COLOR.ACCENT_1, "brightness": 0.4}
        apply_color_spec(sample_shape.fill.fore_color, spec)
        assert sample_shape.fill.fore_color.theme_color == MSO_THEME_COLOR.ACCENT_1
        assert sample_shape.fill.fore_color.brightness == pytest.approx(0.4, abs=0.01)


class TestApplyColorToFill:
    """Tests for apply_color_to_fill."""

    def test_delegates_to_fill_fore_color(self, sample_shape):
        spec = {"rgb": RGBColor(0x00, 0xFF, 0x00)}
        sample_shape.fill.solid()
        apply_color_to_fill(sample_shape.fill, spec)
        assert sample_shape.fill.fore_color.rgb == RGBColor(0x00, 0xFF, 0x00)


class TestApplyColorToFont:
    """Tests for apply_color_to_font."""

    def test_delegates_to_font_color(self, sample_textbox):
        spec = {"rgb": RGBColor(0xFF, 0x00, 0x00)}
        run = sample_textbox.text_frame.paragraphs[0].runs[0]
        apply_color_to_font(run.font.color, spec)
        assert run.font.color.rgb == RGBColor(0xFF, 0x00, 0x00)


class TestExtractColor:
    """Tests for extract_color."""

    def test_rgb_color(self, sample_shape):
        sample_shape.fill.solid()
        sample_shape.fill.fore_color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
        result = extract_color(sample_shape.fill.fore_color)
        assert result == "#AABBCC"

    def test_theme_color(self, sample_shape):
        sample_shape.fill.solid()
        sample_shape.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
        result = extract_color(sample_shape.fill.fore_color)
        assert result == "@accent_1"

    def test_none_type(self, blank_slide):
        from pptx.enum.shapes import MSO_SHAPE

        shape = blank_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 100, 100)
        # A fresh shape may have no explicit fill color type
        result = extract_color(shape.line.color)
        # Result should be None or a valid color string
        assert result is None or isinstance(result, str)


class TestRgbToHex:
    """Tests for rgb_to_hex."""

    @pytest.mark.parametrize(
        "rgb,expected",
        [
            (RGBColor(0x00, 0x78, 0xD4), "#0078D4"),
            (RGBColor(0xFF, 0xFF, 0xFF), "#FFFFFF"),
            (RGBColor(0, 0, 0), "#000000"),
            (None, None),
        ],
    )
    def test_conversion(self, rgb, expected):
        assert rgb_to_hex(rgb) == expected


class TestHexBrightness:
    """Tests for hex_brightness."""

    @pytest.mark.parametrize(
        "hex_val,expected",
        [
            ("#FFFFFF", 255),
            ("#000000", 0),
            ("#FF0000", 76),
            ("#00FF00", 149),
            ("#0000FF", 29),
        ],
    )
    def test_brightness(self, hex_val, expected):
        assert hex_brightness(hex_val) == expected

    def test_mid_gray(self):
        result = hex_brightness("#808080")
        assert 127 <= result <= 129


class TestThemeColorMap:
    """Tests for THEME_COLOR_MAP constant."""

    def test_all_theme_colors_present(self):
        expected_keys = {
            "accent_1",
            "accent_2",
            "accent_3",
            "accent_4",
            "accent_5",
            "accent_6",
            "dark_1",
            "dark_2",
            "light_1",
            "light_2",
            "text_1",
            "text_2",
            "background_1",
            "background_2",
            "hyperlink",
            "followed_hyperlink",
        }
        assert set(THEME_COLOR_MAP.keys()) == expected_keys
