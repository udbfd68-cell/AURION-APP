# Copyright (c) Microsoft Corporation.
# SPDX-License-Identifier: MIT
"""Tests for pptx_text module."""

import pytest
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt
from pptx_text import (
    AUTO_SIZE_MAP,
    AUTO_SIZE_REVERSE,
    SHAPE_KEYS,
    TEXTBOX_KEYS,
    VERTICAL_ANCHOR_MAP,
    VERTICAL_ANCHOR_REVERSE,
    _apply_char_spacing,
    _apply_run_effect,
    _extract_run_effect,
    apply_bullet_properties,
    apply_paragraph_properties,
    apply_run_properties,
    apply_text_properties,
    extract_bullet_properties,
    extract_paragraph_properties,
    extract_run_properties,
    extract_text_frame_properties,
    populate_text_frame,
    split_lines,
)

# ----- Helpers -----


def _make_textbox(slide, text="Test"):
    """Create a textbox on a slide with given text."""
    txBox = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(4), Inches(2))
    txBox.text_frame.text = text
    return txBox


# ----- apply_text_properties -----


class TestApplyTextProperties:
    """Tests for apply_text_properties."""

    @pytest.mark.parametrize(
        "prop,value",
        [
            ("margin_left", 0.5),
            ("margin_right", 0.25),
            ("margin_top", 0.1),
            ("margin_bottom", 0.2),
        ],
    )
    def test_margin(self, blank_slide, prop, value):
        txBox = _make_textbox(blank_slide)
        apply_text_properties(txBox.text_frame, {prop: value})
        assert getattr(txBox.text_frame, prop) == Inches(value)

    @pytest.mark.parametrize(
        "key,auto_size_enum",
        [
            ("none", MSO_AUTO_SIZE.NONE),
            ("fit", MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT),
            ("shrink", MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE),
        ],
    )
    def test_auto_size(self, blank_slide, key, auto_size_enum):
        txBox = _make_textbox(blank_slide)
        apply_text_properties(txBox.text_frame, {"auto_size": key})
        assert txBox.text_frame.auto_size == auto_size_enum

    def test_vertical_anchor_middle(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        apply_text_properties(txBox.text_frame, {"vertical_anchor": "middle"})
        assert txBox.text_frame.vertical_anchor == MSO_VERTICAL_ANCHOR.MIDDLE

    def test_word_wrap(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        apply_text_properties(txBox.text_frame, {"word_wrap": True})
        assert txBox.text_frame.word_wrap is True

    def test_missing_keys_no_change(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        apply_text_properties(txBox.text_frame, {})
        # No error raised, state unchanged


# ----- extract_text_frame_properties -----


class TestExtractTextFrameProperties:
    """Tests for extract_text_frame_properties."""

    def test_roundtrip_margins(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        apply_text_properties(
            txBox.text_frame,
            {
                "margin_left": 0.5,
                "margin_right": 0.25,
                "margin_top": 0.1,
                "margin_bottom": 0.2,
            },
        )
        props = extract_text_frame_properties(txBox.text_frame)
        assert props["margin_left"] == pytest.approx(0.5, abs=0.001)
        assert props["margin_right"] == pytest.approx(0.25, abs=0.001)
        assert props["margin_top"] == pytest.approx(0.1, abs=0.001)
        assert props["margin_bottom"] == pytest.approx(0.2, abs=0.001)

    def test_roundtrip_auto_size(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        apply_text_properties(txBox.text_frame, {"auto_size": "shrink"})
        props = extract_text_frame_properties(txBox.text_frame)
        assert props["auto_size"] == "shrink"

    def test_roundtrip_vertical_anchor(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        apply_text_properties(txBox.text_frame, {"vertical_anchor": "bottom"})
        props = extract_text_frame_properties(txBox.text_frame)
        assert props["vertical_anchor"] == "bottom"

    def test_empty_when_no_props(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        # Textbox has default margins set by python-pptx
        props = extract_text_frame_properties(txBox.text_frame)
        assert isinstance(props, dict)


# ----- apply_paragraph_properties -----


class TestApplyParagraphProperties:
    """Tests for apply_paragraph_properties."""

    def test_space_before(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        para = txBox.text_frame.paragraphs[0]
        apply_paragraph_properties(para, {"space_before": 12})
        assert para.space_before == Pt(12)

    def test_space_after(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        para = txBox.text_frame.paragraphs[0]
        apply_paragraph_properties(para, {"space_after": 6})
        assert para.space_after == Pt(6)

    def test_line_spacing_factor(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        para = txBox.text_frame.paragraphs[0]
        apply_paragraph_properties(para, {"line_spacing": 1.5})
        assert para.line_spacing == pytest.approx(1.5)

    def test_line_spacing_pt(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        para = txBox.text_frame.paragraphs[0]
        apply_paragraph_properties(para, {"line_spacing": 18})
        assert para.line_spacing == Pt(18)

    def test_level(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        para = txBox.text_frame.paragraphs[0]
        apply_paragraph_properties(para, {"level": 2})
        assert para.level == 2

    def test_missing_keys(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        para = txBox.text_frame.paragraphs[0]
        apply_paragraph_properties(para, {})
        # No error


# ----- extract_paragraph_properties -----


class TestExtractParagraphProperties:
    """Tests for extract_paragraph_properties."""

    def test_roundtrip_spacing(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        para = txBox.text_frame.paragraphs[0]
        apply_paragraph_properties(para, {"space_before": 12, "space_after": 6})
        props = extract_paragraph_properties(para)
        assert props["space_before"] == pytest.approx(12.0, abs=0.1)
        assert props["space_after"] == pytest.approx(6.0, abs=0.1)

    def test_roundtrip_line_spacing_factor(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        para = txBox.text_frame.paragraphs[0]
        apply_paragraph_properties(para, {"line_spacing": 1.5})
        props = extract_paragraph_properties(para)
        assert props["line_spacing"] == pytest.approx(1.5)

    def test_level_zero_omitted(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        para = txBox.text_frame.paragraphs[0]
        para.level = 0
        props = extract_paragraph_properties(para)
        assert "level" not in props

    def test_level_nonzero(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        para = txBox.text_frame.paragraphs[0]
        apply_paragraph_properties(para, {"level": 3})
        props = extract_paragraph_properties(para)
        assert props["level"] == 3

    def test_empty_when_no_props(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        para = txBox.text_frame.paragraphs[0]
        props = extract_paragraph_properties(para)
        assert isinstance(props, dict)


# ----- apply_run_properties / extract_run_properties -----


class TestRunProperties:
    """Tests for apply_run_properties and extract_run_properties."""

    def test_underline_roundtrip(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        run = txBox.text_frame.paragraphs[0].runs[0]
        apply_run_properties(run, {"underline": True}, {})
        props = extract_run_properties(run)
        assert props["underline"] is True

    def test_char_spacing_roundtrip(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        run = txBox.text_frame.paragraphs[0].runs[0]
        apply_run_properties(run, {"char_spacing": 1.5}, {})
        props = extract_run_properties(run)
        assert props["char_spacing"] == pytest.approx(1.5)

    def test_hyperlink(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        run = txBox.text_frame.paragraphs[0].runs[0]
        apply_run_properties(run, {"hyperlink": "https://example.com"}, {})
        props = extract_run_properties(run)
        assert props["hyperlink"] == "https://example.com"

    def test_effect_outer_shadow(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        run = txBox.text_frame.paragraphs[0].runs[0]
        effect = {
            "type": "outer_shadow",
            "blurRad": "40000",
            "dist": "20000",
            "dir": "5400000",
            "color": "black",
            "color_type": "preset",
        }
        apply_run_properties(run, {"effect": effect}, {})
        props = extract_run_properties(run)
        assert props["effect"]["type"] == "outer_shadow"
        assert props["effect"]["blurRad"] == "40000"
        assert props["effect"]["color"] == "black"
        assert props["effect"]["color_type"] == "preset"

    def test_empty_when_no_props(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        run = txBox.text_frame.paragraphs[0].runs[0]
        props = extract_run_properties(run)
        assert isinstance(props, dict)


# ----- _apply_run_effect / _extract_run_effect -----


class TestRunEffect:
    """Tests for _apply_run_effect and _extract_run_effect."""

    def test_preset_color(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        run = txBox.text_frame.paragraphs[0].runs[0]
        effect = {
            "type": "outer_shadow",
            "blurRad": "50800",
            "dist": "38100",
            "dir": "2700000",
            "color": "black",
            "color_type": "preset",
        }
        _apply_run_effect(run, effect)
        result = _extract_run_effect(run)
        assert result["type"] == "outer_shadow"
        assert result["color"] == "black"
        assert result["color_type"] == "preset"

    def test_rgb_color(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        run = txBox.text_frame.paragraphs[0].runs[0]
        effect = {
            "type": "outer_shadow",
            "blurRad": "40000",
            "color": "#FF0000",
            "color_type": "rgb",
        }
        _apply_run_effect(run, effect)
        result = _extract_run_effect(run)
        assert result["color"] == "#FF0000"
        assert result["color_type"] == "rgb"

    def test_alpha(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        run = txBox.text_frame.paragraphs[0].runs[0]
        effect = {
            "type": "outer_shadow",
            "blurRad": "40000",
            "color": "black",
            "color_type": "preset",
            "alpha": 50.0,
        }
        _apply_run_effect(run, effect)
        result = _extract_run_effect(run)
        assert result["alpha"] == pytest.approx(50.0, abs=0.1)

    def test_non_shadow_type_noop(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        run = txBox.text_frame.paragraphs[0].runs[0]
        _apply_run_effect(run, {"type": "glow"})
        assert _extract_run_effect(run) is None

    def test_none_effect_noop(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        run = txBox.text_frame.paragraphs[0].runs[0]
        _apply_run_effect(run, None)
        assert _extract_run_effect(run) is None


# ----- _apply_char_spacing -----


class TestApplyCharSpacing:
    """Tests for _apply_char_spacing."""

    @pytest.mark.parametrize(
        "value,expected_spc",
        [
            (2.0, "200"),
            (-1.0, "-100"),
            (0.0, "0"),
        ],
    )
    def test_char_spacing(self, blank_slide, value, expected_spc):
        txBox = _make_textbox(blank_slide)
        run = txBox.text_frame.paragraphs[0].runs[0]
        _apply_char_spacing(run.font, value)
        assert run.font._element.get("spc") == expected_spc


# ----- bullet_properties -----


class TestBulletProperties:
    """Tests for extract_bullet_properties and apply_bullet_properties."""

    def test_bullet_char_roundtrip(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        para = txBox.text_frame.paragraphs[0]
        apply_bullet_properties(para, {"bullet_char": "•"})
        props = extract_bullet_properties(para)
        assert props["bullet_char"] == "•"

    def test_bullet_none_roundtrip(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        para = txBox.text_frame.paragraphs[0]
        apply_bullet_properties(para, {"bullet_none": True})
        props = extract_bullet_properties(para)
        assert props["bullet_none"] is True

    def test_bullet_font(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        para = txBox.text_frame.paragraphs[0]
        apply_bullet_properties(
            para,
            {
                "bullet_char": "→",
                "bullet_font": "Wingdings",
            },
        )
        props = extract_bullet_properties(para)
        assert props["bullet_font"] == "Wingdings"

    def test_bullet_size_pct(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        para = txBox.text_frame.paragraphs[0]
        apply_bullet_properties(
            para,
            {
                "bullet_char": "•",
                "bullet_size_pct": 75000,
            },
        )
        props = extract_bullet_properties(para)
        assert props["bullet_size_pct"] == 75000

    def test_bullet_color(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        para = txBox.text_frame.paragraphs[0]
        apply_bullet_properties(
            para,
            {
                "bullet_char": "•",
                "bullet_color": "#FF0000",
            },
        )
        props = extract_bullet_properties(para)
        assert props["bullet_color"] == "#FF0000"

    def test_bullet_margin_indent(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        para = txBox.text_frame.paragraphs[0]
        apply_bullet_properties(
            para,
            {
                "bullet_margin_left": 228600,
                "bullet_indent": -228600,
            },
        )
        props = extract_bullet_properties(para)
        assert props["bullet_margin_left"] == 228600
        assert props["bullet_indent"] == -228600

    def test_no_bullet_props_noop(self, blank_slide):
        txBox = _make_textbox(blank_slide)
        para = txBox.text_frame.paragraphs[0]
        apply_bullet_properties(para, {"font": "Arial"})
        # No pPr bullet elements should be added
        props = extract_bullet_properties(para)
        assert "bullet_char" not in props

    def test_empty_paragraph_no_pPr(self, blank_slide):
        txBox = _make_textbox(blank_slide, text="")
        para = txBox.text_frame.paragraphs[0]
        props = extract_bullet_properties(para)
        assert isinstance(props, dict)


# ----- Constants -----


class TestTextConstants:
    """Tests for module constants."""

    def test_auto_size_map(self):
        assert set(AUTO_SIZE_MAP.keys()) == {"none", "fit", "shrink"}

    def test_auto_size_reverse(self):
        assert set(AUTO_SIZE_REVERSE.values()) == {"none", "fit", "shrink"}

    def test_vertical_anchor_map(self):
        assert set(VERTICAL_ANCHOR_MAP.keys()) == {"top", "middle", "bottom"}

    def test_vertical_anchor_reverse(self):
        assert set(VERTICAL_ANCHOR_REVERSE.values()) == {"top", "middle", "bottom"}


# ----- split_lines -----


class TestSplitLines:
    """Tests for split_lines helper."""

    @pytest.mark.parametrize(
        "text,expected",
        [
            ("hello", ["hello"]),
            ("line1\nline2", ["line1", "line2"]),
            ("line1\vline2", ["line1", "line2"]),
            ("a\nb\nc", ["a", "b", "c"]),
            ("", [""]),
            ("no breaks", ["no breaks"]),
        ],
    )
    def test_split(self, text, expected):
        assert split_lines(text) == expected


# ----- populate_text_frame -----


class TestPopulateTextFrame:
    """Tests for populate_text_frame shared text-frame population."""

    def test_flat_text(self, blank_slide):
        txBox = _make_textbox(blank_slide, text="")
        populate_text_frame(
            txBox.text_frame,
            {"text": "Hello World"},
            {},
            TEXTBOX_KEYS,
        )
        assert txBox.text_frame.text == "Hello World"

    def test_multiline_text(self, blank_slide):
        txBox = _make_textbox(blank_slide, text="")
        populate_text_frame(
            txBox.text_frame,
            {"text": "Line 1\nLine 2\nLine 3"},
            {},
            TEXTBOX_KEYS,
        )
        paras = txBox.text_frame.paragraphs
        assert len(paras) == 3
        assert paras[0].text == "Line 1"
        assert paras[2].text == "Line 3"

    def test_paragraphs_with_text(self, blank_slide):
        txBox = _make_textbox(blank_slide, text="")
        elem = {
            "paragraphs": [
                {"text": "First paragraph"},
                {"text": "Second paragraph"},
            ],
        }
        populate_text_frame(txBox.text_frame, elem, {}, TEXTBOX_KEYS)
        paras = txBox.text_frame.paragraphs
        assert len(paras) == 2
        assert paras[0].text == "First paragraph"
        assert paras[1].text == "Second paragraph"

    def test_no_text_noop(self, blank_slide):
        txBox = _make_textbox(blank_slide, text="Original")
        populate_text_frame(txBox.text_frame, {}, {}, TEXTBOX_KEYS)
        # Word wrap is set but text frame is unchanged
        assert txBox.text_frame.word_wrap is True

    def test_textbox_keys_apply_font(self, blank_slide):
        txBox = _make_textbox(blank_slide, text="")
        defaults = {"font": "Arial", "size": 20}
        populate_text_frame(
            txBox.text_frame,
            {"text": "Styled"},
            {},
            TEXTBOX_KEYS,
            defaults,
        )
        run = txBox.text_frame.paragraphs[0].runs[0]
        assert run.font.name == "Arial"
        assert run.font.size == Pt(20)

    def test_shape_keys_apply_font(self, blank_slide):
        from pptx.enum.shapes import MSO_SHAPE

        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Inches(4), Inches(2)
        )
        populate_text_frame(
            shape.text_frame,
            {"text": "Shape text", "text_font": "Calibri"},
            {},
            SHAPE_KEYS,
        )
        run = shape.text_frame.paragraphs[0].runs[0]
        assert run.font.name == "Calibri"

    def test_paragraphs_with_runs(self, blank_slide):
        txBox = _make_textbox(blank_slide, text="")
        elem = {
            "paragraphs": [
                {
                    "runs": [
                        {"text": "bold ", "bold": True},
                        {"text": "normal"},
                    ],
                },
            ],
        }
        populate_text_frame(txBox.text_frame, elem, {}, TEXTBOX_KEYS)
        runs = txBox.text_frame.paragraphs[0].runs
        assert len(runs) == 2
        assert runs[0].text == "bold "
        assert runs[0].font.bold is True
        assert runs[1].text == "normal"
