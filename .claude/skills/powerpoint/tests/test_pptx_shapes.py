# Copyright (c) Microsoft Corporation.
# SPDX-License-Identifier: MIT
"""Tests for pptx_shapes module."""

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx_shapes import (
    AUTO_SHAPE_NAME_MAP,
    SHAPE_MAP,
    apply_rotation,
    extract_rotation,
)


class TestShapeMap:
    """Tests for SHAPE_MAP constant."""

    def test_has_basic_shapes(self):
        expected = {"rectangle", "rounded_rectangle", "oval", "diamond", "circle"}
        assert expected.issubset(set(SHAPE_MAP.keys()))

    def test_has_arrow_shapes(self):
        expected = {"right_arrow", "left_arrow", "up_arrow", "down_arrow"}
        assert expected.issubset(set(SHAPE_MAP.keys()))

    def test_has_flowchart_shapes(self):
        expected = {"flowchart_process", "flowchart_decision", "flowchart_terminator"}
        assert expected.issubset(set(SHAPE_MAP.keys()))

    def test_circle_is_oval_alias(self):
        assert SHAPE_MAP["circle"] == SHAPE_MAP["oval"]

    def test_all_values_are_mso_shape(self):
        for name, value in SHAPE_MAP.items():
            assert isinstance(value, int) or hasattr(value, "real"), (
                f"SHAPE_MAP[{name!r}] is not a valid MSO_SHAPE value"
            )


class TestAutoShapeNameMap:
    """Tests for AUTO_SHAPE_NAME_MAP inverse constant."""

    def test_excludes_circle(self):
        assert "circle" not in AUTO_SHAPE_NAME_MAP.values()

    def test_oval_is_canonical(self):
        assert AUTO_SHAPE_NAME_MAP[MSO_SHAPE.OVAL] == "oval"

    def test_rectangle_present(self):
        assert AUTO_SHAPE_NAME_MAP[MSO_SHAPE.RECTANGLE] == "rectangle"

    def test_inverse_consistency(self):
        for name, mso_val in SHAPE_MAP.items():
            if name == "circle":
                continue
            assert AUTO_SHAPE_NAME_MAP[mso_val] == name


class TestApplyRotation:
    """Tests for apply_rotation."""

    def test_apply_nonzero(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
        )
        apply_rotation(shape, 45.0)
        assert shape.rotation == 45.0

    def test_apply_none(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
        )
        original = shape.rotation
        apply_rotation(shape, None)
        assert shape.rotation == original

    def test_apply_zero(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
        )
        shape.rotation = 90.0
        apply_rotation(shape, 0)
        # 0 is treated as no-op per the guard condition
        assert shape.rotation == 90.0


class TestExtractRotation:
    """Tests for extract_rotation."""

    def test_nonzero_rotation(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
        )
        shape.rotation = 90.0
        assert extract_rotation(shape) == 90.0

    def test_zero_rotation(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
        )
        shape.rotation = 0.0
        assert extract_rotation(shape) is None

    def test_default_rotation(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
        )
        assert extract_rotation(shape) is None
