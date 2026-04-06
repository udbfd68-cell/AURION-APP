# Copyright (c) Microsoft Corporation.
# SPDX-License-Identifier: MIT
"""Tests for build_deck module."""

from unittest.mock import MagicMock

import pytest
from build_deck import (
    ContentExtraError,
    _reset_effect_ref,
    _validate_content_extra,
    add_arrow_flow_element,
    add_card_element,
    add_connector_element,
    add_group_element,
    add_image_element,
    add_numbered_step_element,
    add_rich_text_element,
    add_shape_element,
    add_textbox,
    build_element_in_group,
    build_slide,
    clear_slide_shapes,
    discover_slides,
    get_slide_layout,
    main,
    set_slide_bg,
    set_slide_bg_image,
)
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


class TestResetEffectRef:
    """Tests for _reset_effect_ref."""

    def test_resets_idx_to_zero(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(2),
            Inches(2),
        )
        # python-pptx defaults effectRef idx to 2
        _reset_effect_ref(shape)
        pns = "http://schemas.openxmlformats.org/presentationml/2006/main"
        ans = "http://schemas.openxmlformats.org/drawingml/2006/main"
        style_el = shape._element.find(f"{{{pns}}}style")
        if style_el is not None:
            effect_ref = style_el.find(f"{{{ans}}}effectRef")
            if effect_ref is not None:
                assert effect_ref.get("idx") == "0"

    def test_no_style_element(self, blank_slide):
        """Shape without style element doesn't raise."""
        txBox = blank_slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(2),
            Inches(1),
        )
        # Should not raise
        _reset_effect_ref(txBox)


class TestSetSlideBg:
    """Tests for set_slide_bg."""

    def test_solid_fill(self, blank_slide):
        set_slide_bg(blank_slide, "#0078D4", {})
        # Verify background fill was set
        bg = blank_slide.background
        assert bg.fill is not None

    def test_gradient_fill(self, blank_slide):
        grad_spec = {
            "type": "gradient",
            "angle": 90,
            "stops": [
                {"position": 0.0, "color": "#000000"},
                {"position": 1.0, "color": "#FFFFFF"},
            ],
        }
        set_slide_bg(blank_slide, grad_spec, {})


class TestAddTextbox:
    """Tests for add_textbox."""

    def test_basic_textbox(self, blank_slide):
        txBox = add_textbox(blank_slide, 1, 1, 4, 1, "Hello World")
        assert txBox is not None
        assert txBox.text_frame.text == "Hello World"

    def test_textbox_font(self, blank_slide):
        txBox = add_textbox(
            blank_slide,
            1,
            1,
            4,
            1,
            "Styled",
            font_name="Arial",
            font_size=24,
            bold=True,
        )
        runs = txBox.text_frame.paragraphs[0].runs
        assert len(runs) >= 1
        assert runs[0].font.size == Pt(24)
        assert runs[0].font.bold is True

    def test_textbox_multiline(self, blank_slide):
        txBox = add_textbox(blank_slide, 1, 1, 4, 2, "Line1\nLine2\nLine3")
        paras = txBox.text_frame.paragraphs
        assert len(paras) == 3

    def test_textbox_with_name(self, blank_slide):
        txBox = add_textbox(blank_slide, 1, 1, 4, 1, "Named", name="MyBox")
        assert txBox.name == "MyBox"

    def test_textbox_paragraphs_format(self, blank_slide):
        """Per-paragraph formatting via elem dict."""
        elem = {
            "paragraphs": [
                {"text": "Bold para", "font_bold": True, "font_size": 20},
                {"text": "Normal para", "font_size": 14},
            ],
        }
        txBox = add_textbox(
            blank_slide,
            1,
            1,
            4,
            2,
            "",
            elem=elem,
            colors={},
        )
        paras = txBox.text_frame.paragraphs
        assert len(paras) == 2

    def test_textbox_rich_text_runs(self, blank_slide):
        """Per-paragraph runs with mixed formatting."""
        elem = {
            "paragraphs": [
                {
                    "text": "Mixed",
                    "runs": [
                        {"text": "Bold", "bold": True, "size": 16},
                        {"text": " Italic", "italic": True, "size": 16},
                    ],
                },
            ],
        }
        txBox = add_textbox(
            blank_slide,
            1,
            1,
            4,
            1,
            "",
            elem=elem,
            colors={},
        )
        runs = txBox.text_frame.paragraphs[0].runs
        assert len(runs) == 2

    def test_textbox_vertical_tab_split(self, blank_slide):
        """Vertical tab treated as line break."""
        txBox = add_textbox(blank_slide, 1, 1, 4, 2, "Line1\vLine2")
        paras = txBox.text_frame.paragraphs
        assert len(paras) == 2


class TestAddShapeElement:
    """Tests for add_shape_element."""

    def test_basic_rectangle(self, blank_slide):
        elem = {
            "type": "shape",
            "shape": "rectangle",
            "left": 2.0,
            "top": 2.0,
            "width": 3.0,
            "height": 2.0,
        }
        shape = add_shape_element(blank_slide, elem, {}, {})
        assert shape is not None
        assert shape.left == Inches(2.0)

    def test_shape_with_text(self, blank_slide):
        elem = {
            "type": "shape",
            "shape": "rectangle",
            "left": 1.0,
            "top": 1.0,
            "width": 4.0,
            "height": 2.0,
            "text": "Shape Text",
            "text_size": 18,
        }
        shape = add_shape_element(blank_slide, elem, {}, {})
        assert shape.text_frame.text == "Shape Text"

    def test_shape_with_fill(self, blank_slide):
        elem = {
            "type": "shape",
            "shape": "rectangle",
            "left": 1.0,
            "top": 1.0,
            "width": 3.0,
            "height": 2.0,
            "fill": "#0078D4",
        }
        shape = add_shape_element(blank_slide, elem, {}, {})
        assert shape is not None

    def test_shape_with_name(self, blank_slide):
        elem = {
            "type": "shape",
            "shape": "rectangle",
            "left": 1.0,
            "top": 1.0,
            "width": 3.0,
            "height": 2.0,
            "name": "MyShape",
        }
        shape = add_shape_element(blank_slide, elem, {}, {})
        assert shape.name == "MyShape"

    def test_shape_effect_applied(self, blank_slide):
        elem = {
            "type": "shape",
            "shape": "rectangle",
            "left": 1.0,
            "top": 1.0,
            "width": 3.0,
            "height": 2.0,
            "effect": {
                "outer_shadow": {
                    "blur_radius": 50800,
                    "distance": 38100,
                    "direction": 2700000,
                    "color": "#000000",
                    "alpha": 40000,
                }
            },
        }
        shape = add_shape_element(blank_slide, elem, {}, {})
        assert shape is not None

    def test_shape_rotation(self, blank_slide):
        elem = {
            "type": "shape",
            "shape": "rectangle",
            "left": 1.0,
            "top": 1.0,
            "width": 3.0,
            "height": 2.0,
            "rotation": 45.0,
        }
        shape = add_shape_element(blank_slide, elem, {}, {})
        assert shape.rotation == pytest.approx(45.0, abs=0.1)

    def test_shape_paragraphs(self, blank_slide):
        elem = {
            "type": "shape",
            "shape": "rectangle",
            "left": 1.0,
            "top": 1.0,
            "width": 4.0,
            "height": 3.0,
            "text": "multi",
            "paragraphs": [
                {"text": "Para 1", "text_size": 16},
                {"text": "Para 2", "text_size": 14},
            ],
        }
        shape = add_shape_element(blank_slide, elem, {}, {})
        paras = shape.text_frame.paragraphs
        assert len(paras) == 2


class TestAddImageElement:
    """Tests for add_image_element."""

    def test_adds_image(self, blank_slide, sample_image_path):
        content_dir = sample_image_path.parent
        elem = {
            "type": "image",
            "path": sample_image_path.name,
            "left": 1.0,
            "top": 1.0,
            "width": 3.0,
            "height": 2.0,
        }
        pic = add_image_element(blank_slide, elem, content_dir)
        assert pic is not None

    def test_missing_image_fallback(self, blank_slide, tmp_path):
        elem = {
            "type": "image",
            "path": "does_not_exist.png",
            "left": 1.0,
            "top": 1.0,
            "width": 3.0,
            "height": 2.0,
        }
        result = add_image_element(blank_slide, elem, tmp_path)
        # Fallback creates a text placeholder
        assert result is None

    def test_image_with_name(self, blank_slide, sample_image_path):
        content_dir = sample_image_path.parent
        elem = {
            "type": "image",
            "path": sample_image_path.name,
            "left": 1.0,
            "top": 1.0,
            "width": 3.0,
            "height": 2.0,
            "name": "Logo",
        }
        pic = add_image_element(blank_slide, elem, content_dir)
        assert pic.name == "Logo"


class TestAddRichTextElement:
    """Tests for add_rich_text_element."""

    def test_basic_segments(self, blank_slide):
        elem = {
            "type": "rich_text",
            "left": 1.0,
            "top": 1.0,
            "width": 5.0,
            "height": 1.0,
            "segments": [
                {"text": "Bold", "bold": True, "size": 16},
                {"text": " Normal", "size": 16},
            ],
        }
        txBox = add_rich_text_element(blank_slide, elem, {}, {})
        assert txBox is not None
        runs = txBox.text_frame.paragraphs[0].runs
        assert len(runs) == 2
        assert runs[0].text == "Bold"

    def test_segment_with_color(self, blank_slide):
        elem = {
            "type": "rich_text",
            "left": 1.0,
            "top": 1.0,
            "width": 5.0,
            "height": 1.0,
            "segments": [
                {"text": "Red", "color": "#FF0000", "size": 16},
            ],
        }
        txBox = add_rich_text_element(blank_slide, elem, {}, {})
        assert txBox is not None

    def test_rich_text_with_name(self, blank_slide):
        elem = {
            "type": "rich_text",
            "left": 1.0,
            "top": 1.0,
            "width": 5.0,
            "height": 1.0,
            "name": "RichBox",
            "segments": [{"text": "Test", "size": 16}],
        }
        txBox = add_rich_text_element(blank_slide, elem, {}, {})
        assert txBox.name == "RichBox"


class TestAddConnectorElement:
    """Tests for add_connector_element."""

    def test_straight_connector(self, blank_slide):
        elem = {
            "type": "connector",
            "connector_type": "straight",
            "begin_x": 1.0,
            "begin_y": 2.0,
            "end_x": 5.0,
            "end_y": 4.0,
        }
        conn = add_connector_element(blank_slide, elem, {})
        assert conn is not None

    def test_connector_with_name(self, blank_slide):
        elem = {
            "type": "connector",
            "connector_type": "straight",
            "begin_x": 1.0,
            "begin_y": 2.0,
            "end_x": 5.0,
            "end_y": 4.0,
            "name": "Arrow1",
        }
        conn = add_connector_element(blank_slide, elem, {})
        assert conn.name == "Arrow1"


class TestClearSlideShapes:
    """Tests for clear_slide_shapes."""

    def test_clears_all_shapes(self, blank_slide):
        blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(2),
            Inches(2),
        )
        blank_slide.shapes.add_textbox(
            Inches(3),
            Inches(1),
            Inches(2),
            Inches(1),
        )
        assert len(blank_slide.shapes) >= 2
        clear_slide_shapes(blank_slide)
        assert len(blank_slide.shapes) == 0

    def test_clears_empty_slide(self, blank_slide):
        clear_slide_shapes(blank_slide)
        assert len(blank_slide.shapes) == 0


class TestGetSlideLayout:
    """Tests for get_slide_layout."""

    def test_blank_layout(self, blank_presentation):
        style = {}
        content = {"layout": "blank"}
        layout = get_slide_layout(blank_presentation, content, style)
        assert layout is not None

    def test_none_layout(self, blank_presentation):
        style = {}
        content = {}
        layout = get_slide_layout(blank_presentation, content, style)
        assert layout is not None

    def test_named_layout(self, blank_presentation):
        style = {}
        # Use actual layout name from the default template
        first_layout = blank_presentation.slide_layouts[0]
        content = {"layout": first_layout.name}
        layout = get_slide_layout(blank_presentation, content, style)
        assert layout.name == first_layout.name

    def test_index_layout(self, blank_presentation):
        style = {}
        content = {"layout": 0}
        layout = get_slide_layout(blank_presentation, content, style)
        assert layout is not None


class TestDiscoverSlides:
    """Tests for discover_slides."""

    def test_discovers_numbered_dirs(self, tmp_path):
        for i in (1, 3, 2):
            d = tmp_path / f"slide-{i:03d}"
            d.mkdir()
            (d / "content.yaml").write_text(f"slide: {i}\n")
        slides = discover_slides(tmp_path)
        assert len(slides) == 3
        assert [num for num, _ in slides] == [1, 2, 3]

    def test_ignores_non_slide_dirs(self, tmp_path):
        (tmp_path / "global").mkdir()
        (tmp_path / "slide-001").mkdir()
        (tmp_path / "slide-001" / "content.yaml").write_text("slide: 1\n")
        slides = discover_slides(tmp_path)
        assert len(slides) == 1

    def test_empty_dir(self, tmp_path):
        slides = discover_slides(tmp_path)
        assert slides == []

    def test_missing_content_yaml(self, tmp_path):
        (tmp_path / "slide-001").mkdir()
        slides = discover_slides(tmp_path)
        assert slides == []


class TestBuildSlide:
    """Tests for build_slide."""

    def test_build_empty_slide(self, blank_presentation, tmp_path):
        content = {"slide": 1, "elements": []}
        style = {"dimensions": {"width_inches": 13.333, "height_inches": 7.5}}
        slide = build_slide(blank_presentation, content, style, tmp_path)
        assert slide is not None

    def test_build_slide_with_textbox(self, blank_presentation, tmp_path):
        content = {
            "slide": 1,
            "elements": [
                {
                    "type": "textbox",
                    "left": 1.0,
                    "top": 1.0,
                    "width": 4.0,
                    "height": 1.0,
                    "text": "Hello",
                    "font_size": 16,
                },
            ],
        }
        style = {}
        slide = build_slide(blank_presentation, content, style, tmp_path)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert "Hello" in texts

    def test_build_slide_with_shape(self, blank_presentation, tmp_path):
        content = {
            "slide": 1,
            "elements": [
                {
                    "type": "shape",
                    "shape": "rectangle",
                    "left": 2.0,
                    "top": 2.0,
                    "width": 3.0,
                    "height": 2.0,
                    "fill": "#0078D4",
                },
            ],
        }
        style = {}
        slide = build_slide(blank_presentation, content, style, tmp_path)
        assert len(slide.shapes) >= 1

    def test_build_slide_with_speaker_notes(self, blank_presentation, tmp_path):
        content = {
            "slide": 1,
            "elements": [],
            "speaker_notes": "Notes text",
        }
        style = {}
        slide = build_slide(blank_presentation, content, style, tmp_path)
        assert slide.notes_slide.notes_text_frame.text == "Notes text"

    def test_build_slide_empty_notes(self, blank_presentation, tmp_path):
        content = {
            "slide": 1,
            "elements": [],
            "speaker_notes": "",
        }
        style = {}
        slide = build_slide(blank_presentation, content, style, tmp_path)
        assert slide.notes_slide.notes_text_frame.text == ""

    def test_build_existing_slide(self, blank_presentation, tmp_path):
        """Rebuild in-place clears shapes and repopulates."""
        layout = blank_presentation.slide_layouts[6]
        existing = blank_presentation.slides.add_slide(layout)
        existing.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(2),
            Inches(2),
        )
        content = {
            "slide": 1,
            "elements": [
                {
                    "type": "textbox",
                    "left": 1.0,
                    "top": 1.0,
                    "width": 4.0,
                    "height": 1.0,
                    "text": "Rebuilt",
                    "font_size": 16,
                },
            ],
        }
        style = {}
        slide = build_slide(
            blank_presentation,
            content,
            style,
            tmp_path,
            existing_slide=existing,
        )
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert "Rebuilt" in texts

    def test_build_slide_z_order(self, blank_presentation, tmp_path):
        """Elements sorted by z_order."""
        content = {
            "slide": 1,
            "elements": [
                {
                    "type": "textbox",
                    "z_order": 2,
                    "left": 1.0,
                    "top": 1.0,
                    "width": 4.0,
                    "height": 1.0,
                    "text": "Second",
                    "font_size": 16,
                },
                {
                    "type": "textbox",
                    "z_order": 0,
                    "left": 1.0,
                    "top": 2.0,
                    "width": 4.0,
                    "height": 1.0,
                    "text": "First",
                    "font_size": 16,
                },
            ],
        }
        style = {}
        slide = build_slide(blank_presentation, content, style, tmp_path)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert texts == ["First", "Second"]


class TestSetSlideBgImage:
    """Tests for set_slide_bg_image."""

    def test_sets_background_image(self, blank_slide, sample_image_path):
        content_dir = sample_image_path.parent
        set_slide_bg_image(blank_slide, sample_image_path.name, content_dir)
        # Verify bg element was created
        from pptx.oxml.ns import qn

        cSld = blank_slide._element.find(qn("p:cSld"))
        bg = cSld.find(qn("p:bg"))
        assert bg is not None

    def test_missing_image_noop(self, blank_slide, tmp_path):
        set_slide_bg_image(blank_slide, "missing.png", tmp_path)
        # Should not crash


class TestAddCardElement:
    """Tests for add_card_element."""

    def test_basic_card(self, blank_slide):
        elem = {
            "left": 1.0,
            "top": 1.0,
            "width": 5.0,
            "height": 3.0,
            "title": "Card Title",
            "content": [
                {"bullet": "Point 1"},
                {"bullet": "Point 2"},
            ],
        }
        shape = add_card_element(blank_slide, elem, {}, {})
        assert shape is not None

    def test_card_with_accent_bar(self, blank_slide):
        elem = {
            "left": 1.0,
            "top": 1.0,
            "width": 5.0,
            "height": 3.0,
            "accent_bar": True,
            "accent_color": "#FF5733",
        }
        shape = add_card_element(blank_slide, elem, {}, {})
        assert shape is not None

    def test_card_with_border(self, blank_slide):
        elem = {
            "left": 1.0,
            "top": 1.0,
            "width": 5.0,
            "height": 3.0,
            "border_color": "#333333",
            "border_width": 2,
        }
        shape = add_card_element(blank_slide, elem, {}, {})
        assert shape is not None


class TestAddArrowFlowElement:
    """Tests for add_arrow_flow_element."""

    def test_basic_flow(self, blank_slide):
        elem = {
            "left": 1.0,
            "top": 2.0,
            "width": 10.0,
            "height": 1.5,
            "items": [
                {"label": "Step 1", "color": "#0078D4"},
                {"label": "Step 2", "color": "#00B050"},
                {"label": "Step 3", "color": "#FF5733"},
            ],
        }
        add_arrow_flow_element(blank_slide, elem, {}, {})
        # Verify chevrons created
        shapes = [s for s in blank_slide.shapes if s.has_text_frame]
        labels = [s.text_frame.text for s in shapes]
        assert "Step 1" in labels

    def test_empty_items(self, blank_slide):
        elem = {
            "left": 1.0,
            "top": 2.0,
            "width": 10.0,
            "height": 1.5,
            "items": [],
        }
        result = add_arrow_flow_element(blank_slide, elem, {}, {})
        assert result is None


class TestAddNumberedStepElement:
    """Tests for add_numbered_step_element."""

    def test_basic_step(self, blank_slide):
        elem = {
            "left": 1.0,
            "top": 1.0,
            "width": 6.0,
            "height": 1.0,
            "number": 1,
            "label": "First Step",
        }
        add_numbered_step_element(blank_slide, elem, {}, {})
        texts = [s.text_frame.text for s in blank_slide.shapes if s.has_text_frame]
        assert "1" in texts
        assert "First Step" in texts

    def test_step_with_description(self, blank_slide):
        elem = {
            "left": 1.0,
            "top": 1.0,
            "width": 6.0,
            "height": 1.5,
            "number": 2,
            "label": "Second Step",
            "description": "Detailed explanation",
        }
        add_numbered_step_element(blank_slide, elem, {}, {})
        texts = [s.text_frame.text for s in blank_slide.shapes if s.has_text_frame]
        assert "Detailed explanation" in texts


class TestAddGroupElement:
    """Tests for add_group_element."""

    def test_basic_group(self, blank_slide, tmp_path):
        elem = {
            "left": 1.0,
            "top": 1.0,
            "width": 5.0,
            "height": 3.0,
            "elements": [
                {
                    "type": "shape",
                    "shape": "rectangle",
                    "left": 0,
                    "top": 0,
                    "width": 5.0,
                    "height": 3.0,
                    "fill": "#2D2D35",
                },
            ],
        }
        group = add_group_element(blank_slide, elem, {}, {}, tmp_path)
        assert group is not None
        assert len(group.shapes) >= 1

    def test_group_with_textbox(self, blank_slide, tmp_path):
        elem = {
            "left": 1.0,
            "top": 1.0,
            "width": 5.0,
            "height": 3.0,
            "elements": [
                {
                    "type": "textbox",
                    "left": 0.2,
                    "top": 0.2,
                    "width": 4.6,
                    "height": 0.5,
                    "text": "Group Title",
                },
            ],
        }
        group = add_group_element(blank_slide, elem, {}, {}, tmp_path)
        assert len(group.shapes) >= 1

    def test_group_with_name(self, blank_slide, tmp_path):
        elem = {
            "left": 1.0,
            "top": 1.0,
            "width": 5.0,
            "height": 3.0,
            "name": "MyGroup",
            "elements": [],
        }
        group = add_group_element(blank_slide, elem, {}, {}, tmp_path)
        assert group.name == "MyGroup"

    def test_group_with_connector(self, blank_slide, tmp_path):
        elem = {
            "left": 1.0,
            "top": 1.0,
            "width": 5.0,
            "height": 3.0,
            "elements": [
                {
                    "type": "connector",
                    "connector_type": "straight",
                    "begin_x": 0,
                    "begin_y": 0,
                    "end_x": 3.0,
                    "end_y": 2.0,
                },
            ],
        }
        group = add_group_element(blank_slide, elem, {}, {}, tmp_path)
        assert group is not None

    def test_group_shape_with_text(self, blank_slide, tmp_path):
        elem = {
            "left": 1.0,
            "top": 1.0,
            "width": 5.0,
            "height": 3.0,
            "elements": [
                {
                    "type": "shape",
                    "shape": "rectangle",
                    "left": 0,
                    "top": 0,
                    "width": 4.0,
                    "height": 2.0,
                    "text": "Shape in group",
                    "text_size": 14,
                    "text_color": "#FFFFFF",
                },
            ],
        }
        group = add_group_element(blank_slide, elem, {}, {}, tmp_path)
        assert len(group.shapes) >= 1

    def test_group_textbox_with_paragraphs(self, blank_slide, tmp_path):
        elem = {
            "left": 1.0,
            "top": 1.0,
            "width": 5.0,
            "height": 3.0,
            "elements": [
                {
                    "type": "textbox",
                    "left": 0.2,
                    "top": 0.2,
                    "width": 4.6,
                    "height": 2.0,
                    "text": "default",
                    "paragraphs": [
                        {
                            "text": "Para 1",
                            "font_size": 16,
                            "font_bold": True,
                        },
                        {"text": "Para 2", "font_size": 14},
                    ],
                },
            ],
        }
        group = add_group_element(blank_slide, elem, {}, {}, tmp_path)
        assert len(group.shapes) >= 1

    def test_group_shape_with_paragraphs(self, blank_slide, tmp_path):
        elem = {
            "left": 1.0,
            "top": 1.0,
            "width": 5.0,
            "height": 3.0,
            "elements": [
                {
                    "type": "shape",
                    "shape": "rectangle",
                    "left": 0,
                    "top": 0,
                    "width": 4.0,
                    "height": 2.0,
                    "text": "default",
                    "paragraphs": [
                        {
                            "text": "P1",
                            "text_size": 16,
                            "text_bold": True,
                        },
                        {"text": "P2", "text_size": 14},
                    ],
                },
            ],
        }
        group = add_group_element(blank_slide, elem, {}, {}, tmp_path)
        assert len(group.shapes) >= 1

    def test_depth_limit_raises(self, blank_slide, tmp_path):
        """Exceeding max_depth raises ValueError."""
        elem = {
            "left": 1.0,
            "top": 1.0,
            "width": 5.0,
            "height": 3.0,
            "elements": [],
        }
        with pytest.raises(ValueError, match="exceeds limit"):
            add_group_element(
                blank_slide,
                elem,
                {},
                {},
                tmp_path,
                _depth=5,
                max_depth=5,
            )

    def test_nested_group_within_depth_limit(self, blank_slide, tmp_path):
        """Nested groups within the limit build successfully."""
        elem = {
            "left": 1.0,
            "top": 1.0,
            "width": 5.0,
            "height": 3.0,
            "elements": [
                {
                    "type": "group",
                    "left": 0.0,
                    "top": 0.0,
                    "width": 2.0,
                    "height": 1.0,
                    "elements": [],
                },
            ],
        }
        group = add_group_element(
            blank_slide,
            elem,
            {},
            {},
            tmp_path,
            _depth=0,
            max_depth=20,
        )
        assert group is not None

    def test_build_element_in_group_dispatches_group(self, blank_slide, tmp_path):
        """build_element_in_group handles nested group type."""
        parent_group = blank_slide.shapes.add_group_shape()
        child_elem = {
            "type": "group",
            "left": 0.0,
            "top": 0.0,
            "width": 2.0,
            "height": 1.0,
            "elements": [],
        }
        build_element_in_group(
            parent_group,
            child_elem,
            {},
            {},
            tmp_path,
            _depth=0,
            max_depth=20,
        )


class TestAddImageElementExtended:
    """Extended tests for add_image_element covering crop and opacity."""

    def test_image_with_crop(self, blank_slide, sample_image_path):
        content_dir = sample_image_path.parent
        elem = {
            "type": "image",
            "path": sample_image_path.name,
            "left": 1.0,
            "top": 1.0,
            "width": 3.0,
            "height": 2.0,
            "crop": {"l": "10000", "t": "5000", "r": "10000", "b": "5000"},
        }
        pic = add_image_element(blank_slide, elem, content_dir)
        assert pic is not None

    def test_image_with_opacity(self, blank_slide, sample_image_path):
        content_dir = sample_image_path.parent
        elem = {
            "type": "image",
            "path": sample_image_path.name,
            "left": 1.0,
            "top": 1.0,
            "width": 3.0,
            "height": 2.0,
            "opacity": 75.0,
        }
        pic = add_image_element(blank_slide, elem, content_dir)
        assert pic is not None

    def test_image_with_blip_fill_attrs(self, blank_slide, sample_image_path):
        content_dir = sample_image_path.parent
        elem = {
            "type": "image",
            "path": sample_image_path.name,
            "left": 1.0,
            "top": 1.0,
            "width": 3.0,
            "height": 2.0,
            "blip_fill_attrs": {"rotWithShape": "1", "dpi": "0"},
        }
        pic = add_image_element(blank_slide, elem, content_dir)
        assert pic is not None


class TestAddShapeElementExtended:
    """Extended tests for shape element with text branches."""

    def test_shape_multiline_text(self, blank_slide):
        elem = {
            "type": "shape",
            "shape": "rectangle",
            "left": 1.0,
            "top": 1.0,
            "width": 4.0,
            "height": 3.0,
            "text": "Line1\nLine2",
            "text_size": 16,
            "alignment": "center",
        }
        shape = add_shape_element(blank_slide, elem, {}, {})
        paras = shape.text_frame.paragraphs
        assert len(paras) == 2

    def test_shape_with_paragraph_runs(self, blank_slide):
        elem = {
            "type": "shape",
            "shape": "rectangle",
            "left": 1.0,
            "top": 1.0,
            "width": 4.0,
            "height": 3.0,
            "text": "multi",
            "text_color": "#FFFFFF",
            "paragraphs": [
                {
                    "text": "Mixed",
                    "runs": [
                        {"text": "Bold", "bold": True, "size": 16},
                        {
                            "text": " Color",
                            "size": 16,
                            "color": "#FF0000",
                        },
                    ],
                },
            ],
        }
        shape = add_shape_element(blank_slide, elem, {}, {})
        runs = shape.text_frame.paragraphs[0].runs
        assert len(runs) == 2

    def test_shape_with_corner_radius(self, blank_slide):
        elem = {
            "type": "shape",
            "shape": "rounded_rectangle",
            "left": 1.0,
            "top": 1.0,
            "width": 4.0,
            "height": 3.0,
            "corner_radius": 0.05,
        }
        shape = add_shape_element(blank_slide, elem, {}, {})
        assert shape is not None


class TestAddConnectorExtended:
    """Extended tests for connector element with arrow heads."""

    def test_connector_with_arrows(self, blank_slide):
        elem = {
            "type": "connector",
            "connector_type": "straight",
            "begin_x": 1.0,
            "begin_y": 2.0,
            "end_x": 5.0,
            "end_y": 4.0,
            "line_color": "#0078D4",
            "line_width": 2,
            "head_end": "arrow",
            "tail_end": "arrow",
        }
        conn = add_connector_element(blank_slide, elem, {})
        assert conn is not None

    def test_connector_curve_type(self, blank_slide):
        elem = {
            "type": "connector",
            "connector_type": "curve",
            "begin_x": 1.0,
            "begin_y": 2.0,
            "end_x": 5.0,
            "end_y": 4.0,
        }
        conn = add_connector_element(blank_slide, elem, {})
        assert conn is not None


class TestGetSlideLayoutExtended:
    """Extended layout selection tests."""

    def test_style_map_int_layout(self, blank_presentation):
        style = {"layouts": {"title": 0}}
        content = {"layout": "title"}
        layout = get_slide_layout(blank_presentation, content, style)
        assert layout is not None

    def test_style_map_name_layout(self, blank_presentation):
        first_layout = blank_presentation.slide_layouts[0]
        style = {"layouts": {"custom": first_layout.name}}
        content = {"layout": "custom"}
        layout = get_slide_layout(blank_presentation, content, style)
        assert layout.name == first_layout.name

    def test_invalid_index_fallback(self, blank_presentation):
        style = {}
        content = {"layout": 999}
        layout = get_slide_layout(blank_presentation, content, style)
        assert layout is not None

    def test_unknown_name_fallback(self, blank_presentation):
        style = {}
        content = {"layout": "nonexistent_layout_xyz"}
        layout = get_slide_layout(blank_presentation, content, style)
        assert layout is not None


class TestBuildSlideExtended:
    """Extended build_slide tests for more element types."""

    def test_build_slide_with_connector(self, blank_presentation, tmp_path):
        content = {
            "slide": 1,
            "elements": [
                {
                    "type": "connector",
                    "connector_type": "straight",
                    "begin_x": 1.0,
                    "begin_y": 2.0,
                    "end_x": 5.0,
                    "end_y": 4.0,
                },
            ],
        }
        slide = build_slide(blank_presentation, content, {}, tmp_path)
        assert len(slide.shapes) >= 1

    def test_build_slide_with_background_fill(self, blank_presentation, tmp_path):
        content = {
            "slide": 1,
            "elements": [],
            "background": {"fill": "#1A1A2E"},
        }
        slide = build_slide(blank_presentation, content, {}, tmp_path)
        assert slide is not None

    def test_build_slide_with_background_image(
        self, blank_presentation, sample_image_path
    ):
        content_dir = sample_image_path.parent
        content = {
            "slide": 1,
            "elements": [],
            "background": {"image": sample_image_path.name},
        }
        slide = build_slide(blank_presentation, content, {}, content_dir)
        assert slide is not None

    def test_build_slide_with_rich_text(self, blank_presentation, tmp_path):
        content = {
            "slide": 1,
            "elements": [
                {
                    "type": "rich_text",
                    "left": 1.0,
                    "top": 1.0,
                    "width": 5.0,
                    "height": 1.0,
                    "segments": [
                        {"text": "Bold", "bold": True, "size": 16},
                        {"text": " Normal", "size": 16},
                    ],
                },
            ],
        }
        slide = build_slide(blank_presentation, content, {}, tmp_path)
        assert len(slide.shapes) >= 1

    def test_build_slide_with_card(self, blank_presentation, tmp_path):
        content = {
            "slide": 1,
            "elements": [
                {
                    "type": "card",
                    "left": 1.0,
                    "top": 1.0,
                    "width": 5.0,
                    "height": 3.0,
                    "title": "Card",
                    "content": [{"bullet": "Item"}],
                },
            ],
        }
        slide = build_slide(blank_presentation, content, {}, tmp_path)
        assert len(slide.shapes) >= 1

    def test_build_slide_with_arrow_flow(self, blank_presentation, tmp_path):
        content = {
            "slide": 1,
            "elements": [
                {
                    "type": "arrow_flow",
                    "left": 1.0,
                    "top": 2.0,
                    "width": 10.0,
                    "height": 1.5,
                    "items": [
                        {"label": "A", "color": "#0078D4"},
                        {"label": "B", "color": "#00B050"},
                    ],
                },
            ],
        }
        slide = build_slide(blank_presentation, content, {}, tmp_path)
        assert len(slide.shapes) >= 2

    def test_build_slide_with_numbered_step(self, blank_presentation, tmp_path):
        content = {
            "slide": 1,
            "elements": [
                {
                    "type": "numbered_step",
                    "left": 1.0,
                    "top": 1.0,
                    "width": 6.0,
                    "height": 1.0,
                    "number": 1,
                    "label": "Step One",
                },
            ],
        }
        slide = build_slide(blank_presentation, content, {}, tmp_path)
        assert len(slide.shapes) >= 2

    def test_build_slide_with_group(self, blank_presentation, tmp_path):
        content = {
            "slide": 1,
            "elements": [
                {
                    "type": "group",
                    "left": 1.0,
                    "top": 1.0,
                    "width": 5.0,
                    "height": 3.0,
                    "elements": [
                        {
                            "type": "shape",
                            "shape": "rectangle",
                            "left": 0,
                            "top": 0,
                            "width": 5.0,
                            "height": 3.0,
                        },
                    ],
                },
            ],
        }
        slide = build_slide(blank_presentation, content, {}, tmp_path)
        assert len(slide.shapes) >= 1

    def test_build_slide_with_image(self, blank_presentation, sample_image_path):
        content_dir = sample_image_path.parent
        content = {
            "slide": 1,
            "elements": [
                {
                    "type": "image",
                    "path": sample_image_path.name,
                    "left": 1.0,
                    "top": 1.0,
                    "width": 3.0,
                    "height": 2.0,
                },
            ],
        }
        slide = build_slide(blank_presentation, content, {}, content_dir)
        assert len(slide.shapes) >= 1

    def test_build_slide_turbo_mode(self, blank_presentation, tmp_path):
        """20+ elements activate turbo mode."""
        elements = [
            {
                "type": "textbox",
                "left": float(i % 10),
                "top": float(i // 10),
                "width": 1.0,
                "height": 0.5,
                "text": f"Box {i}",
                "font_size": 10,
            }
            for i in range(25)
        ]
        content = {"slide": 1, "elements": elements}
        slide = build_slide(blank_presentation, content, {}, tmp_path)
        assert len(slide.shapes) >= 25

    def test_build_slide_with_placeholders(self, blank_presentation, tmp_path):
        # Use a layout that has placeholders
        layout = blank_presentation.slide_layouts[0]
        slide = blank_presentation.slides.add_slide(layout)
        content = {
            "slide": 1,
            "elements": [],
            "placeholders": {"0": "Title Text"},
        }
        build_slide(
            blank_presentation,
            content,
            {},
            tmp_path,
            existing_slide=slide,
        )
        # Check placeholder was populated
        if 0 in slide.placeholders:
            assert slide.placeholders[0].text == "Title Text"


class TestMain:
    """Tests for main() CLI entry point."""

    def _setup_content_dir(self, tmp_path):
        """Create a content directory with one slide for main() to discover."""
        content_dir = tmp_path / "content"
        slide_dir = content_dir / "slide-001"
        slide_dir.mkdir(parents=True)
        (slide_dir / "content.yaml").write_text("slide: 1\ntitle: Test\nelements: []\n")
        style_file = tmp_path / "style.yaml"
        style_yaml = "dimensions:\n  width_inches: 13.333\n  height_inches: 7.5\n"
        style_file.write_text(style_yaml)
        return content_dir, style_file

    def test_full_build(self, mocker, tmp_path):
        content_dir, style_file = self._setup_content_dir(tmp_path)
        output = tmp_path / "out" / "deck.pptx"

        mock_prs_cls = mocker.patch("build_deck.Presentation")
        mock_build_slide = mocker.patch("build_deck.build_slide")

        mock_prs = MagicMock()
        mock_prs.slides.__len__ = MagicMock(return_value=1)
        mock_prs_cls.return_value = mock_prs

        mocker.patch(
            "sys.argv",
            [
                "build_deck.py",
                "--content-dir",
                str(content_dir),
                "--style",
                str(style_file),
                "--output",
                str(output),
            ],
        )
        main()

        mock_prs_cls.assert_called_once_with()
        mock_build_slide.assert_called_once()
        mock_prs.save.assert_called_once_with(str(output))

    def test_full_build_no_slides(self, mocker, tmp_path):
        """No slide directories found exits with code 1."""
        empty_content = tmp_path / "content"
        empty_content.mkdir()
        style_file = tmp_path / "style.yaml"
        style_yaml = "dimensions:\n  width_inches: 13.333\n  height_inches: 7.5\n"
        style_file.write_text(style_yaml)
        output = tmp_path / "deck.pptx"

        mock_prs_cls = mocker.patch("build_deck.Presentation")
        mock_build_slide = mocker.patch("build_deck.build_slide")

        mock_prs_cls.return_value = MagicMock()

        mocker.patch(
            "sys.argv",
            [
                "build_deck.py",
                "--content-dir",
                str(empty_content),
                "--style",
                str(style_file),
                "--output",
                str(output),
            ],
        )
        with pytest.raises(SystemExit) as exc_info:
            main()

        assert exc_info.value.code == 1
        mock_build_slide.assert_not_called()

    def test_full_build_with_metadata(self, mocker, tmp_path):
        content_dir, _ = self._setup_content_dir(tmp_path)
        style_file = tmp_path / "style.yaml"
        style_file.write_text(
            "dimensions:\n  width_inches: 13.333\n  height_inches: 7.5\n"
            "metadata:\n  title: My Deck\n  author: Tester\n"
        )
        output = tmp_path / "deck.pptx"

        mock_prs_cls = mocker.patch("build_deck.Presentation")
        mocker.patch("build_deck.build_slide")

        mock_prs = MagicMock()
        mock_prs.slides.__len__ = MagicMock(return_value=1)
        mock_props = MagicMock(spec=["title", "author"])
        mock_prs.core_properties = mock_props
        mock_prs_cls.return_value = mock_prs

        mocker.patch(
            "sys.argv",
            [
                "build_deck.py",
                "--content-dir",
                str(content_dir),
                "--style",
                str(style_file),
                "--output",
                str(output),
            ],
        )
        main()

        assert mock_props.title == "My Deck"
        assert mock_props.author == "Tester"

    def test_template_build(self, mocker, tmp_path):
        content_dir, style_file = self._setup_content_dir(tmp_path)
        template = tmp_path / "template.pptx"
        output = tmp_path / "deck.pptx"

        mock_prs_cls = mocker.patch("build_deck.Presentation")
        mock_build_slide = mocker.patch("build_deck.build_slide")

        mock_prs = MagicMock()
        sld_entry = MagicMock()
        sld_entry.rId = "rId1"
        sld_id_lst = [sld_entry]
        mock_prs.slides._sldIdLst = sld_id_lst
        mock_prs.slides.__len__ = MagicMock(side_effect=lambda: len(sld_id_lst))
        mock_prs_cls.return_value = mock_prs

        mocker.patch(
            "sys.argv",
            [
                "build_deck.py",
                "--content-dir",
                str(content_dir),
                "--style",
                str(style_file),
                "--output",
                str(output),
                "--template",
                str(template),
            ],
        )
        main()

        mock_prs_cls.assert_called_once_with(str(template))
        mock_build_slide.assert_called_once()
        mock_prs.save.assert_called_once()

    def test_template_build_no_slides(self, mocker, tmp_path):
        empty_content = tmp_path / "content"
        empty_content.mkdir()
        style_file = tmp_path / "style.yaml"
        style_file.write_text("{}\n")
        template = tmp_path / "template.pptx"
        output = tmp_path / "deck.pptx"

        mock_prs_cls = mocker.patch("build_deck.Presentation")
        mocker.patch("build_deck.build_slide")

        mock_prs = MagicMock()
        mock_prs.slides._sldIdLst = []
        mock_prs_cls.return_value = mock_prs

        mocker.patch(
            "sys.argv",
            [
                "build_deck.py",
                "--content-dir",
                str(empty_content),
                "--style",
                str(style_file),
                "--output",
                str(output),
                "--template",
                str(template),
            ],
        )
        with pytest.raises(SystemExit) as exc_info:
            main()

        assert exc_info.value.code == 1

    def test_partial_rebuild(self, mocker, tmp_path):
        content_dir, style_file = self._setup_content_dir(tmp_path)
        source = tmp_path / "source.pptx"
        output = tmp_path / "deck.pptx"

        mock_prs_cls = mocker.patch("build_deck.Presentation")
        mock_build_slide = mocker.patch("build_deck.build_slide")

        mock_slide = MagicMock()
        mock_prs = MagicMock()
        mock_prs.slides.__len__ = MagicMock(return_value=1)
        mock_prs.slides.__getitem__ = MagicMock(return_value=mock_slide)
        mock_prs_cls.return_value = mock_prs

        mocker.patch(
            "sys.argv",
            [
                "build_deck.py",
                "--content-dir",
                str(content_dir),
                "--style",
                str(style_file),
                "--output",
                str(output),
                "--source",
                str(source),
                "--slides",
                "1",
            ],
        )
        main()

        mock_prs_cls.assert_called_once_with(str(source))
        mock_build_slide.assert_called_once()
        mock_prs.save.assert_called_once()

    def test_partial_rebuild_missing_slide(self, mocker, tmp_path):
        """Slide number not found in content directory prints warning and skips."""
        content_dir, style_file = self._setup_content_dir(tmp_path)
        source = tmp_path / "source.pptx"
        output = tmp_path / "deck.pptx"

        mock_prs_cls = mocker.patch("build_deck.Presentation")
        mock_build_slide = mocker.patch("build_deck.build_slide")

        mock_prs = MagicMock()
        mock_prs.slides.__len__ = MagicMock(return_value=1)
        mock_prs_cls.return_value = mock_prs

        mocker.patch(
            "sys.argv",
            [
                "build_deck.py",
                "--content-dir",
                str(content_dir),
                "--style",
                str(style_file),
                "--output",
                str(output),
                "--source",
                str(source),
                "--slides",
                "99",
            ],
        )
        main()

        mock_build_slide.assert_not_called()
        mock_prs.save.assert_called_once()

    def test_partial_rebuild_out_of_range(self, mocker, tmp_path):
        """Slide index beyond deck length prints warning and skips."""
        content_dir, style_file = self._setup_content_dir(tmp_path)
        source = tmp_path / "source.pptx"
        output = tmp_path / "deck.pptx"

        mock_prs_cls = mocker.patch("build_deck.Presentation")
        mock_build_slide = mocker.patch("build_deck.build_slide")

        mock_prs = MagicMock()
        mock_prs.slides.__len__ = MagicMock(return_value=0)
        mock_prs_cls.return_value = mock_prs

        mocker.patch(
            "sys.argv",
            [
                "build_deck.py",
                "--content-dir",
                str(content_dir),
                "--style",
                str(style_file),
                "--output",
                str(output),
                "--source",
                str(source),
                "--slides",
                "1",
            ],
        )
        main()

        mock_build_slide.assert_not_called()
        mock_prs.save.assert_called_once()


class TestContentExtraValidation:
    """Tests for _validate_content_extra AST security checks."""

    def test_valid_pptx_imports(self, tmp_path):
        """Script using only pptx imports passes validation."""
        script = tmp_path / "content-extra.py"
        script.write_text(
            "from pptx.util import Inches, Pt\n"
            "from pptx.dml.color import RGBColor\n"
            "def render(slide, style, content_dir):\n"
            "    pass\n"
        )
        _validate_content_extra(script)

    def test_valid_safe_stdlib_imports(self, tmp_path):
        """Script using safe stdlib modules passes validation."""
        script = tmp_path / "content-extra.py"
        script.write_text(
            "import math\n"
            "import json\n"
            "from pathlib import Path\n"
            "def render(slide, style, content_dir):\n"
            "    pass\n"
        )
        _validate_content_extra(script)

    def test_blocked_subprocess(self, tmp_path):
        """Script importing subprocess is rejected."""
        script = tmp_path / "content-extra.py"
        script.write_text("import subprocess\n")
        with pytest.raises(ContentExtraError, match="Blocked import 'subprocess'"):
            _validate_content_extra(script)

    def test_blocked_os(self, tmp_path):
        """Script importing os is rejected."""
        script = tmp_path / "content-extra.py"
        script.write_text("import os\n")
        with pytest.raises(ContentExtraError, match="Blocked import 'os'"):
            _validate_content_extra(script)

    def test_blocked_from_os_import(self, tmp_path):
        """from-import of a blocked module is rejected."""
        script = tmp_path / "content-extra.py"
        script.write_text("from os.path import join\n")
        with pytest.raises(ContentExtraError, match="Blocked import 'os.path'"):
            _validate_content_extra(script)

    def test_blocked_shutil(self, tmp_path):
        """Script importing shutil is rejected."""
        script = tmp_path / "content-extra.py"
        script.write_text("import shutil\n")
        with pytest.raises(ContentExtraError, match="Blocked import 'shutil'"):
            _validate_content_extra(script)

    def test_blocked_socket(self, tmp_path):
        """Script importing socket is rejected."""
        script = tmp_path / "content-extra.py"
        script.write_text("import socket\n")
        with pytest.raises(ContentExtraError, match="Blocked import 'socket'"):
            _validate_content_extra(script)

    def test_blocked_ctypes(self, tmp_path):
        """Script importing ctypes is rejected."""
        script = tmp_path / "content-extra.py"
        script.write_text("import ctypes\n")
        with pytest.raises(ContentExtraError, match="Blocked import 'ctypes'"):
            _validate_content_extra(script)

    def test_third_party_import_rejected(self, tmp_path):
        """Script importing a non-stdlib, non-pptx package is rejected."""
        script = tmp_path / "content-extra.py"
        script.write_text("import requests\n")
        with pytest.raises(ContentExtraError, match="Disallowed import 'requests'"):
            _validate_content_extra(script)

    def test_dangerous_eval(self, tmp_path):
        """Script calling eval() is rejected."""
        script = tmp_path / "content-extra.py"
        script.write_text("x = eval('1+1')\n")
        with pytest.raises(ContentExtraError, match="Dangerous builtin 'eval'"):
            _validate_content_extra(script)

    def test_dangerous_exec(self, tmp_path):
        """Script calling exec() is rejected."""
        script = tmp_path / "content-extra.py"
        script.write_text("exec('import os')\n")
        with pytest.raises(ContentExtraError, match="Dangerous builtin 'exec'"):
            _validate_content_extra(script)

    def test_dangerous_dunder_import(self, tmp_path):
        """Script calling __import__() is rejected."""
        script = tmp_path / "content-extra.py"
        script.write_text("m = __import__('os')\n")
        with pytest.raises(ContentExtraError, match="Dangerous builtin '__import__'"):
            _validate_content_extra(script)

    def test_dangerous_compile(self, tmp_path):
        """Script calling compile() is rejected."""
        script = tmp_path / "content-extra.py"
        script.write_text("c = compile('pass', '<str>', 'exec')\n")
        with pytest.raises(ContentExtraError, match="Dangerous builtin 'compile'"):
            _validate_content_extra(script)

    def test_syntax_error_rejected(self, tmp_path):
        """Script with a syntax error is rejected."""
        script = tmp_path / "content-extra.py"
        script.write_text("def render(:\n")
        with pytest.raises(ContentExtraError, match="Syntax error"):
            _validate_content_extra(script)

    def test_build_slide_runs_valid_content_extra(self, blank_presentation, tmp_path):
        """build_slide executes a valid content-extra.py render function."""
        content_dir = tmp_path / "slide-001"
        content_dir.mkdir()
        (content_dir / "content.yaml").write_text("")

        marker_file = tmp_path / "render_called"
        script = content_dir / "content-extra.py"
        script.write_text(
            "from pathlib import Path\n"
            "def render(slide, style, content_dir):\n"
            f"    Path(r'{marker_file}').write_text('yes')\n"
        )

        slide_content = {"layout": "Blank", "elements": []}
        build_slide(blank_presentation, slide_content, {}, content_dir)
        assert marker_file.read_text() == "yes"

    def test_build_slide_rejects_bad_content_extra(self, blank_presentation, tmp_path):
        """build_slide refuses to execute a content-extra.py with blocked imports."""
        content_dir = tmp_path / "slide-001"
        content_dir.mkdir()
        (content_dir / "content.yaml").write_text("")

        script = content_dir / "content-extra.py"
        script.write_text("import subprocess\ndef render(s,st,d): pass\n")

        slide_content = {"layout": "Blank", "elements": []}
        with pytest.raises(ContentExtraError, match="Blocked import 'subprocess'"):
            build_slide(blank_presentation, slide_content, {}, content_dir)

    def test_dangerous_breakpoint(self, tmp_path):
        """Script calling breakpoint() is rejected."""
        script = tmp_path / "content-extra.py"
        script.write_text("breakpoint()\n")
        with pytest.raises(ContentExtraError, match="Dangerous builtin 'breakpoint'"):
            _validate_content_extra(script)

    @pytest.mark.parametrize(
        "builtin_name",
        [
            "delattr",
            "getattr",
            "globals",
            "locals",
            "setattr",
            "vars",
        ],
    )
    def test_indirect_bypass_builtins(self, builtin_name, tmp_path):
        """Indirect bypass builtins are rejected."""
        script = tmp_path / "content-extra.py"
        script.write_text(f"x = {builtin_name}(object)\n")
        with pytest.raises(
            ContentExtraError,
            match=f"Indirect bypass builtin '{builtin_name}'",
        ):
            _validate_content_extra(script)

    def test_allow_scripts_skips_validation(self, blank_presentation, tmp_path):
        """build_slide skips validation when allow_scripts is True."""
        content_dir = tmp_path / "slide-001"
        content_dir.mkdir()
        (content_dir / "content.yaml").write_text("")

        marker_file = tmp_path / "bypass_called"
        script = content_dir / "content-extra.py"
        script.write_text(
            "from pathlib import Path\n"
            "import os\n"
            "def render(slide, style, content_dir):\n"
            f"    Path(r'{marker_file}').write_text('bypassed')\n"
        )

        slide_content = {"layout": "Blank", "elements": []}
        build_slide(
            blank_presentation,
            slide_content,
            {},
            content_dir,
            allow_scripts=True,
        )
        assert marker_file.read_text() == "bypassed"

    def test_allow_scripts_false_still_validates(self, blank_presentation, tmp_path):
        """build_slide validates when allow_scripts is explicitly False."""
        content_dir = tmp_path / "slide-001"
        content_dir.mkdir()
        (content_dir / "content.yaml").write_text("")

        script = content_dir / "content-extra.py"
        script.write_text("import os\ndef render(s,st,d): pass\n")

        slide_content = {"layout": "Blank", "elements": []}
        with pytest.raises(ContentExtraError, match="Blocked import 'os'"):
            build_slide(
                blank_presentation,
                slide_content,
                {},
                content_dir,
                allow_scripts=False,
            )

    def test_restricted_namespace_blocks_eval(self, blank_presentation, tmp_path):
        """Runtime namespace strips dangerous builtins even after AST pass."""
        content_dir = tmp_path / "slide-001"
        content_dir.mkdir()
        (content_dir / "content.yaml").write_text("")

        # Script uses no blocked AST patterns but tries eval at runtime
        # via a string indirection the AST checker cannot catch.
        script = content_dir / "content-extra.py"
        script.write_text(
            "def render(slide, style, content_dir):\n"
            "    fn = __builtins__['eval']\n"
            "    fn('1+1')\n"
        )

        slide_content = {"layout": "Blank", "elements": []}
        with pytest.raises(KeyError, match="eval"):
            build_slide(
                blank_presentation,
                slide_content,
                {},
                content_dir,
                allow_scripts=False,
            )

    def test_restricted_namespace_allows_safe_builtins(
        self, blank_presentation, tmp_path
    ):
        """Safe builtins like len and range remain available."""
        content_dir = tmp_path / "slide-001"
        content_dir.mkdir()
        (content_dir / "content.yaml").write_text("")

        marker = tmp_path / "safe_result"
        script = content_dir / "content-extra.py"
        script.write_text(
            "from pathlib import Path\n"
            "def render(slide, style, content_dir):\n"
            f"    Path(r'{marker}').write_text(str(len(range(5))))\n"
        )

        slide_content = {"layout": "Blank", "elements": []}
        build_slide(
            blank_presentation,
            slide_content,
            {},
            content_dir,
            allow_scripts=False,
        )
        assert marker.read_text() == "5"


class TestAllowScriptsCLI:
    """Integration tests that exercise --allow-scripts through main()."""

    def _setup_with_extra(self, tmp_path, extra_code):
        """Create content dir with a content-extra.py script."""
        content_dir = tmp_path / "content"
        slide_dir = content_dir / "slide-001"
        slide_dir.mkdir(parents=True)
        (slide_dir / "content.yaml").write_text("slide: 1\ntitle: Test\nelements: []\n")
        (slide_dir / "content-extra.py").write_text(extra_code)
        style_file = tmp_path / "style.yaml"
        style_file.write_text(
            "dimensions:\n  width_inches: 13.333\n  height_inches: 7.5\n"
        )
        return content_dir, style_file

    def test_allow_scripts_flag_propagates(self, mocker, tmp_path):
        """--allow-scripts lets a blocked-import script run via main()."""
        marker = tmp_path / "executed"
        extra = (
            "import os\n"
            "from pathlib import Path\n"
            f"def render(slide, style, d): "
            f"Path(r'{marker}').write_text('ok')\n"
        )
        content_dir, style_file = self._setup_with_extra(tmp_path, extra)
        output = tmp_path / "deck.pptx"

        mock_prs_cls = mocker.patch("build_deck.Presentation")

        mock_prs = MagicMock()
        mock_prs.slides.__len__ = MagicMock(return_value=1)
        mock_prs.slide_layouts = MagicMock()
        layout = MagicMock()
        mock_prs.slide_layouts.__getitem__ = MagicMock(return_value=layout)
        mock_prs.slide_layouts.__iter__ = MagicMock(return_value=iter([layout]))
        layout.name = "Blank"
        slide = MagicMock()
        mock_prs.slides.add_slide.return_value = slide
        mock_prs_cls.return_value = mock_prs

        mocker.patch(
            "sys.argv",
            [
                "build_deck.py",
                "--content-dir",
                str(content_dir),
                "--style",
                str(style_file),
                "--output",
                str(output),
                "--allow-scripts",
            ],
        )
        main()

        assert marker.read_text() == "ok"

    def test_no_allow_scripts_rejects_blocked_import(self, mocker, tmp_path):
        """Without --allow-scripts, a blocked-import script fails."""
        extra = "import os\ndef render(s, st, d): pass\n"
        content_dir, style_file = self._setup_with_extra(tmp_path, extra)
        output = tmp_path / "deck.pptx"

        mock_prs_cls = mocker.patch("build_deck.Presentation")

        mock_prs = MagicMock()
        mock_prs.slides.__len__ = MagicMock(return_value=1)
        mock_prs.slide_layouts = MagicMock()
        layout = MagicMock()
        mock_prs.slide_layouts.__getitem__ = MagicMock(return_value=layout)
        mock_prs.slide_layouts.__iter__ = MagicMock(return_value=iter([layout]))
        layout.name = "Blank"
        slide = MagicMock()
        mock_prs.slides.add_slide.return_value = slide
        mock_prs_cls.return_value = mock_prs

        mocker.patch(
            "sys.argv",
            [
                "build_deck.py",
                "--content-dir",
                str(content_dir),
                "--style",
                str(style_file),
                "--output",
                str(output),
            ],
        )
        with pytest.raises(ContentExtraError):
            main()
