# Copyright (c) Microsoft Corporation.
# SPDX-License-Identifier: MIT
"""Tests for pptx_fonts module."""

import pytest
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx_fonts import (
    ALIGNMENT_MAP,
    ALIGNMENT_REVERSE_MAP,
    FONT_WEIGHT_SUFFIXES,
    _extract_char_spacing,
    extract_alignment,
    extract_font_info,
    extract_paragraph_font,
    font_family_matches,
    normalize_font_family,
)


class TestNormalizeFontFamily:
    """Tests for normalize_font_family."""

    @pytest.mark.parametrize(
        "input_name,expected",
        [
            ("Segoe UI", "Segoe UI"),
            ("Segoe UI Semibold", "Segoe UI"),
            ("Segoe UI SemiBold", "Segoe UI"),
            ("Arial Bold", "Arial"),
            ("Segoe UI Light", "Segoe UI"),
            ("Roboto Thin", "Roboto"),
            ("Montserrat Black", "Montserrat"),
            ("Inter Medium", "Inter"),
            ("Open Sans ExtraBold", "Open Sans"),
            ("Noto Sans ExtraLight", "Noto Sans"),
        ],
    )
    def test_normalize(self, input_name, expected):
        assert normalize_font_family(input_name) == expected


class TestFontFamilyMatches:
    """Tests for font_family_matches."""

    @pytest.mark.parametrize(
        "font,expected_set,result",
        [
            ("Arial", {"Arial", "Segoe UI"}, True),
            ("Comic Sans", {"Arial", "Segoe UI"}, False),
            ("Segoe UI Semibold", {"Segoe UI"}, True),
            ("Segoe UI", {"Segoe UI Bold"}, True),
            ("Arial Bold", {"Segoe UI"}, False),
            ("Arial", set(), False),
        ],
    )
    def test_matches(self, font, expected_set, result):
        assert font_family_matches(font, expected_set) is result


class TestExtractFontInfo:
    """Tests for extract_font_info with real python-pptx font objects."""

    def test_name_and_size(self, sample_textbox):
        run = sample_textbox.text_frame.paragraphs[0].runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(24)
        info = extract_font_info(run.font)
        assert info["font"] == "Arial"
        assert info["size"] == 24

    def test_bold(self, sample_textbox):
        run = sample_textbox.text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        info = extract_font_info(run.font)
        assert info["bold"] is True

    def test_italic(self, sample_textbox):
        run = sample_textbox.text_frame.paragraphs[0].runs[0]
        run.font.italic = True
        info = extract_font_info(run.font)
        assert info["italic"] is True

    def test_underline(self, sample_textbox):
        run = sample_textbox.text_frame.paragraphs[0].runs[0]
        run.font.underline = True
        info = extract_font_info(run.font)
        assert info["underline"] is True

    def test_color(self, sample_textbox):
        run = sample_textbox.text_frame.paragraphs[0].runs[0]
        run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
        info = extract_font_info(run.font)
        assert info["color"] == "#FF0000"

    def test_missing_properties_omitted(self, blank_slide):
        txBox = blank_slide.shapes.add_textbox(
            Inches(0), Inches(0), Inches(1), Inches(1)
        )
        txBox.text_frame.text = "X"
        run = txBox.text_frame.paragraphs[0].runs[0]
        # Don't set any properties explicitly
        info = extract_font_info(run.font)
        # Only keys for set values should appear
        assert "bold" not in info or info.get("bold") is not True


class TestExtractCharSpacing:
    """Tests for _extract_char_spacing."""

    def test_with_spacing(self, sample_textbox):
        run = sample_textbox.text_frame.paragraphs[0].runs[0]
        rpr = run.font._element
        rpr.set("spc", "200")
        result = _extract_char_spacing(run.font)
        assert result == 2.0

    def test_without_spacing(self, blank_slide):
        txBox = blank_slide.shapes.add_textbox(
            Inches(0), Inches(0), Inches(1), Inches(1)
        )
        txBox.text_frame.text = "X"
        run = txBox.text_frame.paragraphs[0].runs[0]
        result = _extract_char_spacing(run.font)
        assert result is None

    def test_negative_spacing(self, sample_textbox):
        run = sample_textbox.text_frame.paragraphs[0].runs[0]
        rpr = run.font._element
        rpr.set("spc", "-100")
        result = _extract_char_spacing(run.font)
        assert result == -1.0


class TestExtractParagraphFont:
    """Tests for extract_paragraph_font."""

    def test_paragraph_level_font(self, sample_textbox):
        para = sample_textbox.text_frame.paragraphs[0]
        para.font.name = "Calibri"
        para.font.size = Pt(16)
        para.font.bold = True
        info = extract_paragraph_font(para)
        assert info["font"] == "Calibri"
        assert info["size"] == 16
        assert info["bold"] is True

    def test_empty_paragraph(self, blank_slide):
        txBox = blank_slide.shapes.add_textbox(
            Inches(0), Inches(0), Inches(1), Inches(1)
        )
        txBox.text_frame.text = ""
        para = txBox.text_frame.paragraphs[0]
        info = extract_paragraph_font(para)
        # Without explicit properties, dict should be empty or minimal
        assert isinstance(info, dict)


class TestExtractAlignment:
    """Tests for extract_alignment."""

    @pytest.mark.parametrize(
        "alignment,expected",
        [
            (PP_ALIGN.LEFT, "left"),
            (PP_ALIGN.CENTER, "center"),
            (PP_ALIGN.RIGHT, "right"),
            (PP_ALIGN.JUSTIFY, "justify"),
        ],
    )
    def test_alignment(self, sample_textbox, alignment, expected):
        para = sample_textbox.text_frame.paragraphs[0]
        para.alignment = alignment
        assert extract_alignment(para) == expected

    def test_none(self, blank_slide):
        txBox = blank_slide.shapes.add_textbox(
            Inches(0), Inches(0), Inches(1), Inches(1)
        )
        txBox.text_frame.text = "X"
        para = txBox.text_frame.paragraphs[0]
        assert extract_alignment(para) is None


class TestConstants:
    """Tests for module constants."""

    def test_alignment_map_keys(self):
        assert set(ALIGNMENT_MAP.keys()) == {"left", "center", "right", "justify"}

    def test_alignment_reverse_map_values(self):
        expected = {"left", "center", "right", "justify"}
        assert set(ALIGNMENT_REVERSE_MAP.values()) == expected

    def test_font_weight_suffixes_not_empty(self):
        assert len(FONT_WEIGHT_SUFFIXES) > 0
        for suffix in FONT_WEIGHT_SUFFIXES:
            assert suffix.startswith(" ")
