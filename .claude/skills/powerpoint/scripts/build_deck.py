# Copyright (c) Microsoft Corporation.
# SPDX-License-Identifier: MIT
"""Build a PowerPoint slide deck from YAML content and style definitions.

Usage::

    python build_deck.py --content-dir content/ \
        --style content/global/style.yaml \
        --output slide-deck/presentation.pptx

    python build_deck.py --content-dir content/ \
        --style content/global/style.yaml \
        --source existing.pptx \
        --output slide-deck/presentation.pptx --slides 3,7,15
"""

import argparse
import ast
import builtins
import importlib.util
import re
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from pptx_charts import add_chart_element
from pptx_colors import apply_color_to_font, resolve_color
from pptx_fills import apply_effect_list, apply_fill, apply_line
from pptx_fonts import ALIGNMENT_MAP
from pptx_shapes import SHAPE_MAP, apply_rotation
from pptx_tables import add_table_element
from pptx_text import (
    SHAPE_KEYS,
    TEXTBOX_KEYS,
    apply_run_properties,
    apply_text_properties,
    populate_text_frame,
)
from pptx_utils import load_yaml

CONNECTOR_TYPE_MAP = {
    "straight": MSO_CONNECTOR_TYPE.STRAIGHT,
    "elbow": MSO_CONNECTOR_TYPE.ELBOW,
    "curve": MSO_CONNECTOR_TYPE.CURVE,
}

PNS = "http://schemas.openxmlformats.org/presentationml/2006/main"
ANS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# Stdlib modules blocked in content-extra.py scripts due to security risk.
# content-extra.py may only import from pptx and safe standard-library modules.
_BLOCKED_STDLIB_MODULES = frozenset(
    {
        "code",
        "codeop",
        "compileall",
        "ctypes",
        "dbm",
        "ensurepip",
        "ftplib",
        "http",
        "imaplib",
        "importlib",
        "marshal",
        "multiprocessing",
        "os",
        "pickle",
        "pkgutil",
        "poplib",
        "py_compile",
        "runpy",
        "shelve",
        "shutil",
        "signal",
        "smtplib",
        "socket",
        "sqlite3",
        "subprocess",
        "sys",
        "telnetlib",
        "tempfile",
        "threading",
        "urllib",
        "venv",
        "webbrowser",
        "xmlrpc",
        "zipimport",
    }
)

_DANGEROUS_BUILTINS = frozenset(
    {
        "__import__",
        "breakpoint",
        "compile",
        "eval",
        "exec",
    }
)

# Builtins that can bypass the import allowlist or execute arbitrary strings
# when called indirectly through attribute access or introspection.
_INDIRECT_BYPASS_BUILTINS = frozenset(
    {
        "delattr",
        "getattr",
        "globals",
        "locals",
        "setattr",
        "vars",
    }
)


class ContentExtraError(Exception):
    """A content-extra.py script failed security validation."""


def _check_module_allowed(
    module_name: str, script_path: Path, stdlib_names: frozenset[str]
) -> None:
    """Raise ContentExtraError if *module_name* is not on the allowlist."""
    top_level = module_name.split(".")[0]

    if top_level == "pptx":
        return

    if top_level in _BLOCKED_STDLIB_MODULES:
        raise ContentExtraError(f"Blocked import '{module_name}' in {script_path}")

    if top_level in stdlib_names:
        return

    raise ContentExtraError(
        f"Disallowed import '{module_name}' in {script_path}: "
        "only pptx and safe standard library modules are permitted"
    )


def _validate_content_extra(script_path: Path) -> None:
    """Validate a content-extra.py script's AST before execution.

    Parses the script and rejects imports outside of pptx and safe stdlib
    modules, as well as calls to dangerous builtins (exec, eval, __import__,
    compile, breakpoint).  Raises ContentExtraError on any violation.
    """
    source = script_path.read_text(encoding="utf-8")
    try:
        tree = ast.parse(source, filename=str(script_path))
    except SyntaxError as exc:
        raise ContentExtraError(f"Syntax error in {script_path}: {exc}") from exc

    stdlib_names = sys.stdlib_module_names

    for node in ast.walk(tree):
        if isinstance(node, ast.Import):
            for alias in node.names:
                _check_module_allowed(alias.name, script_path, stdlib_names)
        elif isinstance(node, ast.ImportFrom):
            if node.module:
                _check_module_allowed(node.module, script_path, stdlib_names)
        elif isinstance(node, ast.Call):
            func = node.func
            if isinstance(func, ast.Name):
                if func.id in _DANGEROUS_BUILTINS:
                    raise ContentExtraError(
                        f"Dangerous builtin '{func.id}' in {script_path}"
                    )
                if func.id in _INDIRECT_BYPASS_BUILTINS:
                    raise ContentExtraError(
                        f"Indirect bypass builtin '{func.id}' in {script_path}"
                    )


def _reset_effect_ref(shape):
    """Reset effectRef idx to 0 to prevent theme shadow inheritance.

    python-pptx defaults effectRef idx to 2, which references the theme's
    effectStyleLst[2] that typically includes an outerShdw element.
    """
    style_el = shape._element.find(f"{{{PNS}}}style")
    if style_el is not None:
        effect_ref = style_el.find(f"{{{ANS}}}effectRef")
        if effect_ref is not None:
            effect_ref.set("idx", "0")


def set_slide_bg(slide, fill_spec, colors: dict):
    """Set a background fill on a slide."""
    apply_fill(slide.background, fill_spec, colors)


def set_slide_bg_image(slide, image_path: str, content_dir: Path):
    """Set a background image on a slide using blipFill in the background element."""
    img_file = content_dir / image_path
    if not img_file.exists():
        return

    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.parts.image import Image, ImagePart

    sld = slide._element
    cSld = sld.find(qn("p:cSld"))
    if cSld is None:
        return

    spTree = cSld.find(qn("p:spTree"))

    # Remove existing p:bg element if present
    existing_bg = cSld.find(qn("p:bg"))
    if existing_bg is not None:
        cSld.remove(existing_bg)

    # Create image part and relate to slide
    image = Image.from_file(str(img_file))
    image_part = ImagePart.new(slide.part.package, image)
    rel = slide.part.relate_to(image_part, RT.IMAGE)

    # Build p:bg > p:bgPr > a:blipFill structure
    bg = etree.SubElement(cSld, qn("p:bg"))
    bgPr = etree.SubElement(bg, qn("p:bgPr"))
    blipFill = etree.SubElement(bgPr, qn("a:blipFill"))
    blipFill.set("dpi", "0")
    blipFill.set("rotWithShape", "1")

    blip = etree.SubElement(blipFill, qn("a:blip"))
    blip.set(qn("r:embed"), rel)

    stretch = etree.SubElement(blipFill, qn("a:stretch"))
    etree.SubElement(stretch, qn("a:fillRect"))

    etree.SubElement(bgPr, qn("a:effectLst"))

    # Ensure p:bg appears before p:spTree (required by schema)
    if spTree is not None:
        cSld.remove(bg)
        cSld.insert(list(cSld).index(spTree), bg)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_name=None,
    font_size=16,
    font_color=None,
    bold=False,
    italic=False,
    alignment=None,
    name=None,
    rotation=None,
    elem=None,
    colors=None,
):
    """Add a text box to a slide with font and layout properties.

    Args:
        slide: Target slide object.
        left: Left position in inches.
        top: Top position in inches.
        width: Width in inches.
        height: Height in inches.
        text: Text content for the box.
        font_name: Font family name.
        font_size: Font size in points.
        font_color: Resolved color spec dict.
        bold: Apply bold formatting.
        italic: Apply italic formatting.
        alignment: Paragraph alignment name.
        name: Shape name identifier.
        rotation: Rotation angle in degrees.
        elem: Full element dict from content.yaml for
            paragraph-level and run-level properties.
        colors: Color resolution dict.

    Returns:
        The created textbox shape object.
    """
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if name:
        txBox.name = name
    apply_rotation(txBox, rotation)

    defaults = {
        "font": font_name,
        "size": font_size,
        "color": font_color,
        "bold": bold,
        "italic": italic,
        "alignment": alignment,
    }
    source = elem or {"text": text}
    if "text" not in source:
        source = {**source, "text": text}
    populate_text_frame(txBox.text_frame, source, colors or {}, TEXTBOX_KEYS, defaults)
    return txBox


def add_shape_element(slide, elem, colors, typography):
    """Add a shape element from a content.yaml definition."""
    shape_type = SHAPE_MAP.get(elem.get("shape", "rectangle"), MSO_SHAPE.RECTANGLE)
    left = Inches(elem["left"])
    top = Inches(elem["top"])
    width = Inches(elem["width"])
    height = Inches(elem["height"])

    shape = slide.shapes.add_shape(shape_type, left, top, width, height)

    _reset_effect_ref(shape)

    if "name" in elem:
        shape.name = elem["name"]

    apply_rotation(shape, elem.get("rotation"))
    apply_fill(shape, elem.get("fill"), colors)
    apply_line(shape, elem, colors)

    if "corner_radius" in elem:
        shape.adjustments[0] = elem["corner_radius"]

    if "effect" in elem:
        apply_effect_list(shape, elem["effect"])

    if "text" in elem:
        populate_text_frame(shape.text_frame, elem, colors, SHAPE_KEYS)

    return shape


def add_image_element(slide, elem, content_dir: Path):
    """Add an image element from a content.yaml definition."""
    img_path = content_dir / elem["path"]
    if not img_path.exists():
        # Fallback: add a text box with the path as placeholder
        add_textbox(
            slide,
            elem["left"],
            elem["top"],
            elem["width"],
            elem["height"],
            f"[Image: {elem['path']}]",
            font_size=12,
        )
        return None

    left = Inches(elem["left"])
    top = Inches(elem["top"])
    width = Inches(elem["width"])
    height = Inches(elem["height"])
    pic = slide.shapes.add_picture(str(img_path), left, top, width, height)
    if "name" in elem:
        pic.name = elem["name"]
    apply_rotation(pic, elem.get("rotation"))

    # Restore blipFill attributes (rotWithShape, dpi, etc.)
    if "blip_fill_attrs" in elem:
        blipFill = pic._element.find(qn("p:blipFill"))
        if blipFill is not None:
            for attr_name, attr_val in elem["blip_fill_attrs"].items():
                blipFill.set(attr_name, attr_val)

    # Apply image crop via srcRect on blipFill
    if "crop" in elem:
        blipFill = pic._element.find(qn("p:blipFill"))
        if blipFill is not None:
            srcRect = blipFill.find(qn("a:srcRect"))
            if srcRect is None:
                # Insert srcRect after a:blip
                blip_el = blipFill.find(qn("a:blip"))
                idx = list(blipFill).index(blip_el) + 1 if blip_el is not None else 0
                srcRect = etree.Element(qn("a:srcRect"))
                blipFill.insert(idx, srcRect)
            crop = elem["crop"]
            for side in ("l", "t", "r", "b"):
                if side in crop:
                    srcRect.set(side, str(crop[side]))

    # Apply image opacity via alphaModFix on the blip element
    if "opacity" in elem:
        blip = pic._element.find(".//" + qn("a:blip"))
        if blip is not None:
            amt = str(int(elem["opacity"] * 1000))
            amf = blip.find(qn("a:alphaModFix"))
            if amf is None:
                amf = etree.SubElement(blip, qn("a:alphaModFix"))
            amf.set("amt", amt)

    return pic


def add_rich_text_element(slide, elem, colors, typography):
    """Add a rich text element with mixed font/color segments."""
    txBox = slide.shapes.add_textbox(
        Inches(elem["left"]),
        Inches(elem["top"]),
        Inches(elem["width"]),
        Inches(elem["height"]),
    )
    if "name" in elem:
        txBox.name = elem["name"]
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]

    # Apply text frame-level properties
    apply_text_properties(tf, elem)

    for i, seg in enumerate(elem.get("segments", [])):
        run = p.add_run() if i > 0 else (p.runs[0] if p.runs else p.add_run())
        run.text = seg["text"]
        seg_font = seg.get("font")
        if seg_font:
            run.font.name = seg_font
        run.font.size = Pt(seg.get("size", 16))
        if "color" in seg:
            color_spec = resolve_color(seg["color"])
            apply_color_to_font(run.font.color, color_spec)
        run.font.bold = seg.get("bold", False)
        run.font.italic = seg.get("italic", False)
        apply_run_properties(run, seg, colors)

    return txBox


def add_card_element(slide, elem, colors, typography):
    """Add a card panel with optional title bar and bullet content."""
    left = Inches(elem["left"])
    top = Inches(elem["top"])
    width = Inches(elem["width"])
    height = Inches(elem["height"])

    # Card background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    apply_fill(shape, elem.get("fill", "#2D2D35"), colors)
    if "border_color" in elem:
        apply_line(
            shape,
            {
                "line_color": elem["border_color"],
                "line_width": elem.get("border_width", 1),
            },
            colors,
        )
    else:
        shape.line.fill.background()

    # Accent bar
    if elem.get("accent_bar"):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(elem["left"] + 0.15),
            Inches(elem["top"] + 0.1),
            Inches(elem["width"] - 0.3),
            Inches(0.04),
        )
        apply_fill(bar, elem.get("accent_color", "#0078D4"), colors)
        bar.line.fill.background()

    # Title
    y_offset = 0.2
    if "title" in elem:
        add_textbox(
            slide,
            elem["left"] + 0.2,
            elem["top"] + y_offset,
            elem["width"] - 0.4,
            0.4,
            elem["title"],
            font_name="Segoe UI",
            font_size=elem.get("title_size", 16),
            font_color=resolve_color(elem.get("title_color", "#F8F8FC")),
            bold=elem.get("title_bold", True),
        )
        y_offset += 0.5

    # Content bullets
    for item in elem.get("content", []):
        bullet_text = (
            f"\u2022 {item['bullet']}" if "bullet" in item else item.get("text", "")
        )
        color = resolve_color(item.get("color", "#F8F8FC"))
        add_textbox(
            slide,
            elem["left"] + 0.2,
            elem["top"] + y_offset,
            elem["width"] - 0.4,
            0.35,
            bullet_text,
            font_name="Segoe UI",
            font_size=item.get("size", 14),
            font_color=color,
        )
        y_offset += 0.35

    return shape


def add_arrow_flow_element(slide, elem, colors, typography):
    """Add a horizontal arrow flow diagram."""
    items = elem.get("items", [])
    if not items:
        return

    total_width = elem["width"]
    item_width = total_width / len(items) - 0.3
    x = elem["left"]

    for item in items:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON,
            Inches(x),
            Inches(elem["top"]),
            Inches(item_width),
            Inches(elem["height"]),
        )
        apply_fill(shape, item.get("color", "#0078D4"), colors)
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item["label"]
        p.alignment = ALIGNMENT_MAP["center"]
        run = p.runs[0]
        run.font.name = "Segoe UI"
        run.font.size = Pt(14)
        apply_color_to_font(run.font.color, resolve_color("#F8F8FC"))
        run.font.bold = True

        x += item_width + 0.3


def add_numbered_step_element(slide, elem, colors, typography):
    """Add a numbered step with circle, label, and description."""
    number = elem.get("number", 1)

    # Number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(elem["left"]),
        Inches(elem["top"]),
        Inches(0.5),
        Inches(0.5),
    )
    apply_fill(circle, elem.get("accent_color", "#0078D4"), colors)
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.alignment = ALIGNMENT_MAP["center"]
    run = p.runs[0]
    run.font.name = "Segoe UI"
    run.font.size = Pt(16)
    apply_color_to_font(run.font.color, resolve_color("#F8F8FC"))
    run.font.bold = True

    # Label
    add_textbox(
        slide,
        elem["left"] + 0.6,
        elem["top"],
        elem["width"] - 0.6,
        0.35,
        elem["label"],
        font_name="Segoe UI",
        font_size=16,
        font_color=resolve_color("#F8F8FC"),
        bold=True,
    )

    # Description
    if "description" in elem:
        add_textbox(
            slide,
            elem["left"] + 0.6,
            elem["top"] + 0.35,
            elem["width"] - 0.6,
            0.4,
            elem["description"],
            font_name="Segoe UI",
            font_size=14,
            font_color=resolve_color("#9CA3AF"),
        )


def add_connector_element(slide, elem: dict, colors: dict):
    """Add a connector element from a content.yaml definition.

    YAML schema:
    - type: connector
      connector_type: straight
      begin_x: 3.0
      begin_y: 2.0
      end_x: 7.0
      end_y: 4.0
      line_color: "#0078D4"
      line_width: 2
      dash_style: solid
      head_end: none
      tail_end: arrow
    """
    conn_type = CONNECTOR_TYPE_MAP.get(
        elem.get("connector_type", "straight"), MSO_CONNECTOR_TYPE.STRAIGHT
    )

    connector = slide.shapes.add_connector(
        conn_type,
        Inches(elem["begin_x"]),
        Inches(elem["begin_y"]),
        Inches(elem["end_x"]),
        Inches(elem["end_y"]),
    )

    apply_line(connector, elem, colors)

    # Arrow heads via lxml XML manipulation
    sp_pr = connector._element.find(qn("a:ln"))
    if sp_pr is None:
        ln_parent = connector._element.spPr
        sp_pr = ln_parent.find(qn("a:ln"))
    if sp_pr is None:
        sp_pr = etree.SubElement(connector._element.spPr, qn("a:ln"))

    if "head_end" in elem and elem["head_end"] != "none":
        head = etree.SubElement(sp_pr, qn("a:headEnd"))
        head.set("type", elem["head_end"])
    if "tail_end" in elem and elem["tail_end"] != "none":
        tail = etree.SubElement(sp_pr, qn("a:tailEnd"))
        tail.set("type", elem["tail_end"])

    if "name" in elem:
        connector.name = elem["name"]

    return connector


MAX_GROUP_DEPTH = 20


def add_group_element(
    slide,
    elem: dict,
    colors: dict,
    typography: dict,
    content_dir: Path,
    *,
    _depth: int = 0,
    max_depth: int = MAX_GROUP_DEPTH,
):
    """Add a group element containing nested child elements.

    Raises ValueError when nesting exceeds *max_depth*.

    YAML schema:
    - type: group
      left: 1.0
      top: 2.0
      width: 5.0
      height: 3.0
      elements:
        - type: shape
          shape: rectangle
          left: 0
          top: 0
          width: 5.0
          height: 3.0
          fill: "#2D2D35"
        - type: textbox
          left: 0.2
          top: 0.2
          width: 4.6
          height: 0.5
          text: "Group Title"
    """
    if _depth >= max_depth:
        raise ValueError(f"Group nesting depth {_depth} exceeds limit of {max_depth}")
    group = slide.shapes.add_group_shape()

    group.left = Inches(elem["left"])
    group.top = Inches(elem["top"])
    group.width = Inches(elem["width"])
    group.height = Inches(elem["height"])

    for child_elem in elem.get("elements", []):
        build_element_in_group(
            group,
            child_elem,
            colors,
            typography,
            content_dir,
            _depth=_depth + 1,
            max_depth=max_depth,
        )

    if "name" in elem:
        group.name = elem["name"]

    return group


def build_element_in_group(
    group,
    elem: dict,
    colors: dict,
    typography: dict,
    content_dir: Path,
    *,
    _depth: int = 0,
    max_depth: int = MAX_GROUP_DEPTH,
):
    """Dispatch a child element build within a group shape.

    Reuses top-level builders for shape and textbox. Groups do not support
    table or chart elements.
    """
    elem_type = elem.get("type", "textbox")

    if elem_type == "shape":
        _add_shape_to_collection(group.shapes, elem, colors)
    elif elem_type == "textbox":
        _add_textbox_to_collection(group.shapes, elem, colors)
    elif elem_type == "connector":
        add_connector_element(group, elem, colors)
    elif elem_type == "image":
        add_image_element(group, elem, content_dir)
    elif elem_type == "group":
        add_group_element(
            group,
            elem,
            colors,
            typography,
            content_dir,
            _depth=_depth,
            max_depth=max_depth,
        )


def _add_shape_to_collection(shapes, elem: dict, colors: dict):
    """Add a shape to any shapes collection (slide or group)."""
    shape_type = SHAPE_MAP.get(elem.get("shape", "rectangle"), MSO_SHAPE.RECTANGLE)
    shape = shapes.add_shape(
        shape_type,
        Inches(elem["left"]),
        Inches(elem["top"]),
        Inches(elem["width"]),
        Inches(elem["height"]),
    )
    if "name" in elem:
        shape.name = elem["name"]
    apply_rotation(shape, elem.get("rotation"))
    apply_fill(shape, elem.get("fill"), colors)
    apply_line(shape, elem, colors)
    if "text" in elem:
        populate_text_frame(shape.text_frame, elem, colors, SHAPE_KEYS)
    return shape


def _add_textbox_to_collection(shapes, elem: dict, colors: dict):
    """Add a textbox to any shapes collection (slide or group)."""
    txBox = shapes.add_textbox(
        Inches(elem["left"]),
        Inches(elem["top"]),
        Inches(elem["width"]),
        Inches(elem["height"]),
    )
    if "name" in elem:
        txBox.name = elem["name"]
    populate_text_frame(txBox.text_frame, elem, colors, TEXTBOX_KEYS)
    return txBox


def _build_textbox_element(slide, elem, colors, typography, content_dir):
    """Build a textbox element with full parameter resolution for YAML keys."""
    font_name = elem.get("font")
    font_color = resolve_color(elem["font_color"]) if "font_color" in elem else None
    is_bold = elem.get("font_bold", elem.get("bold", False))
    add_textbox(
        slide,
        elem["left"],
        elem["top"],
        elem["width"],
        elem["height"],
        elem.get("text", ""),
        font_name=font_name,
        font_size=elem.get("font_size", 16),
        font_color=font_color,
        bold=is_bold,
        italic=elem.get("italic", False),
        alignment=elem.get("alignment"),
        name=elem.get("name"),
        rotation=elem.get("rotation"),
        elem=elem,
        colors=colors,
    )


def _build_image_element(slide, elem, colors, typography, content_dir):
    """Delegate image element building to add_image_element."""
    add_image_element(slide, elem, content_dir)


def _build_group_element(slide, elem, colors, typography, content_dir):
    """Delegate group element building to add_group_element."""
    add_group_element(slide, elem, colors, typography, content_dir, _depth=0)


def _build_connector_element(slide, elem, colors, typography, content_dir):
    """Delegate connector building to add_connector_element."""
    add_connector_element(slide, elem, colors)


def _build_chart_element(slide, elem, colors, typography, content_dir):
    """Delegate chart building to add_chart_element."""
    add_chart_element(slide, elem, colors)


def _build_table_element(slide, elem, colors, typography, content_dir):
    """Delegate table building to add_table_element."""
    add_table_element(slide, elem, colors, typography)


# Element builder registry: maps element type names to builder functions.
# All builders share the signature (slide, elem, colors, typography, content_dir).
ELEMENT_BUILDERS = {
    "shape": lambda slide, elem, colors, typography, content_dir: add_shape_element(
        slide, elem, colors, typography
    ),
    "textbox": _build_textbox_element,
    "image": _build_image_element,
    "rich_text": lambda slide, elem, colors, typography, content_dir: (
        add_rich_text_element(slide, elem, colors, typography)
    ),
    "card": lambda slide, elem, colors, typography, content_dir: add_card_element(
        slide, elem, colors, typography
    ),
    "arrow_flow": lambda slide, elem, colors, typography, content_dir: (
        add_arrow_flow_element(slide, elem, colors, typography)
    ),
    "numbered_step": lambda slide, elem, colors, typography, content_dir: (
        add_numbered_step_element(slide, elem, colors, typography)
    ),
    "table": _build_table_element,
    "chart": _build_chart_element,
    "connector": _build_connector_element,
    "group": _build_group_element,
}


def _build_element(
    slide, elem: dict, colors: dict, typography: dict, content_dir: Path
):
    """Dispatch element building via registry lookup."""
    elem_type = elem.get("type", "textbox")
    builder = ELEMENT_BUILDERS.get(elem_type)
    if builder:
        builder(slide, elem, colors, typography, content_dir)


def clear_slide_shapes(slide):
    """Remove all shapes from a slide, preserving the slide itself."""
    sp_tree = slide.shapes._spTree
    shapes_to_remove = [
        sp
        for sp in sp_tree.iterchildren()
        if sp.tag.endswith("}sp")
        or sp.tag.endswith("}pic")
        or sp.tag.endswith("}grpSp")
        or sp.tag.endswith("}cxnSp")
    ]
    for sp in shapes_to_remove:
        sp_tree.remove(sp)


def _all_layouts(prs):
    """Iterate layouts across all slide masters."""
    for master in prs.slide_masters:
        yield from master.slide_layouts


def _find_blank_layout(prs):
    """Find the best blank layout in the presentation, with fallbacks."""
    # Try index 6 first (default blank in standard templates)
    try:
        return prs.slide_layouts[6]
    except IndexError:
        pass
    # Search by name across all masters
    for layout in _all_layouts(prs):
        if layout.name.lower() in ("blank", "blank slide"):
            return layout
    # Fall back to last layout of first master
    return prs.slide_layouts[len(prs.slide_layouts) - 1]


def get_slide_layout(prs, slide_content: dict, style: dict):
    """Select slide layout based on content.yaml or style.yaml configuration."""
    layout_spec = slide_content.get("layout")
    layouts_map = style.get("layouts", {})

    if layout_spec is None or layout_spec == "blank":
        return _find_blank_layout(prs)

    # Resolve through style.yaml layouts map
    if layout_spec in layouts_map:
        layout_ref = layouts_map[layout_spec]
        if isinstance(layout_ref, int):
            try:
                return prs.slide_layouts[layout_ref]
            except IndexError:
                return _find_blank_layout(prs)
        elif isinstance(layout_ref, str):
            for layout in _all_layouts(prs):
                if layout.name == layout_ref:
                    return layout

    # Direct name lookup across all slide masters
    if isinstance(layout_spec, str):
        for layout in _all_layouts(prs):
            if layout.name == layout_spec:
                return layout

    # Direct index lookup
    if isinstance(layout_spec, int):
        try:
            return prs.slide_layouts[layout_spec]
        except IndexError:
            return _find_blank_layout(prs)

    # Fallback to blank
    return _find_blank_layout(prs)


def build_slide(
    prs,
    slide_content: dict,
    style: dict,
    content_dir: Path,
    existing_slide=None,
    *,
    allow_scripts: bool = False,
):
    """Build a single slide from content.yaml data and style context.

    When existing_slide is provided, clears its shapes and rebuilds in place
    instead of appending a new slide.  Set *allow_scripts* to skip AST
    validation of content-extra.py (use only with trusted content).
    """
    colors = {}
    typography = {}

    if existing_slide is not None:
        slide = existing_slide
        clear_slide_shapes(slide)
    else:
        layout = get_slide_layout(prs, slide_content, style)
        slide = prs.slides.add_slide(layout)

    # Populate themed layout placeholders
    placeholders = slide_content.get("placeholders", {})
    for idx_str, value in placeholders.items():
        idx = int(idx_str)
        if idx in slide.placeholders:
            ph = slide.placeholders[idx]
            if isinstance(value, str):
                ph.text = value
            elif isinstance(value, list):
                tf = ph.text_frame
                tf.text = value[0]
                for line in value[1:]:
                    tf.add_paragraph().text = line

    # Remove unused placeholder shapes inherited from the layout
    used_ph_indices = {int(k) for k in placeholders}
    sp_tree = slide.shapes._spTree
    for sp in list(sp_tree.iterchildren()):
        nvSpPr = sp.find(qn("p:nvSpPr"))
        if nvSpPr is None:
            continue
        nvPr = nvSpPr.find(qn("p:nvPr"))
        if nvPr is None:
            continue
        ph = nvPr.find(qn("p:ph"))
        if ph is not None:
            idx = int(ph.get("idx", "0"))
            if idx not in used_ph_indices:
                sp_tree.remove(sp)

    # Set background from per-slide definition only
    bg_block = slide_content.get("background")
    if bg_block and "image" in bg_block:
        set_slide_bg_image(slide, bg_block["image"], content_dir)
    elif bg_block and "fill" in bg_block:
        set_slide_bg(slide, bg_block["fill"], colors)

    # Sort elements by z_order to preserve stacking order
    elements = slide_content.get("elements", [])
    elements = sorted(elements, key=lambda e: e.get("z_order", 0))

    # Filter out empty placeholder elements
    elements = [
        e
        for e in elements
        if not (e.get("_placeholder") and not e.get("text", "").strip())
    ]

    turbo_enabled = len(elements) > 20
    if turbo_enabled:
        slide.shapes.turbo_add_enabled = True

    # Process elements in order
    for elem in elements:
        _build_element(slide, elem, colors, typography, content_dir)

    # Execute content-extra.py if present (validated before loading)
    extra_script = content_dir / "content-extra.py"
    if extra_script.exists():
        if not allow_scripts:
            _validate_content_extra(extra_script)
        spec = importlib.util.spec_from_file_location(
            "content_extra", str(extra_script)
        )
        mod = importlib.util.module_from_spec(spec)
        if not allow_scripts:
            # __import__ is kept because the import machinery needs it;
            # the AST checker already blocks direct __import__() calls.
            stripped = (_DANGEROUS_BUILTINS | _INDIRECT_BYPASS_BUILTINS) - {
                "__import__"
            }
            safe_builtins = {
                k: v for k, v in builtins.__dict__.items() if k not in stripped
            }
            mod.__builtins__ = safe_builtins
        spec.loader.exec_module(mod)
        if hasattr(mod, "render"):
            mod.render(slide, style, content_dir)

    if turbo_enabled:
        slide.shapes.turbo_add_enabled = False

    # Add speaker notes (preserve empty strings when notes slide exists)
    notes = slide_content.get("speaker_notes")
    if notes is not None:
        notes_slide = slide.notes_slide
        notes_text = re.sub(r"\v", "\n", notes) if notes else ""
        notes_slide.notes_text_frame.text = notes_text

    return slide


def discover_slides(content_dir: Path) -> list[tuple[int, Path]]:
    """Discover slide content directories and return sorted (number, path) pairs."""
    slides = []
    for child in content_dir.iterdir():
        if child.is_dir() and child.name.startswith("slide-"):
            match = re.match(r"slide-(\d+)", child.name)
            if match:
                num = int(match.group(1))
                content_yaml = child / "content.yaml"
                if content_yaml.exists():
                    slides.append((num, child))
    return sorted(slides, key=lambda x: x[0])


def main():
    """CLI entry point for building a PowerPoint deck from YAML."""
    parser = argparse.ArgumentParser(
        description="Build a PowerPoint deck from YAML content"
    )
    parser.add_argument(
        "--content-dir", required=True, help="Path to the content/ directory"
    )
    parser.add_argument("--style", required=True, help="Path to the global style.yaml")
    parser.add_argument("--output", required=True, help="Output PPTX file path")
    parser.add_argument("--template", help="Template PPTX file path for themed builds")
    parser.add_argument("--source", help="Source PPTX to update (for partial rebuilds)")
    parser.add_argument(
        "--slides", help="Comma-separated slide numbers to rebuild (requires --source)"
    )
    parser.add_argument(
        "--allow-scripts",
        action="store_true",
        help="Skip AST validation of content-extra.py (trusted content only)",
    )
    args = parser.parse_args()

    content_dir = Path(args.content_dir)
    style = load_yaml(Path(args.style))
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    dims = style.get("dimensions", {})
    width = dims.get("width_inches", 13.333)
    height = dims.get("height_inches", 7.5)

    if args.template:
        # Template build: open template and preserve its theme/layouts
        prs = Presentation(args.template)
        # Only override dimensions when explicitly set in style.yaml
        if "dimensions" in style:
            prs.slide_width = Inches(width)
            prs.slide_height = Inches(height)

        # Remove existing slides from the template — keep only theme/layouts
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

        # Apply presentation metadata from style.yaml
        metadata = style.get("metadata", {})
        if metadata:
            props = prs.core_properties
            for key, value in metadata.items():
                if hasattr(props, key):
                    setattr(props, key, value)

        slides_data = discover_slides(content_dir)
        if not slides_data:
            print("No slide content found in", content_dir)
            sys.exit(1)

        for num, slide_dir in slides_data:
            slide_content = load_yaml(slide_dir / "content.yaml")
            build_slide(
                prs,
                slide_content,
                style,
                slide_dir,
                allow_scripts=args.allow_scripts,
            )
            print(f"Built slide {num}: {slide_content.get('title', 'Untitled')}")
    elif args.source and args.slides:
        # Partial rebuild: open existing deck and replace specific slides
        prs = Presentation(args.source)
        slide_nums = [int(s.strip()) for s in args.slides.split(",")]
        slides_data = discover_slides(content_dir)
        slides_to_rebuild = {
            num: path for num, path in slides_data if num in slide_nums
        }

        for num in slide_nums:
            if num not in slides_to_rebuild:
                print(f"Warning: No content found for slide {num}, skipping")
                continue
            slide_dir = slides_to_rebuild[num]
            slide_content = load_yaml(slide_dir / "content.yaml")
            # Rebuild in-place: clear shapes on the existing slide and repopulate
            idx = num - 1
            if idx < len(prs.slides):
                existing_slide = prs.slides[idx]
                build_slide(
                    prs,
                    slide_content,
                    style,
                    slide_dir,
                    existing_slide=existing_slide,
                    allow_scripts=args.allow_scripts,
                )
                print(f"Rebuilt slide {num} in-place")
            else:
                slide_count = len(prs.slides)
                print(
                    f"Warning: Slide {num} does not exist"
                    f" in deck (has {slide_count} slides),"
                    f" skipping"
                )
    else:
        # Full build
        prs = Presentation()
        prs.slide_width = Inches(width)
        prs.slide_height = Inches(height)

        # Apply presentation metadata from style.yaml
        metadata = style.get("metadata", {})
        if metadata:
            props = prs.core_properties
            for key, value in metadata.items():
                if hasattr(props, key):
                    setattr(props, key, value)

        slides_data = discover_slides(content_dir)
        if not slides_data:
            print("No slide content found in", content_dir)
            sys.exit(1)

        for num, slide_dir in slides_data:
            slide_content = load_yaml(slide_dir / "content.yaml")
            build_slide(
                prs,
                slide_content,
                style,
                slide_dir,
                allow_scripts=args.allow_scripts,
            )
            print(f"Built slide {num}: {slide_content.get('title', 'Untitled')}")

    prs.save(str(output_path))
    print(f"\nDeck saved to {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
