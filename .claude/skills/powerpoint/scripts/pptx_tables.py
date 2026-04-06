# Copyright (c) Microsoft Corporation.
# SPDX-License-Identifier: MIT
"""Table build and extract utilities for PowerPoint skill scripts.

Provides add_table_element() for building tables from YAML definitions
and extract_table() for extracting table data from existing presentations.
"""

from pptx.util import Inches, Pt
from pptx_colors import apply_color_to_font, resolve_color
from pptx_fills import apply_fill, extract_fill
from pptx_fonts import extract_font_info
from pptx_text import VERTICAL_ANCHOR_MAP, VERTICAL_ANCHOR_REVERSE
from pptx_utils import emu_to_inches


def add_table_element(slide, elem: dict, colors: dict, typography: dict):
    """Add a table element from a content.yaml definition.

    YAML schema:
    - type: table
      left: 1.0
      top: 2.0
      width: 10.0
      height: 3.0
      first_row: true
      horz_banding: true
      columns:
        - width: 2.5
        - width: 3.75
      rows:
        - cells:
            - text: "Header"
              fill: "#0078D4"
              font_color: "#F8F8FC"
              font_bold: true
              merge_right: 2
    """
    rows_data = elem.get("rows", [])
    cols_data = elem.get("columns", [])
    n_rows = len(rows_data)
    n_cols = (
        len(cols_data)
        if cols_data
        else max((len(r.get("cells", [])) for r in rows_data), default=1)
    )

    table_shape = slide.shapes.add_table(
        n_rows,
        n_cols,
        Inches(elem["left"]),
        Inches(elem["top"]),
        Inches(elem["width"]),
        Inches(elem["height"]),
    )
    table = table_shape.table

    # Table properties
    if "first_row" in elem:
        table.first_row = elem["first_row"]
    if "last_row" in elem:
        table.last_row = elem["last_row"]
    if "first_col" in elem:
        table.first_col = elem["first_col"]
    if "last_col" in elem:
        table.last_col = elem["last_col"]
    if "horz_banding" in elem:
        table.horz_banding = elem["horz_banding"]
    if "vert_banding" in elem:
        table.vert_banding = elem["vert_banding"]

    # Column widths
    for i, col_spec in enumerate(cols_data):
        if i < n_cols and "width" in col_spec:
            table.columns[i].width = Inches(col_spec["width"])

    # Cell population
    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_data in enumerate(row_data.get("cells", [])):
            if col_idx >= n_cols:
                break
            cell = table.cell(row_idx, col_idx)

            # Handle cell merging
            if "merge_right" in cell_data and cell_data["merge_right"] > 0:
                merge_target = table.cell(row_idx, col_idx + cell_data["merge_right"])
                cell.merge(merge_target)
            if "merge_down" in cell_data and cell_data["merge_down"] > 0:
                merge_target = table.cell(row_idx + cell_data["merge_down"], col_idx)
                cell.merge(merge_target)

            # Set text
            text = cell_data.get("text", "")
            if text:
                cell.text = str(text)

            # Cell fill
            if "fill" in cell_data:
                apply_fill(cell, cell_data["fill"], colors)

            # Cell text formatting
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if "font_color" in cell_data:
                        color_spec = resolve_color(cell_data["font_color"], colors)
                        apply_color_to_font(run.font.color, color_spec)
                    if cell_data.get("font_bold"):
                        run.font.bold = True
                    if "font_size" in cell_data:
                        run.font.size = Pt(cell_data["font_size"])
                    if "font" in cell_data:
                        run.font.name = cell_data["font"]

            # Vertical anchor
            if "vertical_anchor" in cell_data:
                anchor = VERTICAL_ANCHOR_MAP.get(cell_data["vertical_anchor"])
                if anchor is not None:
                    cell.vertical_anchor = anchor

    if "name" in elem:
        table_shape.name = elem["name"]

    return table_shape


def extract_table(shape, colors: dict | None = None) -> dict:
    """Extract a table element definition from a GraphicFrame shape."""
    table = shape.table
    elem = {
        "type": "table",
        "left": emu_to_inches(shape.left),
        "top": emu_to_inches(shape.top),
        "width": emu_to_inches(shape.width),
        "height": emu_to_inches(shape.height),
        "name": shape.name,
    }

    # Table properties
    if table.first_row:
        elem["first_row"] = True
    if table.last_row:
        elem["last_row"] = True
    if table.first_col:
        elem["first_col"] = True
    if table.last_col:
        elem["last_col"] = True
    if table.horz_banding:
        elem["horz_banding"] = True
    if table.vert_banding:
        elem["vert_banding"] = True

    # Column widths
    elem["columns"] = [{"width": emu_to_inches(col.width)} for col in table.columns]

    # Rows and cells
    elem["rows"] = []
    for row in table.rows:
        row_data = {"cells": []}
        for cell in row.cells:
            cell_data = {"text": cell.text}

            # Cell fill
            try:
                cell_fill = extract_fill(cell.fill)
                if cell_fill:
                    cell_data["fill"] = cell_fill
            except (AttributeError, TypeError):
                pass

            # Merge info
            if cell.is_merge_origin:
                if cell.span_width > 1:
                    cell_data["merge_right"] = cell.span_width - 1
                if cell.span_height > 1:
                    cell_data["merge_down"] = cell.span_height - 1
            elif cell.is_spanned:
                cell_data["_spanned"] = True

            # Vertical anchor
            if cell.vertical_anchor is not None:
                label = VERTICAL_ANCHOR_REVERSE.get(cell.vertical_anchor)
                if label:
                    cell_data["vertical_anchor"] = label

            # Font info from first run
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    font_info = extract_font_info(run.font)
                    if font_info.get("bold"):
                        cell_data["font_bold"] = True
                    if "color" in font_info:
                        cell_data["font_color"] = font_info["color"]
                    if "size" in font_info:
                        cell_data["font_size"] = font_info["size"]
                    if "font" in font_info:
                        cell_data["font"] = font_info["font"]
                    break
                break

            row_data["cells"].append(cell_data)
        elem["rows"].append(row_data)

    return elem
