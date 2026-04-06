# Copyright (c) Microsoft Corporation.
# SPDX-License-Identifier: MIT
"""Tests for extract_content module."""

from collections import Counter
from pathlib import Path
from unittest.mock import MagicMock

import pytest
from extract_content import (
    _build_color_map,
    _classify_slide_brightness,
    _cluster_themes,
    _convert_svg_to_png,
    _has_formatting_variation,
    _ImageSecurityError,
    _is_background_image,
    _is_freeform,
    _resolve_theme_colors,
    _resolve_theme_refs_in_content,
    _sanitize_svg,
    _save_image_blob,
    _validate_emf_magic_bytes,
    detect_global_style,
    extract_child_shape,
    extract_connector,
    extract_freeform,
    extract_group,
    extract_image,
    extract_shape,
    extract_slide,
    extract_textbox,
)
from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


class TestHasFormattingVariation:
    """Tests for _has_formatting_variation."""

    def test_single_run_no_variation(self):
        assert _has_formatting_variation([{"text": "hi"}]) is False

    def test_empty_list(self):
        assert _has_formatting_variation([]) is False

    def test_same_formatting(self):
        runs = [
            {"text": "a", "font": "Arial", "bold": False},
            {"text": "b", "font": "Arial", "bold": False},
        ]
        assert _has_formatting_variation(runs) is False

    @pytest.mark.parametrize(
        "diff_key,val_a,val_b",
        [
            ("font", "Arial", "Calibri"),
            ("size", 12, 16),
            ("color", "#000000", "#FF0000"),
            ("bold", True, False),
            ("italic", True, False),
            ("underline", True, False),
        ],
    )
    def test_variation_detected(self, diff_key, val_a, val_b):
        runs = [
            {"text": "a", diff_key: val_a},
            {"text": "b", diff_key: val_b},
        ]
        assert _has_formatting_variation(runs) is True


class TestExtractConnector:
    """Tests for extract_connector."""

    def test_basic_connector(self, blank_slide):
        from pptx.enum.shapes import MSO_CONNECTOR_TYPE

        connector = blank_slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            Inches(1),
            Inches(2),
            Inches(5),
            Inches(4),
        )
        result = extract_connector(connector)
        assert result["type"] == "connector"
        assert result["begin_x"] == pytest.approx(1.0, abs=0.01)
        assert result["begin_y"] == pytest.approx(2.0, abs=0.01)
        assert result["end_x"] == pytest.approx(5.0, abs=0.01)
        assert result["end_y"] == pytest.approx(4.0, abs=0.01)
        assert "name" in result


class TestIsFreeform:
    """Tests for _is_freeform."""

    def test_rectangle_is_not_freeform(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(2),
            Inches(2),
        )
        assert _is_freeform(shape) is False

    def test_oval_is_not_freeform(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1),
            Inches(1),
            Inches(2),
            Inches(2),
        )
        assert _is_freeform(shape) is False


class TestIsBackgroundImage:
    """Tests for _is_background_image."""

    @pytest.mark.parametrize(
        "w_factor,h_factor,expected",
        [
            (1.0, 1.0, True),
            (0.375, 0.4, False),
            (0.96, 0.96, True),
        ],
    )
    def test_coverage(self, w_factor, h_factor, expected):
        shape = MagicMock()
        shape.width = Inches(13.333 * w_factor)
        shape.height = Inches(7.5 * h_factor)
        assert _is_background_image(shape, 13.333, 7.5) is expected


class TestExtractShape:
    """Tests for extract_shape."""

    def test_basic_shape(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(2),
            Inches(3),
            Inches(4),
            Inches(2),
        )
        result = extract_shape(shape)
        assert result["type"] == "shape"
        assert result["shape"] == "rectangle"
        assert result["left"] == pytest.approx(2.0, abs=0.01)
        assert result["top"] == pytest.approx(3.0, abs=0.01)
        assert result["width"] == pytest.approx(4.0, abs=0.01)
        assert result["height"] == pytest.approx(2.0, abs=0.01)
        assert "name" in result

    def test_shape_with_text(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(3),
            Inches(2),
        )
        tf = shape.text_frame
        tf.text = "Hello"
        result = extract_shape(shape)
        assert result["text"] == "Hello"

    def test_shape_with_fill(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(3),
            Inches(2),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x00, 0x78, 0xD4)
        result = extract_shape(shape)
        assert result.get("fill") == "#0078D4"

    def test_oval_shape_type(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1),
            Inches(1),
            Inches(2),
            Inches(2),
        )
        result = extract_shape(shape)
        assert result["shape"] == "oval"

    def test_shape_multi_paragraph(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(3),
        )
        tf = shape.text_frame
        tf.text = "Line 1"
        tf.add_paragraph().text = "Line 2"
        result = extract_shape(shape)
        assert "paragraphs" in result
        assert len(result["paragraphs"]) == 2

    def test_rounded_rectangle_corner_radius(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(3),
            Inches(2),
        )
        result = extract_shape(shape)
        # Rounded rectangles have adjustments for corner radius
        assert result["shape"] == "rounded_rectangle"


class TestExtractTextbox:
    """Tests for extract_textbox."""

    def test_basic_textbox(self, blank_slide):
        txBox = blank_slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(1),
        )
        txBox.text_frame.text = "Sample"
        result = extract_textbox(txBox)
        assert result["type"] == "textbox"
        assert result["text"] == "Sample"
        assert result["left"] == pytest.approx(1.0, abs=0.01)

    def test_textbox_with_formatting(self, blank_slide):
        txBox = blank_slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(1),
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Bold text"
        run.font.bold = True
        run.font.size = Pt(24)
        run.font.name = "Arial"
        result = extract_textbox(txBox)
        assert result.get("font_bold") is True or (
            "paragraphs" in result and result["paragraphs"][0].get("font_bold") is True
        )

    def test_textbox_multi_paragraph(self, blank_slide):
        txBox = blank_slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(2),
        )
        tf = txBox.text_frame
        tf.text = "Para 1"
        tf.add_paragraph().text = "Para 2"
        result = extract_textbox(txBox)
        assert "paragraphs" in result
        assert len(result["paragraphs"]) == 2

    def test_textbox_empty(self, blank_slide):
        txBox = blank_slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(1),
        )
        result = extract_textbox(txBox)
        assert result["text"] == ""

    def test_textbox_rich_text_detection(self, blank_slide):
        """Rich text detected when runs have different formatting."""
        txBox = blank_slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(1),
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = "Bold"
        run1.font.bold = True
        run2 = p.add_run()
        run2.text = " Normal"
        run2.font.bold = False
        result = extract_textbox(txBox)
        # Should detect runs variation and include "runs" in paragraph
        if "paragraphs" in result:
            assert any("runs" in pd for pd in result["paragraphs"])


class TestExtractImage:
    """Tests for extract_image."""

    def test_extract_embedded_image(self, blank_slide, sample_image_path, tmp_path):
        pic = blank_slide.shapes.add_picture(
            str(sample_image_path),
            Inches(1),
            Inches(1),
            Inches(3),
            Inches(2),
        )
        output_dir = tmp_path / "output"
        output_dir.mkdir()
        result = extract_image(pic, output_dir, 1, 1)
        assert result["type"] == "image"
        assert "path" in result
        assert result["left"] == pytest.approx(1.0, abs=0.01)

    def test_extract_image_position(self, blank_slide, sample_image_path, tmp_path):
        pic = blank_slide.shapes.add_picture(
            str(sample_image_path),
            Inches(2.5),
            Inches(3.0),
            Inches(4.0),
            Inches(2.5),
        )
        output_dir = tmp_path / "output"
        output_dir.mkdir()
        result = extract_image(pic, output_dir, 1, 1)
        assert result["left"] == pytest.approx(2.5, abs=0.01)
        assert result["top"] == pytest.approx(3.0, abs=0.01)
        assert result["width"] == pytest.approx(4.0, abs=0.01)
        assert result["height"] == pytest.approx(2.5, abs=0.01)

    def test_extract_image_saves_file(self, blank_slide, sample_image_path, tmp_path):
        pic = blank_slide.shapes.add_picture(
            str(sample_image_path),
            Inches(1),
            Inches(1),
            Inches(3),
            Inches(2),
        )
        output_dir = tmp_path / "output"
        output_dir.mkdir()
        result = extract_image(pic, output_dir, 1, 1)
        img_file = output_dir / result["path"]
        assert img_file.exists()

    def test_extract_image_delegates_to_save_image_blob(
        self, blank_slide, sample_image_path, tmp_path, mocker
    ):
        pic = blank_slide.shapes.add_picture(
            str(sample_image_path),
            Inches(1),
            Inches(1),
            Inches(3),
            Inches(2),
        )
        output_dir = tmp_path / "output"
        output_dir.mkdir()
        mock_save = mocker.patch(
            "extract_content._save_image_blob",
            return_value={"path": "images/image-01.png"},
        )
        result = extract_image(pic, output_dir, 2, 5)
        mock_save.assert_called_once_with(pic, output_dir, 2, 5)
        assert result["path"] == "images/image-01.png"


class TestExtractSlide:
    """Tests for extract_slide."""

    def test_empty_slide(self, blank_slide, tmp_path):
        content, slide_dir = extract_slide(blank_slide, 1, tmp_path)
        assert content["slide"] == 1
        assert content["elements"] == []

    def test_slide_with_textbox(self, blank_slide, tmp_path):
        txBox = blank_slide.shapes.add_textbox(
            Inches(1),
            Inches(0.5),
            Inches(4),
            Inches(1),
        )
        txBox.text_frame.text = "Slide Title"
        content, slide_dir = extract_slide(blank_slide, 1, tmp_path)
        assert len(content["elements"]) == 1
        assert content["elements"][0]["type"] == "textbox"
        # Title detection: textbox near top with short text
        assert content["title"] == "Slide Title"

    def test_slide_with_shape(self, blank_slide, tmp_path):
        blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(2),
            Inches(2),
            Inches(3),
            Inches(2),
        )
        content, slide_dir = extract_slide(blank_slide, 1, tmp_path)
        shape_elems = [e for e in content["elements"] if e["type"] == "shape"]
        assert len(shape_elems) == 1

    def test_slide_with_image(self, blank_slide, sample_image_path, tmp_path):
        blank_slide.shapes.add_picture(
            str(sample_image_path),
            Inches(1),
            Inches(1),
            Inches(3),
            Inches(2),
        )
        content, slide_dir = extract_slide(blank_slide, 1, tmp_path)
        img_elems = [e for e in content["elements"] if e["type"] == "image"]
        assert len(img_elems) == 1

    def test_slide_z_order(self, blank_slide, tmp_path):
        """Elements preserve z_order from slide shape ordering."""
        blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(2),
            Inches(2),
        )
        blank_slide.shapes.add_textbox(
            Inches(3),
            Inches(1),
            Inches(2),
            Inches(1),
        )
        content, _ = extract_slide(blank_slide, 1, tmp_path)
        for elem in content["elements"]:
            assert "z_order" in elem

    def test_slide_speaker_notes(self, blank_slide, tmp_path):
        notes_slide = blank_slide.notes_slide
        notes_slide.notes_text_frame.text = "Speaker notes here"
        content, _ = extract_slide(blank_slide, 1, tmp_path)
        assert content["speaker_notes"] == "Speaker notes here"


class TestResolveThemeRefsInContent:
    """Tests for _resolve_theme_refs_in_content."""

    def test_resolve_string(self):
        theme = {"accent_1": "#0078D4"}
        result = _resolve_theme_refs_in_content("@accent_1", theme)
        assert result == "#0078D4"

    def test_no_theme_ref(self):
        result = _resolve_theme_refs_in_content("#FF0000", {})
        assert result == "#FF0000"

    def test_unresolved_ref(self):
        result = _resolve_theme_refs_in_content("@missing", {})
        assert result == "@missing"

    def test_nested_dict(self):
        theme = {"text_1": "#FFFFFF"}
        content = {"fill": "@text_1", "size": 12}
        result = _resolve_theme_refs_in_content(content, theme)
        assert result["fill"] == "#FFFFFF"
        assert result["size"] == 12

    def test_list(self):
        theme = {"accent_1": "#0078D4"}
        content = ["@accent_1", "#FF0000"]
        result = _resolve_theme_refs_in_content(content, theme)
        assert result == ["#0078D4", "#FF0000"]

    def test_depth_limit_raises(self):
        """Deeply nested dicts exceed the depth limit."""
        nested = "leaf"
        for _ in range(10):
            nested = {"wrap": nested}
        with pytest.raises(ValueError, match="exceeds limit"):
            _resolve_theme_refs_in_content(nested, {}, max_depth=5)

    def test_depth_limit_allows_normal_nesting(self):
        content = {"a": {"b": {"c": "@x"}}}
        theme = {"x": "#AABBCC"}
        result = _resolve_theme_refs_in_content(content, theme, max_depth=50)
        assert result["a"]["b"]["c"] == "#AABBCC"


class TestDetectGlobalStyle:
    """Tests for detect_global_style."""

    def test_basic_presentation(self, blank_presentation):
        style = detect_global_style(blank_presentation)
        assert "dimensions" in style
        assert style["dimensions"]["width_inches"] == pytest.approx(13.333, abs=0.01)
        assert style["dimensions"]["height_inches"] == pytest.approx(7.5, abs=0.01)

    def test_style_has_defaults(self, blank_presentation):
        style = detect_global_style(blank_presentation)
        assert "defaults" in style


class TestExtractGroup:
    """Tests for extract_group."""

    def test_group_structure(self, blank_slide, tmp_path):
        group = blank_slide.shapes.add_group_shape()
        group.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(2),
            Inches(1),
        )
        result = extract_group(group, 1, tmp_path, 0)
        assert result["type"] == "group"
        assert "elements" in result
        assert len(result["elements"]) >= 1

    def test_depth_limit_raises(self, blank_slide, tmp_path):
        """Exceeding max_depth raises ValueError."""
        group = blank_slide.shapes.add_group_shape()
        with pytest.raises(ValueError, match="exceeds limit"):
            extract_group(group, 1, tmp_path, 0, _depth=5, max_depth=5)

    def test_depth_limit_allows_normal_nesting(self, blank_slide, tmp_path):
        """Groups within the depth limit extract successfully."""
        group = blank_slide.shapes.add_group_shape()
        group.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(2),
            Inches(1),
        )
        result = extract_group(group, 1, tmp_path, 0, _depth=0, max_depth=20)
        assert result["type"] == "group"


class TestExtractFreeform:
    """Tests for extract_freeform."""

    def test_freeform_structure(self, blank_slide):
        """Create a shape and verify freeform extraction produces expected keys."""
        # Use a regular shape to test extraction output structure
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(2),
            Inches(2),
        )
        result = extract_freeform(shape)
        assert result["type"] == "freeform"
        assert "left" in result
        assert "top" in result
        assert "width" in result
        assert "height" in result
        assert "name" in result


class TestExtractConnectorExtended:
    """Extended tests for extract_connector."""

    def test_connector_with_line(self, blank_slide):
        from pptx.enum.shapes import MSO_CONNECTOR_TYPE

        connector = blank_slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            Inches(1),
            Inches(1),
            Inches(5),
            Inches(3),
        )
        connector.line.color.rgb = RGBColor(0xFF, 0x00, 0x00)
        result = extract_connector(connector)
        assert result["type"] == "connector"
        assert "line_color" in result or "name" in result


class TestExtractShapeExtended:
    """Extended extract_shape tests for more branches."""

    def test_shape_with_solid_fill(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(3),
            Inches(2),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0x00)
        result = extract_shape(shape)
        assert result["fill"] == "#FF0000"

    def test_shape_with_line(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(3),
            Inches(2),
        )
        shape.line.color.rgb = RGBColor(0x00, 0x00, 0xFF)
        from pptx.util import Pt

        shape.line.width = Pt(2)
        result = extract_shape(shape)
        assert "line_color" in result

    def test_shape_rotation(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(3),
            Inches(2),
        )
        shape.rotation = 45.0
        result = extract_shape(shape)
        assert result.get("rotation") == pytest.approx(45.0, abs=0.1)


class TestExtractSlideExtended:
    """Extended tests for extract_slide."""

    def test_slide_with_connector(self, blank_slide, tmp_path):
        from pptx.enum.shapes import MSO_CONNECTOR_TYPE

        blank_slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            Inches(1),
            Inches(1),
            Inches(5),
            Inches(3),
        )
        content, _ = extract_slide(blank_slide, 1, tmp_path)
        conn_elems = [e for e in content["elements"] if e["type"] == "connector"]
        assert len(conn_elems) == 1

    def test_slide_multiple_shapes(self, blank_slide, tmp_path):
        for i in range(3):
            blank_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(i),
                Inches(i),
                Inches(2),
                Inches(1),
            )
        content, _ = extract_slide(blank_slide, 1, tmp_path)
        assert len(content["elements"]) == 3

    def test_slide_with_group(self, blank_slide, tmp_path):
        group = blank_slide.shapes.add_group_shape()
        group.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(2),
            Inches(1),
        )
        content, _ = extract_slide(blank_slide, 1, tmp_path)
        group_elems = [e for e in content["elements"] if e["type"] == "group"]
        assert len(group_elems) == 1


class TestDetectGlobalStyleExtended:
    """Extended tests for detect_global_style."""

    def test_style_has_fonts(self, blank_presentation):
        style = detect_global_style(blank_presentation)
        assert "defaults" in style

    def test_style_dimensions_match(self, blank_presentation):
        style = detect_global_style(blank_presentation)
        dims = style["dimensions"]
        expected_w = blank_presentation.slide_width / 914400
        expected_h = blank_presentation.slide_height / 914400
        assert dims["width_inches"] == pytest.approx(expected_w, abs=0.01)
        assert dims["height_inches"] == pytest.approx(expected_h, abs=0.01)


class TestSaveImageBlob:
    """Tests for _save_image_blob."""

    def test_embedded_image(self, blank_slide, sample_image_path, tmp_path):
        pic = blank_slide.shapes.add_picture(
            str(sample_image_path),
            Inches(1),
            Inches(1),
            Inches(3),
            Inches(2),
        )
        result = _save_image_blob(pic, tmp_path, 1, 1)
        assert "path" in result
        assert result["path"].startswith("images/")
        assert (tmp_path / result["path"]).exists()

    def test_linked_image_fallback(self):
        shape = MagicMock()
        type(shape).image = property(
            lambda self: (_ for _ in ()).throw(ValueError("no blob"))
        )
        result = _save_image_blob(shape, Path("/tmp"), 1, 1)
        assert result["path"] == "LINKED_IMAGE_NOT_EMBEDDED"

    def test_jpeg_extension_normalized(self, blank_slide, sample_image_path, tmp_path):
        pic = blank_slide.shapes.add_picture(
            str(sample_image_path),
            Inches(1),
            Inches(1),
        )
        result = _save_image_blob(pic, tmp_path, 1, 1)
        assert result["path"].endswith(".png")

    @pytest.mark.parametrize(
        "content_type,expected_ext",
        [
            ("image/bmp", "bmp"),
            ("image/gif", "gif"),
            ("image/jpeg", "jpg"),
            ("image/png", "png"),
            ("image/tiff", "tiff"),
        ],
    )
    def test_allowed_content_type_produces_correct_extension(
        self, content_type, expected_ext, tmp_path
    ):
        """Each allowed content type maps to its correct file extension."""
        mock_shape = MagicMock()
        mock_img = MagicMock()
        mock_img.content_type = content_type
        mock_img.blob = b"\x00" * 100
        mock_shape.image = mock_img

        result = _save_image_blob(mock_shape, tmp_path, 1, 1)
        assert result["path"].endswith(f".{expected_ext}")

    @pytest.mark.parametrize(
        "content_type",
        [
            "image/vnd.ms-photo",
            "application/octet-stream",
            "text/html",
        ],
    )
    def test_unsupported_content_type_rejected(self, content_type, tmp_path):
        """Unsupported content types raise ValueError."""
        mock_shape = MagicMock()
        mock_img = MagicMock()
        mock_img.content_type = content_type
        mock_img.blob = b"\x00" * 100
        mock_shape.image = mock_img

        with pytest.raises(ValueError, match="Unsupported image content type"):
            _save_image_blob(mock_shape, tmp_path, 1, 1)

    def test_oversized_blob_rejected(self, tmp_path, monkeypatch):
        """Blobs exceeding MAX_IMAGE_BLOB_BYTES are rejected."""
        monkeypatch.setattr("extract_content.MAX_IMAGE_BLOB_BYTES", 1024)
        mock_shape = MagicMock()
        mock_img = MagicMock()
        mock_img.content_type = "image/png"
        mock_img.blob = b"\x00" * 1025
        mock_shape.image = mock_img

        with pytest.raises(ValueError, match="exceeds"):
            _save_image_blob(mock_shape, tmp_path, 1, 1)

    def test_blob_at_size_limit_accepted(self, tmp_path, monkeypatch):
        """Blobs at exactly MAX_IMAGE_BLOB_BYTES are accepted."""
        monkeypatch.setattr("extract_content.MAX_IMAGE_BLOB_BYTES", 1024)
        mock_shape = MagicMock()
        mock_img = MagicMock()
        mock_img.content_type = "image/png"
        mock_img.blob = b"\x00" * 1024
        mock_shape.image = mock_img

        result = _save_image_blob(mock_shape, tmp_path, 1, 1)
        assert "images/" in result["path"]

    def test_path_within_output_directory(self, tmp_path):
        """Saved image resolves to within the output directory."""
        mock_shape = MagicMock()
        mock_img = MagicMock()
        mock_img.content_type = "image/png"
        mock_img.blob = b"\x00" * 100
        mock_shape.image = mock_img

        result = _save_image_blob(mock_shape, tmp_path, 1, 1)
        img_path = tmp_path / result["path"]
        assert img_path.resolve().is_relative_to(tmp_path.resolve())

    def test_path_traversal_blocked(self, tmp_path, monkeypatch):
        """Blob write rejects paths that escape the output directory."""
        mock_shape = MagicMock()
        mock_img = MagicMock()
        mock_img.content_type = "image/png"
        mock_img.blob = b"\x00" * 100
        mock_shape.image = mock_img

        original_resolve = Path.resolve

        def patched_resolve(self, strict=False):
            resolved = original_resolve(self, strict=strict)
            if "images" in str(self):
                return Path("/outside") / self.name
            return resolved

        monkeypatch.setattr(Path, "resolve", patched_resolve)

        with pytest.raises(_ImageSecurityError, match="escapes output directory"):
            _save_image_blob(mock_shape, tmp_path, 1, 1)

    def test_wmf_aldus_magic_accepted(self, tmp_path):
        """WMF blob with Aldus Placeable signature is accepted."""
        mock_shape = MagicMock()
        mock_img = MagicMock()
        mock_img.content_type = "image/x-wmf"
        mock_img.blob = b"\xd7\xcd\xc6\x9a" + b"\x00" * 96
        mock_shape.image = mock_img

        result = _save_image_blob(mock_shape, tmp_path, 1, 1)
        assert result["path"].endswith(".wmf")

    def test_wmf_standard_magic_accepted(self, tmp_path):
        """WMF blob with standard header signature is accepted."""
        mock_shape = MagicMock()
        mock_img = MagicMock()
        mock_img.content_type = "image/x-wmf"
        mock_img.blob = b"\x01\x00\x09\x00" + b"\x00" * 96
        mock_shape.image = mock_img

        result = _save_image_blob(mock_shape, tmp_path, 1, 1)
        assert result["path"].endswith(".wmf")

    def test_wmf_invalid_magic_rejected(self, tmp_path):
        """WMF blob without a recognized signature is rejected."""
        mock_shape = MagicMock()
        mock_img = MagicMock()
        mock_img.content_type = "image/x-wmf"
        mock_img.blob = b"\x00" * 100
        mock_shape.image = mock_img

        with pytest.raises(_ImageSecurityError, match="recognized file signature"):
            _save_image_blob(mock_shape, tmp_path, 1, 1)

    def test_wmf_blob_too_short_rejected(self, tmp_path):
        """WMF blob shorter than 4 bytes is rejected."""
        mock_shape = MagicMock()
        mock_img = MagicMock()
        mock_img.content_type = "image/x-wmf"
        mock_img.blob = b"\xd7\xcd"
        mock_shape.image = mock_img

        with pytest.raises(_ImageSecurityError, match="too short"):
            _save_image_blob(mock_shape, tmp_path, 1, 1)

    def test_emf_valid_blob_accepted(self, tmp_path):
        """EMF blob with correct EMR_HEADER record type and signature is accepted."""
        mock_shape = MagicMock()
        mock_img = MagicMock()
        mock_img.content_type = "image/emf"
        blob = b"\x01\x00\x00\x00" + b"\x00" * 36 + b" EMF" + b"\x00" * 56
        mock_img.blob = blob
        mock_shape.image = mock_img

        result = _save_image_blob(mock_shape, tmp_path, 1, 1)
        assert result["path"].endswith(".emf")

    def test_emf_invalid_magic_rejected(self, tmp_path):
        """EMF blob without the expected signature is rejected."""
        mock_shape = MagicMock()
        mock_img = MagicMock()
        mock_img.content_type = "image/emf"
        mock_img.blob = b"\x00" * 100
        mock_shape.image = mock_img

        with pytest.raises(_ImageSecurityError, match="expected file signature"):
            _save_image_blob(mock_shape, tmp_path, 1, 1)

    def test_emf_blob_too_short_rejected(self, tmp_path):
        """EMF blob shorter than 44 bytes is rejected."""
        mock_shape = MagicMock()
        mock_img = MagicMock()
        mock_img.content_type = "image/emf"
        mock_img.blob = b"\x01\x00\x00\x00" + b"\x00" * 10
        mock_shape.image = mock_img

        with pytest.raises(_ImageSecurityError, match="too short"):
            _save_image_blob(mock_shape, tmp_path, 1, 1)

    def test_emf_x_emf_content_type_accepted(self, tmp_path):
        """image/x-emf content type is accepted and validated like image/emf."""
        mock_shape = MagicMock()
        mock_img = MagicMock()
        mock_img.content_type = "image/x-emf"
        blob = b"\x01\x00\x00\x00" + b"\x00" * 36 + b" EMF" + b"\x00" * 56
        mock_img.blob = blob
        mock_shape.image = mock_img

        result = _save_image_blob(mock_shape, tmp_path, 1, 1)
        assert result["path"].endswith(".emf")

    def test_svg_blob_saved_as_png(self, tmp_path):
        """SVG content is converted to PNG and saved with .png extension."""
        mock_shape = MagicMock()
        mock_img = MagicMock()
        mock_img.content_type = "image/svg+xml"
        mock_img.blob = (
            b'<svg xmlns="http://www.w3.org/2000/svg" width="10" height="10">'
            b'<rect width="10" height="10" fill="red"/>'
            b"</svg>"
        )
        mock_shape.image = mock_img

        result = _save_image_blob(mock_shape, tmp_path, 1, 1)
        assert result["path"].endswith(".png")
        saved = (tmp_path / result["path"]).read_bytes()
        assert saved[:4] == b"\x89PNG"

    def test_svg_xxe_rejected(self, tmp_path):
        """SVG blob containing XXE entity raises _ImageSecurityError."""
        mock_shape = MagicMock()
        mock_img = MagicMock()
        mock_img.content_type = "image/svg+xml"
        mock_img.blob = (
            b'<?xml version="1.0"?>'
            b'<!DOCTYPE svg [<!ENTITY xxe SYSTEM "file:///etc/passwd">]>'
            b'<svg xmlns="http://www.w3.org/2000/svg">&xxe;</svg>'
        )
        mock_shape.image = mock_img

        with pytest.raises(_ImageSecurityError, match="DTD declaration"):
            _save_image_blob(mock_shape, tmp_path, 1, 1)


class TestValidateEmfMagicBytes:
    """Tests for _validate_emf_magic_bytes."""

    def test_valid_emf_passes(self):
        blob = b"\x01\x00\x00\x00" + b"\x00" * 36 + b" EMF" + b"\x00" * 56
        _validate_emf_magic_bytes(blob)

    def test_wrong_record_type_rejected(self):
        blob = b"\x02\x00\x00\x00" + b"\x00" * 36 + b" EMF" + b"\x00" * 56
        with pytest.raises(_ImageSecurityError, match="expected file signature"):
            _validate_emf_magic_bytes(blob)

    def test_missing_emf_signature_rejected(self):
        blob = b"\x01\x00\x00\x00" + b"\x00" * 36 + b"XXXX" + b"\x00" * 56
        with pytest.raises(_ImageSecurityError, match="expected file signature"):
            _validate_emf_magic_bytes(blob)

    def test_too_short_rejected(self):
        with pytest.raises(_ImageSecurityError, match="too short"):
            _validate_emf_magic_bytes(b"\x01\x00\x00\x00" + b"\x00" * 10)


class TestSanitizeSvg:
    """Tests for _sanitize_svg."""

    def test_valid_svg_passes(self):
        svg = b'<svg xmlns="http://www.w3.org/2000/svg"><rect/></svg>'
        result = _sanitize_svg(svg)
        assert b"<svg" in result

    def test_xxe_entity_rejected(self):
        svg = (
            b'<?xml version="1.0"?>'
            b'<!DOCTYPE svg [<!ENTITY xxe SYSTEM "file:///etc/passwd">]>'
            b'<svg xmlns="http://www.w3.org/2000/svg">&xxe;</svg>'
        )
        with pytest.raises(_ImageSecurityError, match="DTD declaration"):
            _sanitize_svg(svg)

    def test_non_xml_blob_rejected(self):
        with pytest.raises(_ImageSecurityError, match="not well-formed XML"):
            _sanitize_svg(b"\x89PNG\r\n\x1a\n not xml")


class TestConvertSvgToPng:
    """Tests for _convert_svg_to_png."""

    def test_valid_svg_produces_png(self):
        svg = (
            b'<svg xmlns="http://www.w3.org/2000/svg" width="10" height="10">'
            b'<rect width="10" height="10" fill="red"/>'
            b"</svg>"
        )
        result = _convert_svg_to_png(svg)
        assert result[:4] == b"\x89PNG"

    def test_xxe_rejected_before_conversion(self):
        svg = (
            b'<?xml version="1.0"?>'
            b'<!DOCTYPE svg [<!ENTITY xxe SYSTEM "file:///etc/passwd">]>'
            b'<svg xmlns="http://www.w3.org/2000/svg">&xxe;</svg>'
        )
        with pytest.raises(_ImageSecurityError, match="DTD declaration"):
            _convert_svg_to_png(svg)


class TestExtractChildShape:
    """Tests for extract_child_shape dispatch."""

    def test_textbox_type(self, blank_slide, tmp_path):
        txBox = blank_slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(1)
        )
        txBox.text_frame.text = "child text"
        result = extract_child_shape(txBox, 1, tmp_path, 0)
        assert result is not None
        assert result["type"] == "textbox"

    def test_auto_shape_type(self, blank_slide, tmp_path):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(2),
            Inches(2),
        )
        result = extract_child_shape(shape, 1, tmp_path, 0)
        assert result is not None
        assert result["type"] == "shape"

    def test_connector_type(self, blank_slide, tmp_path):
        from pptx.enum.shapes import MSO_CONNECTOR_TYPE

        conn = blank_slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            Inches(1),
            Inches(1),
            Inches(5),
            Inches(3),
        )
        result = extract_child_shape(conn, 1, tmp_path, 0)
        assert result is not None
        assert result["type"] == "connector"

    def test_picture_type(self, blank_slide, sample_image_path, tmp_path):
        pic = blank_slide.shapes.add_picture(
            str(sample_image_path),
            Inches(1),
            Inches(1),
            Inches(3),
            Inches(2),
        )
        result = extract_child_shape(pic, 1, tmp_path, 1)
        assert result is not None
        assert result["type"] == "image"

    def test_group_type(self, blank_slide, tmp_path):
        group = blank_slide.shapes.add_group_shape()
        result = extract_child_shape(group, 1, tmp_path, 0)
        assert result is not None
        assert result["type"] == "group"

    def test_unrecognized_type(self, tmp_path):
        shape = MagicMock()
        shape.shape_type = 99
        shape.left = Inches(1)
        shape.top = Inches(1)
        shape.width = Inches(2)
        shape.height = Inches(2)
        shape.name = "Unknown"
        shape.has_table = False
        shape.has_chart = False
        shape._element = MagicMock()
        shape._element.find.return_value = None
        result = extract_child_shape(shape, 1, tmp_path, 0)
        assert result is not None
        assert result.get("_unrecognized_shape_type") == 99


class TestClassifySlideBrightness:
    """Tests for _classify_slide_brightness."""

    def test_dark_background(self):
        result = _classify_slide_brightness("#1A1A2E", Counter(), False)
        assert result == "dark"

    def test_light_background(self):
        result = _classify_slide_brightness("#FFFFFF", Counter(), False)
        assert result == "light"

    def test_bg_image_dark_text(self):
        text_colors = Counter({"#333333": 5, "#222222": 3})
        result = _classify_slide_brightness(None, text_colors, True)
        assert result == "light"

    def test_bg_image_light_text(self):
        text_colors = Counter({"#FFFFFF": 5, "#EEEEEE": 3})
        result = _classify_slide_brightness(None, text_colors, True)
        assert result == "dark"

    def test_no_bg_infer_from_text_dark(self):
        text_colors = Counter({"#333333": 10})
        result = _classify_slide_brightness(None, text_colors, False)
        assert result == "light"

    def test_no_bg_infer_from_text_light(self):
        text_colors = Counter({"#FFFFFF": 10})
        result = _classify_slide_brightness(None, text_colors, False)
        assert result == "dark"

    def test_no_bg_no_text_defaults_dark(self):
        result = _classify_slide_brightness(None, Counter(), False)
        assert result == "dark"

    def test_equal_text_defaults_dark(self):
        text_colors = Counter({"#333333": 5, "#FFFFFF": 5})
        result = _classify_slide_brightness(None, text_colors, False)
        assert result == "dark"


class TestBuildColorMap:
    """Tests for _build_color_map."""

    def test_empty_counters(self):
        result = _build_color_map(Counter(), Counter(), Counter(), Counter())
        assert result == {}

    def test_bg_dark(self):
        bg = Counter({"#1A1A2E": 5})
        result = _build_color_map(bg, Counter(), Counter(), Counter())
        assert result["bg_dark"] == "#1A1A2E"

    def test_fill_card(self):
        fill = Counter({"#2D2D44": 3})
        result = _build_color_map(Counter(), fill, Counter(), Counter())
        assert result["bg_card"] == "#2D2D44"

    def test_text_colors_white(self):
        text = Counter({"#FAFAFA": 10})
        result = _build_color_map(Counter(), Counter(), text, Counter())
        assert result.get("text_white") == "#FAFAFA"

    def test_text_colors_dark(self):
        text = Counter({"#222222": 10})
        result = _build_color_map(Counter(), Counter(), text, Counter())
        assert result.get("text_dark") == "#222222"

    def test_text_colors_gray(self):
        text = Counter({"#999999": 10})
        result = _build_color_map(Counter(), Counter(), text, Counter())
        assert result.get("text_gray") == "#999999"

    def test_accent_colors(self):
        accents = Counter({"#0078D4": 5, "#00B294": 3, "#40A040": 2})
        result = _build_color_map(Counter(), Counter(), Counter(), accents)
        assert "accent_blue" in result
        assert "accent_teal" in result
        assert "accent_green" in result


class TestClusterThemes:
    """Tests for _cluster_themes."""

    def test_empty_profiles(self):
        result = _cluster_themes([], Counter(), Counter(), Counter())
        assert result == []

    def test_all_dark_no_clusters(self):
        profiles = [
            {
                "slide": 1,
                "bg_brightness": "dark",
                "bg_color": "#000",
                "text_colors": {},
                "fill_colors": {},
                "has_bg_image": False,
            },
        ]
        result = _cluster_themes(profiles, Counter(), Counter(), Counter())
        assert result == []

    def test_mixed_themes(self):
        profiles = [
            {
                "slide": 1,
                "bg_brightness": "light",
                "bg_color": "#FFF",
                "text_colors": {"#333333": 5},
                "fill_colors": {"#EEEEEE": 2},
                "has_bg_image": False,
            },
            {
                "slide": 2,
                "bg_brightness": "dark",
                "bg_color": "#111",
                "text_colors": {"#FFFFFF": 5},
                "fill_colors": {"#333333": 2},
                "has_bg_image": False,
            },
        ]
        text = Counter({"#333333": 5, "#FFFFFF": 5})
        fills = Counter({"#EEEEEE": 2, "#333333": 2})
        result = _cluster_themes(profiles, text, fills, Counter())
        assert len(result) == 2
        names = {t["name"] for t in result}
        assert "light" in names
        assert "dark" in names


class TestResolveThemeColors:
    """Tests for _resolve_theme_colors."""

    def test_basic_presentation(self, blank_presentation):
        colors = _resolve_theme_colors(blank_presentation)
        # Default presentation has a theme with some color mappings
        assert isinstance(colors, dict)

    def test_empty_presentation_has_colors(self, blank_presentation):
        colors = _resolve_theme_colors(blank_presentation)
        # A standard blank presentation should have some theme colors
        # (dk1/lt1 etc. are always present in the default theme)
        if colors:
            assert any(k.startswith("dark") or k.startswith("light") for k in colors)

    @pytest.mark.parametrize(
        "entity_uri",
        [
            "file:///dev/null",
            "file:///etc/passwd",
        ],
    )
    def test_xxe_system_entity_blocked(self, entity_uri):
        """SYSTEM entities in theme XML are blocked by the hardened parser."""
        xxe_xml = (
            b'<?xml version="1.0"?>'
            b'<!DOCTYPE foo [<!ENTITY xxe SYSTEM "' + entity_uri.encode() + b'">]>'
            b'<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            b"<a:themeElements><a:clrScheme name='Test'>"
            b'<a:dk1><a:srgbClr val="&xxe;"/></a:dk1>'
            b"</a:clrScheme></a:themeElements></a:theme>"
        )
        rel = MagicMock()
        rel.reltype = (
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        rel.target_part.blob = xxe_xml
        prs = MagicMock()
        prs.slide_masters = [MagicMock()]
        prs.slide_masters[0].part.rels.values.return_value = [rel]
        colors = _resolve_theme_colors(prs)
        assert isinstance(colors, dict)
        assert "dark_1" not in colors
        for val in colors.values():
            assert entity_uri not in val

    def test_predefined_entities_preserved(self):
        """Predefined XML entities still resolve with the hardened parser."""
        parser = etree.XMLParser(resolve_entities=False, no_network=True)
        root = etree.fromstring(b'<root attr="&lt;value&gt;"/>', parser=parser)
        assert root.get("attr") == "<value>"


class TestExtractShapeFormatting:
    """Tests for extract_shape deeper formatting branches."""

    def test_shape_text_with_color(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(3),
            Inches(2),
        )
        tf = shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Colored"
        run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
        result = extract_shape(shape)
        assert result.get("text") == "Colored"

    def test_shape_text_bold_italic(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(3),
            Inches(2),
        )
        tf = shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Styled"
        run.font.bold = True
        run.font.italic = True
        result = extract_shape(shape)
        assert result.get("text_bold") is True or (
            "paragraphs" in result and result["paragraphs"][0].get("text_bold") is True
        )

    def test_shape_text_with_font_name_and_size(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(3),
            Inches(2),
        )
        tf = shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Formatted"
        run.font.name = "Arial"
        run.font.size = Pt(18)
        result = extract_shape(shape)
        # Should capture font info at element or paragraph level
        has_font = result.get("text_font") or (
            "paragraphs" in result and result["paragraphs"][0].get("text_font")
        )
        assert has_font

    def test_shape_corner_radius(self, blank_slide):
        shape = blank_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(3),
            Inches(2),
        )
        result = extract_shape(shape)
        assert "corner_radius" in result


class TestExtractTextboxFormatting:
    """Tests for extract_textbox deeper formatting branches."""

    def test_textbox_with_color(self, blank_slide):
        txBox = blank_slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(1),
        )
        p = txBox.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Red"
        run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
        result = extract_textbox(txBox)
        assert result.get("font_color") or (
            "paragraphs" in result and result["paragraphs"][0].get("font_color")
        )

    def test_textbox_with_underline(self, blank_slide):
        txBox = blank_slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(1),
        )
        p = txBox.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Underlined"
        run.font.underline = True
        result = extract_textbox(txBox)
        has_underline = result.get("underline") or (
            "paragraphs" in result and result["paragraphs"][0].get("underline")
        )
        assert has_underline

    def test_textbox_font_name_and_size(self, blank_slide):
        txBox = blank_slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(1),
        )
        p = txBox.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Styled"
        run.font.name = "Calibri"
        run.font.size = Pt(20)
        result = extract_textbox(txBox)
        has_font = result.get("font") or (
            "paragraphs" in result and result["paragraphs"][0].get("font")
        )
        assert has_font

    def test_textbox_multiline_with_formatting(self, blank_slide):
        txBox = blank_slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(2),
        )
        tf = txBox.text_frame
        p1 = tf.paragraphs[0]
        run1 = p1.add_run()
        run1.text = "Para 1"
        run1.font.bold = True

        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = "Para 2"
        run2.font.italic = True

        result = extract_textbox(txBox)
        assert "paragraphs" in result
        assert len(result["paragraphs"]) == 2


class TestExtractImageExtended:
    """Extended tests for extract_image."""

    def test_linked_image(self, tmp_path):
        shape = MagicMock()
        type(shape).image = property(
            lambda self: (_ for _ in ()).throw(ValueError("linked"))
        )
        shape.left = Inches(1)
        shape.top = Inches(1)
        shape.width = Inches(3)
        shape.height = Inches(2)
        shape.name = "LinkedPic"
        shape._element = MagicMock()
        shape._element.find.return_value = None
        result = extract_image(shape, tmp_path, 1, 1)
        assert result["path"] == "LINKED_IMAGE_NOT_EMBEDDED"
        assert result["type"] == "image"

    def test_unsupported_content_type_skipped(self, tmp_path):
        """Benign ValueError from unsupported content type is caught and skipped."""
        mock_shape = MagicMock()
        mock_img = MagicMock()
        mock_img.content_type = "image/vnd.ms-photo"
        mock_img.blob = b"\x00" * 100
        mock_shape.image = mock_img
        mock_shape.left = Inches(1)
        mock_shape.top = Inches(1)
        mock_shape.width = Inches(3)
        mock_shape.height = Inches(2)
        mock_shape.name = "SkippedPic"

        result = extract_image(mock_shape, tmp_path, 1, 1)
        assert result["path"] == "SKIPPED"
        assert "Unsupported image content type" in result["_skipped_reason"]

    def test_security_error_propagates(self, tmp_path):
        """_ImageSecurityError from path traversal is not caught by extract_image."""
        mock_shape = MagicMock()
        mock_img = MagicMock()
        mock_img.content_type = "image/x-wmf"
        mock_img.blob = b"\x00" * 100  # invalid magic bytes
        mock_shape.image = mock_img
        mock_shape.left = Inches(1)
        mock_shape.top = Inches(1)
        mock_shape.width = Inches(3)
        mock_shape.height = Inches(2)
        mock_shape.name = "SecurityPic"

        with pytest.raises(_ImageSecurityError):
            extract_image(mock_shape, tmp_path, 1, 1)

    def test_svg_xxe_propagates(self, tmp_path):
        """_ImageSecurityError from SVG XXE is not caught by extract_image."""
        mock_shape = MagicMock()
        mock_img = MagicMock()
        mock_img.content_type = "image/svg+xml"
        mock_img.blob = (
            b'<?xml version="1.0"?>'
            b'<!DOCTYPE svg [<!ENTITY xxe SYSTEM "file:///etc/passwd">]>'
            b'<svg xmlns="http://www.w3.org/2000/svg">&xxe;</svg>'
        )
        mock_shape.image = mock_img
        mock_shape.left = Inches(1)
        mock_shape.top = Inches(1)
        mock_shape.width = Inches(3)
        mock_shape.height = Inches(2)
        mock_shape.name = "SvgXxePic"

        with pytest.raises(_ImageSecurityError, match="DTD declaration"):
            extract_image(mock_shape, tmp_path, 1, 1)

    def test_oversized_blob_skipped(self, tmp_path, monkeypatch):
        """Benign ValueError from oversized blob is caught and skipped."""
        monkeypatch.setattr("extract_content.MAX_IMAGE_BLOB_BYTES", 1024)
        mock_shape = MagicMock()
        mock_img = MagicMock()
        mock_img.content_type = "image/png"
        mock_img.blob = b"\x00" * 1025
        mock_shape.image = mock_img
        mock_shape.left = Inches(1)
        mock_shape.top = Inches(1)
        mock_shape.width = Inches(3)
        mock_shape.height = Inches(2)
        mock_shape.name = "OversizedPic"

        result = extract_image(mock_shape, tmp_path, 1, 1)
        assert result["path"] == "SKIPPED"
        assert "exceeds" in result["_skipped_reason"]


class TestExtractSlideDeep:
    """Tests for extract_slide covering more shape dispatch branches."""

    def test_slide_with_placeholder(self, blank_presentation, tmp_path):
        # Use a layout that has title placeholder
        layout = blank_presentation.slide_layouts[0]
        slide = blank_presentation.slides.add_slide(layout)
        # Set title
        if slide.placeholders:
            slide.placeholders[0].text = "Title Text"
        content, _ = extract_slide(slide, 1, tmp_path)
        placeholders = [e for e in content["elements"] if e.get("_placeholder")]
        assert len(placeholders) >= 1

    def test_slide_layout_name(self, blank_slide, tmp_path):
        content, _ = extract_slide(blank_slide, 1, tmp_path)
        assert "layout" in content

    def test_slide_background_fill(self, blank_slide, tmp_path):
        content, _ = extract_slide(blank_slide, 1, tmp_path)
        # Background fill extraction is tested through the slide path;
        # follow_master_background is read-only, so we verify structure
        assert "elements" in content


class TestDetectGlobalStyleDeep:
    """Tests for detect_global_style with styled slides."""

    def test_style_with_text(self, blank_presentation):
        layout = blank_presentation.slide_layouts[6]
        slide = blank_presentation.slides.add_slide(layout)
        txBox = slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(1),
        )
        run = txBox.text_frame.paragraphs[0].add_run()
        run.text = "Sample text"
        run.font.name = "Arial"
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        style = detect_global_style(blank_presentation)
        assert "dimensions" in style
        assert "defaults" in style

    def test_style_with_shapes_and_fills(self, blank_presentation):
        layout = blank_presentation.slide_layouts[6]
        slide = blank_presentation.slides.add_slide(layout)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(4),
            Inches(3),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x44)
        style = detect_global_style(blank_presentation)
        assert "dimensions" in style

    def test_style_with_metadata(self, blank_presentation):
        blank_presentation.core_properties.title = "Test Deck"
        blank_presentation.core_properties.author = "Tester"
        style = detect_global_style(blank_presentation)
        assert "metadata" in style
        assert style["metadata"]["title"] == "Test Deck"

    def test_style_with_accent_shape(self, blank_presentation):
        """A thin shape produces an accent color."""
        layout = blank_presentation.slide_layouts[6]
        slide = blank_presentation.slides.add_slide(layout)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(3),
            Inches(10),
            Inches(0.05),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x00, 0x78, 0xD4)
        style = detect_global_style(blank_presentation)
        assert "dimensions" in style


class TestMainExtractContent:
    """Tests for main() entry point."""

    def test_main_basic(self, tmp_path, mocker):
        from extract_content import main

        mock_prs_cls = mocker.patch("extract_content.Presentation")
        mock_prs = MagicMock()
        mock_prs.slides = []
        mock_prs.slide_width = Inches(13.333)
        mock_prs.slide_height = Inches(7.5)
        mock_prs.slide_masters = []
        mock_prs.core_properties = MagicMock(
            title=None,
            author=None,
            subject=None,
            keywords=None,
            description=None,
            category=None,
        )
        mock_prs_cls.return_value = mock_prs

        pptx_file = tmp_path / "deck.pptx"
        pptx_file.write_bytes(b"PK")
        out_dir = tmp_path / "output"

        mocker.patch(
            "sys.argv",
            [
                "extract_content.py",
                "--input",
                str(pptx_file),
                "--output-dir",
                str(out_dir),
            ],
        )
        main()
        assert (out_dir / "global" / "style.yaml").exists()
