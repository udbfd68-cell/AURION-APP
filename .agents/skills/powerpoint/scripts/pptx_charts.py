# Copyright (c) Microsoft Corporation.
# SPDX-License-Identifier: MIT
"""Chart build and extract utilities for PowerPoint skill scripts.

Provides add_chart_element() for building charts from YAML definitions
and extract_chart() for extracting chart data from existing presentations.
"""

from pptx.chart.data import BubbleChartData, CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx_colors import apply_color_to_fill, resolve_color
from pptx_utils import emu_to_inches

CHART_TYPE_MAP = {
    "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
    "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA,
    "radar": XL_CHART_TYPE.RADAR,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "bubble": XL_CHART_TYPE.BUBBLE,
}

CHART_TYPE_REVERSE = {v: k for k, v in CHART_TYPE_MAP.items()}

SCATTER_CHART_TYPES = {"scatter", "scatter_lines", "scatter_smooth"}
BUBBLE_CHART_TYPES = {"bubble"}


def add_chart_element(slide, elem: dict, colors: dict):
    """Add a chart element from a content.yaml definition.

    YAML schema:
    - type: chart
      chart_type: column_clustered
      left: 1.0
      top: 2.0
      width: 8.0
      height: 4.5
      title: "Quarterly Sales"
      has_legend: true
      chart_style: 10
      categories: ["Q1", "Q2", "Q3", "Q4"]
      series:
        - name: "East"
          values: [19.2, 22.3, 18.4, 23.1]
          color: "#0078D4"
    """
    chart_type_name = elem.get("chart_type", "column_clustered")
    chart_type = CHART_TYPE_MAP.get(chart_type_name, XL_CHART_TYPE.COLUMN_CLUSTERED)

    # Choose data class based on chart type
    if chart_type_name in SCATTER_CHART_TYPES:
        chart_data = XyChartData()
        for series_spec in elem.get("series", []):
            series = chart_data.add_series(series_spec.get("name", ""))
            x_values = series_spec.get("x_values", [])
            y_values = series_spec.get("y_values", [])
            for x_val, y_val in zip(x_values, y_values):
                series.add_data_point(x_val, y_val)
    elif chart_type_name in BUBBLE_CHART_TYPES:
        chart_data = BubbleChartData()
        for series_spec in elem.get("series", []):
            series = chart_data.add_series(series_spec.get("name", ""))
            x_values = series_spec.get("x_values", [])
            y_values = series_spec.get("y_values", [])
            sizes = series_spec.get("sizes", [])
            for x, y, size in zip(x_values, y_values, sizes):
                series.add_data_point(x, y, size)
    else:
        chart_data = CategoryChartData()
        chart_data.categories = elem.get("categories", [])
        for series_spec in elem.get("series", []):
            chart_data.add_series(
                series_spec.get("name", ""),
                series_spec.get("values", []),
            )

    chart_shape = slide.shapes.add_chart(
        chart_type,
        Inches(elem["left"]),
        Inches(elem["top"]),
        Inches(elem["width"]),
        Inches(elem["height"]),
        chart_data,
    )
    chart = chart_shape.chart

    # Chart properties
    if "title" in elem:
        chart.has_title = True
        chart.chart_title.text_frame.text = elem["title"]
    if "has_legend" in elem:
        chart.has_legend = elem["has_legend"]
    if "chart_style" in elem:
        chart.style = elem["chart_style"]

    # Series coloring
    for i, series_spec in enumerate(elem.get("series", [])):
        if "color" in series_spec and i < len(chart.series):
            series = chart.series[i]
            series.format.fill.solid()
            color_spec = resolve_color(series_spec["color"], colors)
            apply_color_to_fill(series.format.fill, color_spec)

    if "name" in elem:
        chart_shape.name = elem["name"]

    return chart_shape


def extract_chart(shape) -> dict:
    """Extract a chart element definition from a GraphicFrame shape."""
    chart = shape.chart

    elem = {
        "type": "chart",
        "chart_type": CHART_TYPE_REVERSE.get(chart.chart_type, "column_clustered"),
        "left": emu_to_inches(shape.left),
        "top": emu_to_inches(shape.top),
        "width": emu_to_inches(shape.width),
        "height": emu_to_inches(shape.height),
        "name": shape.name,
    }

    if chart.has_title:
        try:
            elem["title"] = chart.chart_title.text_frame.text
        except (AttributeError, TypeError):
            pass
    elem["has_legend"] = chart.has_legend

    # Extract categories and series data
    try:
        plot = chart.plots[0]
        if hasattr(plot, "categories") and plot.categories:
            elem["categories"] = list(plot.categories)
        elem["series"] = []
        for series in plot.series:
            series_data = {
                "name": series.name or "",
                "values": list(series.values),
            }
            elem["series"].append(series_data)
    except (IndexError, AttributeError):
        pass

    return elem
