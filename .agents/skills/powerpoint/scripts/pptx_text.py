# Copyright (c) Microsoft Corporation.
# SPDX-License-Identifier: MIT
"""Text frame, paragraph, and run property utilities for PowerPoint skill scripts.

Centralizes text properties (margins, auto-size, spacing, underline,
hyperlinks, bullets) and shared text-frame population used by build_deck.py
and extract_content.py.
"""

import re

from lxml import etree
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from pptx_colors import apply_color_to_font, resolve_color
from pptx_fills import build_shadow_xml, parse_shadow_xml
from pptx_fonts import ALIGNMENT_MAP, _extract_char_spacing

AUTO_SIZE_MAP = {
    "none": MSO_AUTO_SIZE.NONE,
    "fit": MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
    "shrink": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
}

AUTO_SIZE_REVERSE = {
    MSO_AUTO_SIZE.NONE: "none",
    MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT: "fit",
    MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE: "shrink",
}

VERTICAL_ANCHOR_MAP = {
    "top": MSO_VERTICAL_ANCHOR.TOP,
    "middle": MSO_VERTICAL_ANCHOR.MIDDLE,
    "bottom": MSO_VERTICAL_ANCHOR.BOTTOM,
}

VERTICAL_ANCHOR_REVERSE = {
    MSO_VERTICAL_ANCHOR.TOP: "top",
    MSO_VERTICAL_ANCHOR.MIDDLE: "middle",
    MSO_VERTICAL_ANCHOR.BOTTOM: "bottom",
}

# EMU per inch constant for margin conversions
_EMU_PER_INCH = 914400

# Key-mapping specifications for YAML schema differences per element type.
# Maps canonical keys (font, size, color, bold) to the actual YAML key names.
TEXTBOX_KEYS = {
    "font": "font",
    "size": "font_size",
    "color": "font_color",
    "bold": "font_bold",
}
SHAPE_KEYS = {
    "font": "text_font",
    "size": "text_size",
    "color": "text_color",
    "bold": "text_bold",
}


def split_lines(text: str) -> list[str]:
    """Split text on newline and vertical-tab characters."""
    if "\n" in text or "\v" in text:
        return re.split(r"\n|\v", text)
    return [text]


def _apply_run_formatting(run, elem: dict, keys: dict, defaults: dict, colors: dict):
    """Apply font properties to a single run using key-mapped element values.

    Args:
        run: python-pptx Run object.
        elem: Dict with run or paragraph properties.
        keys: Key-mapping dict (TEXTBOX_KEYS or SHAPE_KEYS).
        defaults: Fallback values for font, size, color, bold, italic.
        colors: Color resolution dict.
    """
    font_name = elem.get(keys["font"], defaults.get("font"))
    if font_name:
        run.font.name = font_name
    run.font.size = Pt(elem.get(keys["size"], defaults.get("size", 16)))

    color_val = elem.get(keys["color"])
    if color_val:
        apply_color_to_font(run.font.color, resolve_color(color_val, colors))
    elif defaults.get("color"):
        apply_color_to_font(run.font.color, defaults["color"])

    run.font.bold = elem.get(keys["bold"], defaults.get("bold", False))
    run.font.italic = elem.get("italic", defaults.get("italic", False))
    apply_run_properties(run, elem, colors)


def _apply_rich_run_formatting(run, seg: dict, defaults: dict, colors: dict):
    """Apply formatting from a rich-text run segment.

    Rich-text segments use short keys: font, size, color, bold, italic.
    """
    seg_font = seg.get("font", defaults.get("font"))
    if seg_font:
        run.font.name = seg_font
    run.font.size = Pt(seg.get("size", defaults.get("size", 16)))

    if "color" in seg:
        apply_color_to_font(run.font.color, resolve_color(seg["color"], colors))
    elif defaults.get("color"):
        apply_color_to_font(run.font.color, defaults["color"])

    run.font.bold = seg.get("bold", False)
    run.font.italic = seg.get("italic", False)
    apply_run_properties(run, seg, colors)


def populate_text_frame(
    tf, elem: dict, colors: dict, keys: dict, defaults: dict | None = None
):
    """Populate a text frame from an element definition.

    Handles three text layouts:
    1. Per-paragraph with optional rich-text runs (elem has "paragraphs" key)
    2. Flat text with newline splitting (elem has "text" key)
    3. No text (no-op)

    Args:
        tf: python-pptx TextFrame object.
        elem: Element definition dict from content.yaml.
        colors: Color resolution dict.
        keys: Key-mapping dict (TEXTBOX_KEYS or SHAPE_KEYS).
        defaults: Fallback values for font, size, color, bold, italic, alignment.
    """
    defaults = defaults or {}
    tf.word_wrap = True
    apply_text_properties(tf, elem)

    paragraphs = elem.get("paragraphs")
    if paragraphs:
        _populate_paragraphs(tf, paragraphs, elem, colors, keys, defaults)
        return

    text = elem.get("text")
    if text is None:
        return

    _populate_flat_text(tf, text, elem, colors, keys, defaults)


def _populate_paragraphs(
    tf, paragraphs: list[dict], elem: dict, colors: dict, keys: dict, defaults: dict
):
    """Populate text frame from per-paragraph definitions."""
    alignment = defaults.get("alignment")

    for i, p_def in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p_align = p_def.get("alignment", alignment)
        if p_align:
            p.alignment = ALIGNMENT_MAP.get(p_align, ALIGNMENT_MAP["left"])
        apply_paragraph_properties(p, p_def)
        apply_bullet_properties(p, p_def)

        runs = p_def.get("runs")
        if runs:
            for j, seg in enumerate(runs):
                run = p.add_run() if j > 0 else (p.runs[0] if p.runs else p.add_run())
                run.text = seg.get("text", "")
                _apply_rich_run_formatting(run, seg, defaults, colors)
        else:
            run = p.add_run()
            run.text = p_def.get("text", "")
            _apply_run_formatting(run, p_def, keys, defaults, colors)


def _populate_flat_text(
    tf, text: str, elem: dict, colors: dict, keys: dict, defaults: dict
):
    """Populate text frame with flat text split on newlines."""
    alignment = defaults.get("alignment") or elem.get("alignment")
    lines = split_lines(text)

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if alignment:
            p.alignment = ALIGNMENT_MAP.get(alignment, ALIGNMENT_MAP["left"])
        apply_paragraph_properties(p, elem)
        apply_bullet_properties(p, elem)
        run = p.add_run()
        run.text = line
        _apply_run_formatting(run, elem, keys, defaults, colors)


def apply_text_properties(text_frame, elem: dict):
    """Apply text frame-level properties from element definition.

    Supports: margin_left/right/top/bottom (inches), auto_size,
    vertical_anchor, word_wrap.
    """
    if "margin_left" in elem:
        text_frame.margin_left = Inches(elem["margin_left"])
    if "margin_right" in elem:
        text_frame.margin_right = Inches(elem["margin_right"])
    if "margin_top" in elem:
        text_frame.margin_top = Inches(elem["margin_top"])
    if "margin_bottom" in elem:
        text_frame.margin_bottom = Inches(elem["margin_bottom"])
    if "auto_size" in elem:
        text_frame.auto_size = AUTO_SIZE_MAP.get(elem["auto_size"], MSO_AUTO_SIZE.NONE)
    if "vertical_anchor" in elem:
        text_frame.vertical_anchor = VERTICAL_ANCHOR_MAP.get(elem["vertical_anchor"])
    if "word_wrap" in elem:
        text_frame.word_wrap = elem["word_wrap"]


def apply_paragraph_properties(paragraph, elem: dict):
    """Apply paragraph-level properties.

    Supports: space_before, space_after (pts), line_spacing (pts or factor), level.
    """
    if "space_before" in elem:
        paragraph.space_before = Pt(elem["space_before"])
    if "space_after" in elem:
        paragraph.space_after = Pt(elem["space_after"])
    if "line_spacing" in elem:
        val = elem["line_spacing"]
        if isinstance(val, float) and val < 10:
            # Factor-based spacing (e.g. 1.5 = 150%)
            paragraph.line_spacing = val
        else:
            paragraph.line_spacing = Pt(val)
    if "level" in elem:
        paragraph.level = elem["level"]


def apply_run_properties(run, elem: dict, colors: dict):
    """Apply run-level font properties beyond basic font/size/color/bold/italic.

    Supports: underline, hyperlink, char_spacing, effect (outer shadow).
    When a hyperlink is set, the font color is re-applied afterward to prevent
    the automatic theme hyperlink color from overriding the intended color.
    """
    if elem.get("underline"):
        run.font.underline = True
    if "hyperlink" in elem:
        run.hyperlink.address = elem["hyperlink"]
        # Re-apply font color after hyperlink to override auto-coloring
        color_key = next(
            (k for k in ("font_color", "text_color", "color") if k in elem), None
        )
        if color_key:
            apply_color_to_font(run.font.color, resolve_color(elem[color_key], colors))
    if "char_spacing" in elem:
        _apply_char_spacing(run.font, elem["char_spacing"])
    effect = elem.get("effect") or elem.get("text_effect")
    if effect:
        _apply_run_effect(run, effect)


def _apply_run_effect(run, effect: dict):
    """Apply outer shadow effect to a run's rPr element."""
    if not effect or effect.get("type") != "outer_shadow":
        return
    rpr = run.font._element
    existing = rpr.find(qn("a:effectLst"))
    if existing is not None:
        rpr.remove(existing)

    effect_lst = etree.SubElement(rpr, qn("a:effectLst"))
    build_shadow_xml(effect_lst, effect)


def _apply_char_spacing(font, spacing_pt: float):
    """Apply character spacing to a font via the spc attribute on a:rPr.

    Args:
        font: python-pptx font object.
        spacing_pt: Spacing in points (converted to hundredths of a point for XML).
    """
    rpr = font._element
    spc_val = str(int(spacing_pt * 100))
    rpr.set("spc", spc_val)


def extract_text_frame_properties(text_frame) -> dict:
    """Extract text frame-level properties (margins, auto_size, vertical_anchor)."""
    props = {}
    if text_frame.margin_left is not None:
        props["margin_left"] = round(text_frame.margin_left / _EMU_PER_INCH, 3)
    if text_frame.margin_right is not None:
        props["margin_right"] = round(text_frame.margin_right / _EMU_PER_INCH, 3)
    if text_frame.margin_top is not None:
        props["margin_top"] = round(text_frame.margin_top / _EMU_PER_INCH, 3)
    if text_frame.margin_bottom is not None:
        props["margin_bottom"] = round(text_frame.margin_bottom / _EMU_PER_INCH, 3)
    if text_frame.auto_size is not None:
        label = AUTO_SIZE_REVERSE.get(text_frame.auto_size)
        if label:
            props["auto_size"] = label
    if text_frame.vertical_anchor is not None:
        label = VERTICAL_ANCHOR_REVERSE.get(text_frame.vertical_anchor)
        if label:
            props["vertical_anchor"] = label
    return props


def extract_paragraph_properties(paragraph) -> dict:
    """Extract paragraph-level spacing properties."""
    props = {}
    if paragraph.space_before is not None:
        props["space_before"] = round(paragraph.space_before.pt, 1)
    if paragraph.space_after is not None:
        props["space_after"] = round(paragraph.space_after.pt, 1)
    if paragraph.line_spacing is not None:
        if isinstance(paragraph.line_spacing, float):
            props["line_spacing"] = paragraph.line_spacing
        else:
            props["line_spacing"] = round(paragraph.line_spacing.pt, 1)
    if paragraph.level and paragraph.level > 0:
        props["level"] = paragraph.level
    return props


def extract_run_properties(run) -> dict:
    """Extract run-level properties beyond basic font info.

    Extracts underline, hyperlink, char_spacing, and effects.
    """
    props = {}
    if run.font.underline:
        props["underline"] = True
    try:
        if run.hyperlink and run.hyperlink.address:
            props["hyperlink"] = run.hyperlink.address
    except (AttributeError, TypeError):
        pass
    spc = _extract_char_spacing(run.font)
    if spc is not None:
        props["char_spacing"] = spc
    effect = _extract_run_effect(run)
    if effect:
        props["effect"] = effect
    return props


def _extract_run_effect(run) -> dict | None:
    """Extract outer shadow effect from a run's rPr effectLst."""
    try:
        rpr = run.font._element
        effect_lst = rpr.find(qn("a:effectLst"))
        if effect_lst is None or len(effect_lst) == 0:
            return None
        shadow = effect_lst.find(qn("a:outerShdw"))
        if shadow is None:
            return None
        return parse_shadow_xml(shadow)
    except (AttributeError, TypeError, IndexError):
        return None


def extract_bullet_properties(paragraph) -> dict:
    """Extract bullet properties from a paragraph's pPr element.

    Returns dict with bullet_char, bullet_font, bullet_size_pct, bullet_color
    when present. Returns {"bullet_none": True} when buNone is set.
    """
    props = {}
    pPr = paragraph._p.find(qn("a:pPr"))
    if pPr is None:
        return props

    buNone = pPr.find(qn("a:buNone"))
    if buNone is not None:
        props["bullet_none"] = True
        return props

    buChar = pPr.find(qn("a:buChar"))
    if buChar is not None:
        props["bullet_char"] = buChar.get("char", "•")

    buFont = pPr.find(qn("a:buFont"))
    if buFont is not None:
        typeface = buFont.get("typeface")
        if typeface:
            props["bullet_font"] = typeface

    buSzPct = pPr.find(qn("a:buSzPct"))
    if buSzPct is not None:
        val = buSzPct.get("val")
        if val:
            props["bullet_size_pct"] = int(val)

    buClr = pPr.find(qn("a:buClr"))
    if buClr is not None:
        srgb = buClr.find(qn("a:srgbClr"))
        if srgb is not None:
            props["bullet_color"] = f"#{srgb.get('val', '000000')}"

    # Extract paragraph margin and indent (controls bullet-to-text spacing)
    marL = pPr.get("marL")
    if marL is not None:
        props["bullet_margin_left"] = int(marL)

    indent = pPr.get("indent")
    if indent is not None:
        props["bullet_indent"] = int(indent)

    return props


def apply_bullet_properties(paragraph, elem: dict):
    """Apply bullet properties to a paragraph via lxml.

    Reads bullet_char, bullet_font, bullet_size_pct, bullet_color,
    bullet_margin_left, bullet_indent from elem.
    """
    has_bullet_props = (
        "bullet_char" in elem
        or "bullet_none" in elem
        or "bullet_margin_left" in elem
        or "bullet_indent" in elem
    )
    if not has_bullet_props:
        return

    pPr = paragraph._p.find(qn("a:pPr"))
    if pPr is None:
        pPr = etree.SubElement(paragraph._p, qn("a:pPr"))
        paragraph._p.insert(0, pPr)

    # Apply margin and indent (controls bullet-to-text spacing)
    if "bullet_margin_left" in elem:
        pPr.set("marL", str(elem["bullet_margin_left"]))
    if "bullet_indent" in elem:
        pPr.set("indent", str(elem["bullet_indent"]))

    if elem.get("bullet_none"):
        etree.SubElement(pPr, qn("a:buNone"))
        return

    if "bullet_font" in elem:
        buFont = etree.SubElement(pPr, qn("a:buFont"))
        buFont.set("typeface", elem["bullet_font"])

    if "bullet_size_pct" in elem:
        buSzPct = etree.SubElement(pPr, qn("a:buSzPct"))
        buSzPct.set("val", str(elem["bullet_size_pct"]))

    if "bullet_color" in elem:
        buClr = etree.SubElement(pPr, qn("a:buClr"))
        srgb = etree.SubElement(buClr, qn("a:srgbClr"))
        srgb.set("val", elem["bullet_color"].lstrip("#"))

    if "bullet_char" in elem:
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", elem["bullet_char"])
