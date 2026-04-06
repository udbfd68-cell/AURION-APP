# Copyright (c) Microsoft Corporation.
# SPDX-License-Identifier: MIT
"""Validate PPTX properties that cannot be detected from rendered images.

Checks speaker notes and slide count. Visual checks (overlay, overflow,
spacing, contrast, margins, placeholders) are performed by validate_slides.py
through Copilot SDK vision model inspection of rendered slide images.

Usage:
    python validate_deck.py --input slide-deck/presentation.pptx --content-dir content/
    python validate_deck.py --input deck.pptx --output results.json --report report.md
    python validate_deck.py --input deck.pptx --slides "1,3,5" --output results.json
"""

import argparse
import json
import logging
import sys
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx_utils import (
    EXIT_ERROR,
    EXIT_FAILURE,
    EXIT_SUCCESS,
    configure_logging,
    parse_slide_filter,
)

logger = logging.getLogger(__name__)

SEVERITY_ICON = {"error": "❌", "warning": "⚠️", "info": "ℹ️"}
QUALITY_ICON = {"good": "✅", "needs-attention": "⚠️"}


def check_speaker_notes(slide, slide_num: int) -> list[dict]:
    """Check for missing or empty speaker notes.

    Distinguishes between: notes slide present with content, notes slide
    present but empty, and notes slide absent entirely.

    Returns:
        List of issue dicts with check_type, severity, description, location.
    """
    issues = []
    try:
        if not slide.has_notes_slide:
            issues.append(
                {
                    "check_type": "speaker_notes",
                    "severity": "warning",
                    "description": "Missing speaker notes (no notes slide)",
                    "location": "notes",
                }
            )
            return issues
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if not notes:
            issues.append(
                {
                    "check_type": "speaker_notes",
                    "severity": "info",
                    "description": "Speaker notes present but empty",
                    "location": "notes",
                }
            )
    except (AttributeError, TypeError):
        issues.append(
            {
                "check_type": "speaker_notes",
                "severity": "warning",
                "description": "Missing speaker notes",
                "location": "notes",
            }
        )
    return issues


def validate_deck(
    pptx_path: Path,
    content_dir: Path | None = None,
    slide_filter: set[int] | None = None,
) -> dict:
    """Run PPTX-only validation checks (speaker notes, slide count).

    Returns:
        Dict with source, slide_count, slides (per-slide issues), and
        optional top-level issues for slide count findings.
    """
    prs = Presentation(str(pptx_path))
    total_slides = len(prs.slides)
    slides = []
    top_level_issues = []

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        if slide_filter and slide_num not in slide_filter:
            continue
        issues = check_speaker_notes(slide, slide_num)
        quality = "good" if not issues else "needs-attention"
        slides.append(
            {
                "slide_number": slide_num,
                "issues": issues,
                "overall_quality": quality,
            }
        )

    if content_dir and not slide_filter:
        slide_dirs = sorted(
            [
                d
                for d in content_dir.iterdir()
                if d.is_dir() and d.name.startswith("slide-")
            ]
        )
        if len(slide_dirs) != total_slides:
            if len(slide_dirs) < total_slides:
                top_level_issues.append(
                    {
                        "check_type": "slide_count",
                        "severity": "info",
                        "description": (
                            "Partial content detected"
                            f" — {total_slides} slides in PPTX, "
                            f"{len(slide_dirs)} content directories"
                            " (expected for incremental updates)"
                        ),
                        "location": "deck",
                    }
                )
            else:
                top_level_issues.append(
                    {
                        "check_type": "slide_count",
                        "severity": "warning",
                        "description": (
                            f"Slide count mismatch: {total_slides} slides in PPTX, "
                            f"{len(slide_dirs)} content directories"
                        ),
                        "location": "deck",
                    }
                )

    result = {
        "source": "pptx-properties",
        "slide_count": total_slides,
        "slides": slides,
    }
    if top_level_issues:
        result["deck_issues"] = top_level_issues
    return result


def generate_report(results: dict) -> str:
    """Generate a Markdown validation report from results.

    Args:
        results: Validation results dict from validate_deck().

    Returns:
        Markdown report string.
    """
    lines = ["# PPTX Property Validation Report", ""]
    ts = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
    lines.append(f"**Generated**: {ts}  ")
    lines.append(f"**Source**: {results['source']}  ")
    lines.append(f"**Slides**: {results['slide_count']}")
    lines.append("")

    # Count severities across all slides and deck-level issues
    error_count = 0
    warning_count = 0
    info_count = 0
    for slide in results["slides"]:
        for issue in slide.get("issues", []):
            sev = issue.get("severity", "info")
            if sev == "error":
                error_count += 1
            elif sev == "warning":
                warning_count += 1
            else:
                info_count += 1
    for issue in results.get("deck_issues", []):
        sev = issue.get("severity", "info")
        if sev == "error":
            error_count += 1
        elif sev == "warning":
            warning_count += 1
        else:
            info_count += 1

    lines.append("## Summary")
    lines.append("")
    lines.append("| Severity | Count |")
    lines.append("|-|-|")
    lines.append(f"| ❌ Errors | {error_count} |")
    lines.append(f"| ⚠️ Warnings | {warning_count} |")
    lines.append(f"| ℹ️ Info | {info_count} |")
    lines.append("")

    # Deck-level issues
    deck_issues = results.get("deck_issues", [])
    if deck_issues:
        lines.append("## Deck-Level Findings")
        lines.append("")
        lines.append("| Severity | Check | Description |")
        lines.append("|-|-|-|")
        for issue in deck_issues:
            sev = issue.get("severity", "info")
            sev_icon = SEVERITY_ICON.get(sev, "")
            check = issue.get("check_type", "")
            desc = issue.get("description", "")
            lines.append(f"| {sev_icon} {sev} | {check} | {desc} |")
        lines.append("")

    # Per-slide details
    lines.append("## Per-Slide Findings")
    lines.append("")
    for slide in results["slides"]:
        num = slide.get("slide_number", "?")
        quality = slide.get("overall_quality", "unknown")
        icon = QUALITY_ICON.get(quality, "❓")
        lines.append(f"### Slide {num} {icon} {quality}")
        lines.append("")

        issues = slide.get("issues", [])
        if not issues:
            lines.append("No issues found.")
            lines.append("")
            continue

        lines.append("| Severity | Check | Location | Description |")
        lines.append("|-|-|-|-|")
        for issue in issues:
            sev = issue.get("severity", "info")
            sev_icon = SEVERITY_ICON.get(sev, "")
            check = issue.get("check_type", "")
            loc = issue.get("location", "")
            desc = issue.get("description", "")
            lines.append(f"| {sev_icon} {sev} | {check} | {loc} | {desc} |")
        lines.append("")

    return "\n".join(lines)


def max_severity(results: dict) -> str:
    """Return the highest severity found across all issues."""
    severities = set()
    for slide in results["slides"]:
        for issue in slide.get("issues", []):
            severities.add(issue.get("severity", "info"))
    for issue in results.get("deck_issues", []):
        severities.add(issue.get("severity", "info"))
    if "error" in severities:
        return "error"
    if "warning" in severities:
        return "warning"
    if "info" in severities:
        return "info"
    return "none"


def create_parser() -> argparse.ArgumentParser:
    """Create and configure argument parser."""
    parser = argparse.ArgumentParser(
        description="Validate PPTX-only properties (speaker notes, slide count)"
    )
    parser.add_argument(
        "--input", required=True, type=Path, help="Input PPTX file path"
    )
    parser.add_argument(
        "--content-dir", type=Path, help="Content directory for slide count comparison"
    )
    parser.add_argument(
        "--slides", help="Comma-separated slide numbers to validate (default: all)"
    )
    parser.add_argument(
        "--output", type=Path, help="Output JSON file path (default: stdout)"
    )
    parser.add_argument("--report", type=Path, help="Output Markdown report file path")
    parser.add_argument(
        "--per-slide-dir",
        type=Path,
        help="Directory for per-slide JSON files (slide-NNN-deck-validation.json)",
    )
    parser.add_argument(
        "-v", "--verbose", action="store_true", help="Enable verbose logging"
    )
    return parser


def main() -> int:
    """Main entry point."""
    parser = create_parser()
    args = parser.parse_args()
    configure_logging(getattr(args, "verbose", False))

    pptx_path = args.input
    if not pptx_path.exists():
        logger.error("File not found: %s", pptx_path)
        return EXIT_ERROR

    slide_filter = parse_slide_filter(args.slides)

    logger.info("Validating PPTX properties: %s", pptx_path)
    results = validate_deck(pptx_path, args.content_dir, slide_filter=slide_filter)

    # Write per-slide deck validation JSON files
    if args.per_slide_dir:
        args.per_slide_dir.mkdir(parents=True, exist_ok=True)
        for slide_result in results["slides"]:
            slide_num = slide_result.get("slide_number", 0)
            per_slide_path = (
                args.per_slide_dir / f"slide-{slide_num:03d}-deck-validation.json"
            )
            per_slide_json = json.dumps(slide_result, indent=2)
            per_slide_path.write_text(per_slide_json, encoding="utf-8")
            logger.debug("Per-slide deck results written to %s", per_slide_path)

    # Output JSON
    output_json = json.dumps(results, indent=2)
    if args.output:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(output_json, encoding="utf-8")
        logger.info("Results written to %s", args.output)
    else:
        print(output_json)

    # Generate Markdown report
    if args.report:
        report_md = generate_report(results)
        args.report.parent.mkdir(parents=True, exist_ok=True)
        args.report.write_text(report_md, encoding="utf-8")
        logger.info("Report written to %s", args.report)

    # Report summary
    total_issues = sum(len(s.get("issues", [])) for s in results["slides"])
    total_issues += len(results.get("deck_issues", []))
    severity = max_severity(results)
    slide_count = results["slide_count"]
    logger.info(
        "Validation complete: %d issue(s) across %d slide(s)",
        total_issues,
        slide_count,
    )

    # Exit code: info-only → success, warning/error → failure
    if severity in ("error", "warning"):
        return EXIT_FAILURE
    return EXIT_SUCCESS


if __name__ == "__main__":
    sys.exit(main())
