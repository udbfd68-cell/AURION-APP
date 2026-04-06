# Copyright (c) Microsoft Corporation.
# SPDX-License-Identifier: MIT
"""Extract content from an existing PPTX into YAML content and style definitions.

Usage::

    python extract_content.py \
        --input existing-deck.pptx --output-dir content/

    python extract_content.py \
        --input existing-deck.pptx --output-dir content/ \
        --slides 3,7,15
"""

import argparse
import logging
from collections import Counter
from pathlib import Path

import cairosvg
import yaml
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx_charts import extract_chart
from pptx_colors import extract_color, hex_brightness
from pptx_fills import extract_effect_list, extract_fill, extract_line
from pptx_fonts import (
    extract_alignment,
    extract_font_info,
    extract_paragraph_font,
    normalize_font_family,
)
from pptx_shapes import AUTO_SHAPE_NAME_MAP, extract_rotation
from pptx_tables import extract_table
from pptx_text import (
    extract_bullet_properties,
    extract_paragraph_properties,
    extract_run_properties,
    extract_text_frame_properties,
)
from pptx_utils import emu_to_inches

MAX_IMAGE_BLOB_BYTES = 100 * 1024 * 1024  # 100 MB


class _ImageSecurityError(ValueError):
    """Security-critical image validation failure that must not be suppressed."""


_CONTENT_TYPE_TO_EXT: dict[str, str] = {
    "image/bmp": "bmp",
    "image/gif": "gif",
    "image/jpeg": "jpg",
    "image/png": "png",
    "image/tiff": "tiff",
    # WMF retained for legitimate PPTX files; validated by magic-byte check.
    # See CVE-2005-4560 for historical WMF risk context.
    "image/x-wmf": "wmf",
    # EMF retained for charts, SmartArt, and diagrams; validated by magic-byte check.
    "image/emf": "emf",
    "image/x-emf": "emf",
    # SVG sanitized via hardened XMLParser and converted to PNG by cairosvg.
    "image/svg+xml": "svg",
}

# WMF file signatures used for magic-byte validation.
_WMF_ALDUS_MAGIC = b"\xd7\xcd\xc6\x9a"
_WMF_STANDARD_PREFIXES = (b"\x01\x00\x09\x00", b"\x02\x00\x09\x00")

# EMF file signatures: EMR_HEADER record type at offset 0, " EMF" at offset 40.
_EMF_RECORD_TYPE = b"\x01\x00\x00\x00"
_EMF_SIGNATURE = b" EMF"


def _validate_wmf_magic_bytes(blob: bytes) -> None:
    """Reject WMF blobs that lack a recognized file signature."""
    if len(blob) < 4:
        raise _ImageSecurityError("WMF blob too short for magic-byte validation")
    head = blob[:4]
    if head == _WMF_ALDUS_MAGIC or head in _WMF_STANDARD_PREFIXES:
        return
    raise _ImageSecurityError(
        "WMF blob does not start with a recognized file signature"
    )


def _validate_emf_magic_bytes(blob: bytes) -> None:
    """Reject EMF blobs that lack the expected EMR_HEADER and signature."""
    if len(blob) < 44:
        raise _ImageSecurityError("EMF blob too short for magic-byte validation")
    if blob[:4] != _EMF_RECORD_TYPE or blob[40:44] != _EMF_SIGNATURE:
        raise _ImageSecurityError("EMF blob does not match expected file signature")


def _sanitize_svg(blob: bytes) -> bytes:
    """Parse SVG through a hardened XMLParser to block XXE and DTD attacks.

    Returns re-serialized XML bytes.  Raises *_ImageSecurityError* when
    the blob is not well-formed XML or contains prohibited constructs.
    """
    parser = etree.XMLParser(
        resolve_entities=False,
        no_network=True,
        dtd_validation=False,
        load_dtd=False,
    )
    try:
        root = etree.fromstring(blob, parser=parser)
    except etree.XMLSyntaxError as exc:
        raise _ImageSecurityError(f"SVG blob is not well-formed XML: {exc}") from exc
    if root.getroottree().docinfo.internalDTD is not None:
        raise _ImageSecurityError("SVG blob contains a DTD declaration")
    return etree.tostring(root, xml_declaration=True, encoding="UTF-8")


def _convert_svg_to_png(blob: bytes) -> bytes:
    """Sanitize an SVG blob and convert it to PNG via cairosvg."""
    clean_svg = _sanitize_svg(blob)
    return cairosvg.svg2png(bytestring=clean_svg)


def extract_connector(shape) -> dict:
    """Extract a connector element definition."""
    elem = {
        "type": "connector",
        "begin_x": emu_to_inches(shape.begin_x),
        "begin_y": emu_to_inches(shape.begin_y),
        "end_x": emu_to_inches(shape.end_x),
        "end_y": emu_to_inches(shape.end_y),
        "name": shape.name,
    }
    line_props = extract_line(shape)
    if line_props:
        elem.update(line_props)
    return elem


def _is_freeform(shape) -> bool:
    """Check whether a shape is a freeform with custom geometry."""
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    return shape._element.find(".//a:custGeom", nsmap) is not None


def _is_background_image(shape, slide_w: float, slide_h: float) -> bool:
    """Detect whether a PICTURE shape covers the full slide as a background.

    A shape qualifies if it covers at least 95% of slide dimensions.
    """
    w = emu_to_inches(shape.width)
    h = emu_to_inches(shape.height)
    return (w >= slide_w * 0.95) and (h >= slide_h * 0.95)


def _save_image_blob(shape, output_dir: Path, slide_num: int, img_count: int) -> dict:
    """Save an embedded image blob to disk with security validation.

    Validates content type against an allowlist, enforces a size limit,
    and checks that the resolved output path stays within *output_dir*.
    """
    try:
        img = shape.image
    except ValueError:
        return {"path": "LINKED_IMAGE_NOT_EMBEDDED"}

    ext = _CONTENT_TYPE_TO_EXT.get(img.content_type)
    if ext is None:
        raise ValueError(f"Unsupported image content type: {img.content_type}")

    blob = img.blob
    if len(blob) > MAX_IMAGE_BLOB_BYTES:
        raise ValueError(
            f"Image blob size {len(blob)} exceeds limit of {MAX_IMAGE_BLOB_BYTES} bytes"
        )

    if ext == "wmf":
        _validate_wmf_magic_bytes(blob)
    elif ext == "emf":
        _validate_emf_magic_bytes(blob)
    elif ext == "svg":
        blob = _convert_svg_to_png(blob)
        ext = "png"

    img_name = f"image-{img_count:02d}.{ext}"
    img_path = output_dir / "images" / img_name

    if not img_path.resolve().is_relative_to(output_dir.resolve()):
        raise _ImageSecurityError(
            f"Image path {img_path} escapes output directory {output_dir}"
        )

    img_path.parent.mkdir(parents=True, exist_ok=True)
    with open(img_path, "wb") as f:
        f.write(blob)
    return {"path": f"images/{img_name}"}


def extract_freeform(shape) -> dict:
    """Extract a freeform shape with its path vertices."""
    elem = {
        "type": "freeform",
        "left": emu_to_inches(shape.left),
        "top": emu_to_inches(shape.top),
        "width": emu_to_inches(shape.width),
        "height": emu_to_inches(shape.height),
        "name": shape.name,
    }

    rot = extract_rotation(shape)
    if rot is not None:
        elem["rotation"] = rot

    # Extract fill and line properties
    try:
        fill_result = extract_fill(shape.fill)
        if fill_result is not None:
            elem["fill"] = fill_result
    except (AttributeError, TypeError):
        pass

    line_props = extract_line(shape)
    if line_props:
        elem.update(line_props)

    # Extract path vertices from custGeom XML
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    paths = []
    for path_el in shape._element.findall(".//a:custGeom/a:pathLst/a:path", nsmap):
        commands = []
        for child in path_el:
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag == "moveTo":
                pt = child.find("a:pt", nsmap)
                if pt is not None:
                    commands.append(
                        {
                            "cmd": "moveTo",
                            "x": int(pt.get("x", 0)),
                            "y": int(pt.get("y", 0)),
                        }
                    )
            elif tag == "lnTo":
                pt = child.find("a:pt", nsmap)
                if pt is not None:
                    commands.append(
                        {
                            "cmd": "lineTo",
                            "x": int(pt.get("x", 0)),
                            "y": int(pt.get("y", 0)),
                        }
                    )
            elif tag == "cubicBezTo":
                pts = child.findall("a:pt", nsmap)
                commands.append(
                    {
                        "cmd": "cubicBezTo",
                        "pts": [
                            {"x": int(p.get("x", 0)), "y": int(p.get("y", 0))}
                            for p in pts
                        ],
                    }
                )
            elif tag == "close":
                commands.append({"cmd": "close"})
        if commands:
            paths.append(commands)

    if paths:
        elem["paths"] = paths

    return elem


MAX_GROUP_DEPTH = 20


def extract_group(
    shape,
    slide_num: int,
    output_dir,
    img_count: int,
    *,
    _depth: int = 0,
    max_depth: int = MAX_GROUP_DEPTH,
) -> dict:
    """Extract a group shape and its nested child elements.

    Raises ValueError when nesting exceeds *max_depth*.
    """
    if _depth >= max_depth:
        raise ValueError(f"Group nesting depth {_depth} exceeds limit of {max_depth}")
    elem = {
        "type": "group",
        "left": emu_to_inches(shape.left),
        "top": emu_to_inches(shape.top),
        "width": emu_to_inches(shape.width),
        "height": emu_to_inches(shape.height),
        "name": shape.name,
        "elements": [],
    }
    for child in shape.shapes:
        child_elem = extract_child_shape(
            child,
            slide_num,
            output_dir,
            img_count,
            _depth=_depth + 1,
            max_depth=max_depth,
        )
        if child_elem:
            elem["elements"].append(child_elem)
    return elem


def _extract_shape_by_type(
    shape,
    slide_num: int,
    output_dir,
    img_count: int,
    *,
    _depth: int = 0,
    max_depth: int = MAX_GROUP_DEPTH,
) -> dict | None:
    """Dispatch extraction based on shape_type, table/chart, or freeform."""
    shape_type = shape.shape_type

    # Simple shape_type dispatch (these extractors need no extra context)
    _SIMPLE_EXTRACTORS = {
        17: extract_textbox,  # TEXT_BOX
        1: extract_shape,  # AUTO_SHAPE
        9: extract_connector,  # LINE / CONNECTOR
    }
    extractor = _SIMPLE_EXTRACTORS.get(shape_type)
    if extractor:
        return extractor(shape)

    if shape_type == 13:  # PICTURE
        return extract_image(shape, output_dir, slide_num, img_count)
    if shape_type == 6:  # GROUP
        return extract_group(
            shape,
            slide_num,
            output_dir,
            img_count,
            _depth=_depth,
            max_depth=max_depth,
        )

    # Table and chart detection via attribute check
    if hasattr(shape, "has_table") and shape.has_table:
        return extract_table(shape)
    if hasattr(shape, "has_chart") and shape.has_chart:
        return extract_chart(shape)
    if _is_freeform(shape):
        return extract_freeform(shape)

    return None


def extract_child_shape(
    shape,
    slide_num: int,
    output_dir,
    img_count: int,
    *,
    _depth: int = 0,
    max_depth: int = MAX_GROUP_DEPTH,
) -> dict | None:
    """Extract a single child shape within a group."""
    result = _extract_shape_by_type(
        shape,
        slide_num,
        output_dir,
        img_count,
        _depth=_depth,
        max_depth=max_depth,
    )
    if result is not None:
        return result

    # Fallback for unrecognized shape types
    elem = {
        "type": "shape",
        "shape": "rectangle",
        "left": emu_to_inches(shape.left),
        "top": emu_to_inches(shape.top),
        "width": emu_to_inches(shape.width),
        "height": emu_to_inches(shape.height),
        "name": shape.name,
    }
    if shape.shape_type is not None:
        elem["_unrecognized_shape_type"] = int(shape.shape_type)
    return elem


def _has_formatting_variation(runs: list) -> bool:
    """Check if multiple runs have different formatting properties."""
    if len(runs) <= 1:
        return False
    fonts = {r.get("font") for r in runs if "font" in r}
    sizes = {r.get("size") for r in runs if "size" in r}
    colors = {r.get("color") for r in runs if "color" in r}
    bolds = {r.get("bold", False) for r in runs}
    italics = {r.get("italic", False) for r in runs}
    underlines = {r.get("underline", False) for r in runs}
    return (
        len(fonts) > 1
        or len(sizes) > 1
        or len(colors) > 1
        or len(bolds) > 1
        or len(italics) > 1
        or len(underlines) > 1
    )


# Key-mapping for extraction: maps canonical keys to output YAML key names
_SHAPE_EXTRACT_KEYS = {
    "font": "text_font",
    "size": "text_size",
    "color": "text_color",
    "bold": "text_bold",
}
_TEXTBOX_EXTRACT_KEYS = {
    "font": "font",
    "size": "font_size",
    "color": "font_color",
    "bold": "font_bold",
}

# Keys to promote from first paragraph to element level
_SHAPE_PROMOTE_KEYS = (
    "text_font",
    "text_size",
    "text_color",
    "text_bold",
    "italic",
    "alignment",
    "char_spacing",
)
_TEXTBOX_PROMOTE_KEYS = (
    "font",
    "font_size",
    "font_color",
    "font_bold",
    "italic",
    "alignment",
    "char_spacing",
)


def _extract_text_content(text_frame, keys: dict, promote_keys: tuple) -> dict:
    """Extract text content from a text frame into an element dict fragment.

    Handles paragraph iteration, run extraction, rich-text detection, and
    paragraph/element-level key promotion.

    Args:
        text_frame: python-pptx TextFrame object.
        keys: Key-mapping dict for font/size/color/bold output names.
        promote_keys: Tuple of keys to promote from first paragraph to element level.

    Returns:
        Dict with text, text frame properties, paragraph data, and promoted defaults.
    """
    result = {}
    text = text_frame.text.strip()
    if not text:
        return result

    result["text"] = text

    tf_props = extract_text_frame_properties(text_frame)
    if tf_props:
        result.update(tf_props)

    para_dicts = []
    for para in text_frame.paragraphs:
        run_info = {}
        para_runs = []
        for run in para.runs:
            font_info = extract_font_info(run.font)
            run_extra = extract_run_properties(run)
            para_runs.append({"text": run.text, **font_info, **run_extra})
            if not run_info:
                run_info = {**font_info, **run_extra}

        para_info = extract_paragraph_font(para)
        para_spacing = extract_paragraph_properties(para)
        bullet_props = extract_bullet_properties(para)
        alignment = extract_alignment(para)
        merged = {**para_info, **run_info}

        p_dict = {"text": para.text}
        if "font" in merged:
            p_dict[keys["font"]] = merged["font"]
        if "size" in merged:
            p_dict[keys["size"]] = merged["size"]
        if "color" in merged:
            p_dict[keys["color"]] = merged["color"]
        if merged.get("bold"):
            p_dict[keys["bold"]] = True
        if merged.get("italic"):
            p_dict["italic"] = True
        if merged.get("underline"):
            p_dict["underline"] = True
        if merged.get("hyperlink"):
            p_dict["hyperlink"] = merged["hyperlink"]
        if "char_spacing" in merged:
            p_dict["char_spacing"] = merged["char_spacing"]
        if "effect" in merged:
            p_dict["text_effect"] = merged["effect"]
        if alignment:
            p_dict["alignment"] = alignment
        if para_spacing:
            p_dict.update(para_spacing)
        if bullet_props:
            p_dict.update(bullet_props)
        if _has_formatting_variation(para_runs):
            p_dict["runs"] = para_runs
        para_dicts.append(p_dict)

    non_empty = [p for p in para_dicts if p["text"].strip()]
    any_has_runs = any("runs" in p for p in para_dicts)
    if len(para_dicts) > 1 or any_has_runs:
        result["paragraphs"] = para_dicts
        if non_empty:
            first = non_empty[0]
            for key in promote_keys:
                if key in first:
                    result[key] = first[key]
    elif non_empty:
        first = non_empty[0]
        for key, val in first.items():
            if key != "text":
                result[key] = val

    return result


def extract_shape(shape) -> dict:
    """Extract a shape element definition."""
    elem = {
        "type": "shape",
        "shape": "rectangle",
        "left": emu_to_inches(shape.left),
        "top": emu_to_inches(shape.top),
        "width": emu_to_inches(shape.width),
        "height": emu_to_inches(shape.height),
        "name": shape.name,
    }

    rot = extract_rotation(shape)
    if rot is not None:
        elem["rotation"] = rot

    # Detect shape type from auto_shape_type enum
    try:
        elem["shape"] = AUTO_SHAPE_NAME_MAP.get(shape.auto_shape_type, "rectangle")
    except (AttributeError, TypeError):
        elem["shape"] = "rectangle"

    # Extract corner radius (adjustment values) for rounded rectangles
    try:
        if shape.adjustments and len(shape.adjustments) > 0:
            elem["corner_radius"] = round(shape.adjustments[0], 5)
    except (AttributeError, TypeError, IndexError):
        pass

    # Extract fill
    try:
        fill_result = extract_fill(shape.fill)
        if fill_result is not None:
            elem["fill"] = fill_result
    except (AttributeError, TypeError):
        pass

    # Extract line properties
    line_props = extract_line(shape)
    if line_props:
        elem.update(line_props)

    # Extract effect list (outer shadow)
    effect = extract_effect_list(shape)
    if effect:
        elem["effect"] = effect

    # Extract text if present
    if shape.has_text_frame:
        text_data = _extract_text_content(
            shape.text_frame, _SHAPE_EXTRACT_KEYS, _SHAPE_PROMOTE_KEYS
        )
        elem.update(text_data)

    return elem


def extract_textbox(shape) -> dict:
    """Extract a text box element definition."""
    elem = {
        "type": "textbox",
        "left": emu_to_inches(shape.left),
        "top": emu_to_inches(shape.top),
        "width": emu_to_inches(shape.width),
        "height": emu_to_inches(shape.height),
        "text": shape.text_frame.text.strip() if shape.has_text_frame else "",
        "name": shape.name,
    }

    rot = extract_rotation(shape)
    if rot is not None:
        elem["rotation"] = rot

    if shape.has_text_frame:
        text_data = _extract_text_content(
            shape.text_frame, _TEXTBOX_EXTRACT_KEYS, _TEXTBOX_PROMOTE_KEYS
        )
        elem.update(text_data)

    return elem


def extract_image(shape, output_dir: Path, slide_num: int, img_count: int) -> dict:
    """Extract an image element and save the image file."""
    try:
        blob_result = _save_image_blob(shape, output_dir, slide_num, img_count)
    except _ImageSecurityError:
        raise
    except ValueError as exc:
        logging.warning("Skipping image on slide %d: %s", slide_num, exc)
        return {
            "type": "image",
            "path": "SKIPPED",
            "left": emu_to_inches(shape.left),
            "top": emu_to_inches(shape.top),
            "width": emu_to_inches(shape.width),
            "height": emu_to_inches(shape.height),
            "name": shape.name,
            "_skipped_reason": str(exc),
        }

    if blob_result["path"] == "LINKED_IMAGE_NOT_EMBEDDED":
        elem = {
            "type": "image",
            "path": "LINKED_IMAGE_NOT_EMBEDDED",
            "left": emu_to_inches(shape.left),
            "top": emu_to_inches(shape.top),
            "width": emu_to_inches(shape.width),
            "height": emu_to_inches(shape.height),
            "name": shape.name,
            "_note": "Image was linked, not embedded in the PPTX",
        }
        rot = extract_rotation(shape)
        if rot is not None:
            elem["rotation"] = rot
        return elem

    elem = {
        "type": "image",
        "path": blob_result["path"],
        "left": emu_to_inches(shape.left),
        "top": emu_to_inches(shape.top),
        "width": emu_to_inches(shape.width),
        "height": emu_to_inches(shape.height),
        "name": shape.name,
    }
    rot = extract_rotation(shape)
    if rot is not None:
        elem["rotation"] = rot

    # Extract image crop from srcRect on blipFill
    blipFill = shape._element.find(qn("p:blipFill"))
    if blipFill is not None:
        # Preserve blipFill attributes (rotWithShape, dpi, etc.)
        blip_fill_attrs = {}
        for attr_name in ("rotWithShape", "dpi"):
            val = blipFill.get(attr_name)
            if val is not None:
                blip_fill_attrs[attr_name] = val
        if blip_fill_attrs:
            elem["blip_fill_attrs"] = blip_fill_attrs

        srcRect = blipFill.find(qn("a:srcRect"))
        if srcRect is not None and srcRect.attrib:
            crop = {}
            for side in ("l", "t", "r", "b"):
                val = srcRect.get(side)
                if val is not None:
                    crop[side] = int(val)
            if crop:
                elem["crop"] = crop

    # Extract image opacity from alphaModFix on the blip element
    blip = shape._element.find(".//" + qn("a:blip"))
    if blip is not None:
        amf = blip.find(qn("a:alphaModFix"))
        if amf is not None:
            amt = int(amf.get("amt", "100000"))
            elem["opacity"] = round(amt / 1000, 1)

    return elem


def detect_global_style(prs) -> dict:
    """Analyze the presentation to detect common styling patterns.

    Detects multiple theme zones (e.g., light and dark slides) by clustering
    slides based on background brightness and dominant text colors.
    """
    bg_colors = Counter()
    text_colors = Counter()
    accent_colors = Counter()
    fill_colors = Counter()
    font_names = Counter()
    font_sizes = Counter()

    # Per-slide analysis for theme clustering
    slide_profiles = []

    slide_w = emu_to_inches(prs.slide_width)
    slide_h = emu_to_inches(prs.slide_height)

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        slide_bg = None
        slide_text_colors = Counter()
        slide_fill_colors = Counter()
        has_bg_image = False

        # Detect background colors
        try:
            fill_result = extract_fill(slide.background.fill)
            if isinstance(fill_result, str):
                bg_colors[fill_result] += 1
                slide_bg = fill_result
        except (AttributeError, TypeError):
            pass

        for i, shape in enumerate(slide.shapes):
            # Detect full-slide background images
            if (
                i == 0
                and shape.shape_type == 13
                and _is_background_image(shape, slide_w, slide_h)
            ):
                has_bg_image = True
                continue

            # Collect fill colors
            try:
                fill_result = extract_fill(shape.fill)
                if isinstance(fill_result, str):
                    h = emu_to_inches(shape.height)
                    if h < 0.1:
                        accent_colors[fill_result] += 1
                    else:
                        fill_colors[fill_result] += 1
                        slide_fill_colors[fill_result] += 1
            except (AttributeError, TypeError):
                pass

            # Collect font information
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            base_name = normalize_font_family(run.font.name)
                            font_names[base_name] += 1
                        if run.font.size:
                            font_sizes[int(run.font.size.pt)] += 1
                        try:
                            color = extract_color(run.font.color)
                            if isinstance(color, str) and color.startswith("#"):
                                text_colors[color] += 1
                                slide_text_colors[color] += 1
                        except (AttributeError, TypeError):
                            pass

        # Classify slide brightness
        bg_brightness = _classify_slide_brightness(
            slide_bg, slide_text_colors, has_bg_image
        )
        slide_profiles.append(
            {
                "slide": slide_num,
                "bg_color": slide_bg,
                "bg_brightness": bg_brightness,
                "has_bg_image": has_bg_image,
                "text_colors": dict(slide_text_colors),
                "fill_colors": dict(slide_fill_colors),
            }
        )

    # Build global color map from frequency analysis
    colors = _build_color_map(bg_colors, fill_colors, text_colors, accent_colors)

    # Detect themes by clustering slides into light/dark groups
    themes = _cluster_themes(slide_profiles, text_colors, fill_colors, accent_colors)

    # Determine primary fonts
    body_font = "Segoe UI"
    code_font = "Cascadia Code"
    for f, _count in font_names.most_common():
        if any(kw in f.lower() for kw in ("cascadia", "consolas", "mono", "courier")):
            code_font = f
        else:
            body_font = f
            break

    # Determine font sizes
    heading_size = 28
    body_size = 16
    if font_sizes:
        filtered = {s: c for s, c in font_sizes.items() if 8 < s < 60}
        if filtered:
            sorted_sizes = sorted(filtered.keys())
            body_size = sorted_sizes[len(sorted_sizes) // 2]
            heading_size = sorted_sizes[int(len(sorted_sizes) * 0.85)]

    style = {
        "dimensions": {
            "width_inches": emu_to_inches(prs.slide_width),
            "height_inches": emu_to_inches(prs.slide_height),
            "format": "16:9",
        },
        "defaults": {
            "speaker_notes_required": True,
        },
        "typography": {
            "body_font": body_font,
            "code_font": code_font,
            "heading_size": heading_size,
            "body_size": body_size,
        },
    }

    if colors:
        style["colors"] = colors

    if themes:
        style["themes"] = themes

    # Extract presentation metadata
    metadata = {}
    props = prs.core_properties
    for attr in ("title", "author", "subject", "keywords", "description", "category"):
        val = getattr(props, attr, None)
        if val:
            metadata[attr] = val
    if metadata:
        style["metadata"] = metadata

    return style


def _classify_slide_brightness(
    bg_color: str | None, text_colors: Counter, has_bg_image: bool
) -> str:
    """Classify a slide as 'light' or 'dark' based on background and text colors."""
    if has_bg_image and bg_color is None:
        # Slides with background images and no solid bg — infer from text colors
        dark_text = sum(
            c for hex_c, c in text_colors.items() if hex_brightness(hex_c) < 100
        )
        light_text = sum(
            c for hex_c, c in text_colors.items() if hex_brightness(hex_c) > 150
        )
        return "light" if dark_text >= light_text else "dark"

    if bg_color and isinstance(bg_color, str) and bg_color.startswith("#"):
        return "light" if hex_brightness(bg_color) > 128 else "dark"

    # Default: infer from text colors
    dark_text = sum(
        c for hex_c, c in text_colors.items() if hex_brightness(hex_c) < 100
    )
    light_text = sum(
        c for hex_c, c in text_colors.items() if hex_brightness(hex_c) > 150
    )
    if dark_text > light_text:
        return "light"
    if light_text > dark_text:
        return "dark"
    return "dark"


def _build_color_map(
    bg_colors: Counter,
    fill_colors: Counter,
    text_colors: Counter,
    accent_colors: Counter,
) -> dict:
    """Build the global color map from frequency analysis."""
    colors = {}
    if bg_colors:
        colors["bg_dark"] = bg_colors.most_common(1)[0][0]
    if fill_colors:
        colors["bg_card"] = fill_colors.most_common(1)[0][0]

    for color_hex, _count in text_colors.most_common(5):
        brightness = hex_brightness(color_hex)
        if brightness > 200 and "text_white" not in colors:
            colors["text_white"] = color_hex
        elif brightness < 80 and "text_dark" not in colors:
            colors["text_dark"] = color_hex
        elif 80 <= brightness <= 200 and "text_gray" not in colors:
            colors["text_gray"] = color_hex

    accent_names = ["accent_blue", "accent_teal", "accent_green"]
    for i, (color_hex, _count) in enumerate(accent_colors.most_common(3)):
        if i < len(accent_names):
            colors[accent_names[i]] = color_hex

    return colors


def _cluster_themes(
    slide_profiles: list[dict],
    text_colors: Counter,
    fill_colors: Counter,
    accent_colors: Counter,
) -> list[dict]:
    """Cluster slides into theme groups based on brightness classification."""
    light_slides = [p for p in slide_profiles if p["bg_brightness"] == "light"]
    dark_slides = [p for p in slide_profiles if p["bg_brightness"] == "dark"]

    # Only produce themes when both light and dark groups exist
    if not light_slides or not dark_slides:
        return []

    themes = []

    # Light theme
    light_text = Counter()
    light_fills = Counter()
    for p in light_slides:
        light_text.update(p["text_colors"])
        light_fills.update(p["fill_colors"])

    light_colors = {}
    for color_hex, _count in light_text.most_common(5):
        brightness = hex_brightness(color_hex)
        if brightness < 80 and "text_primary" not in light_colors:
            light_colors["text_primary"] = color_hex
        elif 80 <= brightness <= 200 and "text_secondary" not in light_colors:
            light_colors["text_secondary"] = color_hex
    if light_fills:
        light_colors["bg_card"] = light_fills.most_common(1)[0][0]

    themes.append(
        {
            "name": "light",
            "slides": sorted(p["slide"] for p in light_slides),
            "colors": light_colors,
        }
    )

    # Dark theme
    dark_text = Counter()
    dark_fills = Counter()
    dark_bgs = Counter()
    for p in dark_slides:
        dark_text.update(p["text_colors"])
        dark_fills.update(p["fill_colors"])
        if p["bg_color"]:
            dark_bgs[p["bg_color"]] += 1

    dark_colors = {}
    if dark_bgs:
        dark_colors["bg_dark"] = dark_bgs.most_common(1)[0][0]
    for color_hex, _count in dark_text.most_common(5):
        brightness = hex_brightness(color_hex)
        if brightness > 200 and "text_primary" not in dark_colors:
            dark_colors["text_primary"] = color_hex
        elif 80 <= brightness <= 200 and "text_secondary" not in dark_colors:
            dark_colors["text_secondary"] = color_hex
    if dark_fills:
        dark_colors["bg_card"] = dark_fills.most_common(1)[0][0]

    themes.append(
        {
            "name": "dark",
            "slides": sorted(p["slide"] for p in dark_slides),
            "colors": dark_colors,
        }
    )

    return themes


def extract_slide(
    slide,
    slide_num: int,
    output_dir: Path,
    slide_dims: tuple[float, float] | None = None,
) -> dict:
    """Extract all elements from a slide into a content.yaml structure."""
    slide_dir = output_dir / f"slide-{slide_num:03d}"
    slide_dir.mkdir(parents=True, exist_ok=True)

    content = {
        "slide": slide_num,
        "title": "",
        "elements": [],
    }

    # Extract layout name
    try:
        layout_name = slide.slide_layout.name
        if layout_name:
            content["layout"] = layout_name
    except (AttributeError, TypeError):
        pass

    # Extract slide background
    try:
        if not slide.follow_master_background:
            fill_result = extract_fill(slide.background.fill)
            if fill_result is not None:
                content["background"] = {"fill": fill_result}
    except (AttributeError, TypeError):
        pass

    # Extract speaker notes (include empty string when notes slide exists)
    try:
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            content["speaker_notes"] = notes
    except (AttributeError, TypeError):
        pass

    img_count = 0

    for z_index, shape in enumerate(list(slide.shapes)):
        shape_type = shape.shape_type

        # Track image count for filename generation
        if shape_type == 13:
            img_count += 1

        # Handle placeholders specially (extract as textbox with marker)
        if shape_type == 14:
            if not shape.has_text_frame:
                continue
            elem = extract_textbox(shape)
            elem["_placeholder"] = True
            elem["z_order"] = z_index
            content["elements"].append(elem)
            continue

        # Use shared dispatcher for all other shape types
        elem = _extract_shape_by_type(shape, slide_num, slide_dir, img_count)
        if elem is not None:
            elem["z_order"] = z_index
            content["elements"].append(elem)

            # Detect title from textbox near top of slide
            if (
                shape_type == 17
                and not content["title"]
                and emu_to_inches(shape.top) < 1.5
            ):
                text = shape.text_frame.text.strip() if shape.has_text_frame else ""
                if text and len(text) < 100:
                    content["title"] = text
            continue

        # Fallback for unrecognized shape types
        elem_data = {
            "type": "shape",
            "shape": "rectangle",
            "left": emu_to_inches(shape.left),
            "top": emu_to_inches(shape.top),
            "width": emu_to_inches(shape.width),
            "height": emu_to_inches(shape.height),
            "name": shape.name,
            "z_order": z_index,
        }
        if shape_type is not None:
            elem_data["_unrecognized_shape_type"] = int(shape_type)
        content["elements"].append(elem_data)

    return content, slide_dir


def _resolve_theme_colors(prs) -> dict:
    """Extract theme color name→hex mappings from the presentation's theme XML.

    Reads clrScheme from the slide master's theme and maps theme names
    (background_1, text_1, accent_1, etc.) to their actual hex values.
    """
    color_map = {}
    scheme_names = {
        "dk1": "dark_1",
        "dk2": "dark_2",
        "lt1": "light_1",
        "lt2": "light_2",
        "accent1": "accent_1",
        "accent2": "accent_2",
        "accent3": "accent_3",
        "accent4": "accent_4",
        "accent5": "accent_5",
        "accent6": "accent_6",
        "hlink": "hyperlink",
        "folHlink": "followed_hyperlink",
    }
    # Map canonical aliases
    aliases = {
        "dark_1": "text_1",
        "dark_2": "text_2",
        "light_1": "background_1",
        "light_2": "background_2",
    }
    try:
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        master = prs.slide_masters[0]
        theme_el = None
        # Theme is stored as a related part (generic Part, not XmlPart),
        # so parse its blob directly with lxml.
        for rel in master.part.rels.values():
            if "theme" in rel.reltype:
                parser = etree.XMLParser(resolve_entities=False, no_network=True)
                theme_el = etree.fromstring(rel.target_part.blob, parser=parser)
                break

        if theme_el is not None:
            clr_scheme = theme_el.find(f".//{{{ns_a}}}clrScheme")
            if clr_scheme is not None:
                for child in clr_scheme:
                    tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                    theme_name = scheme_names.get(tag)
                    if theme_name is None:
                        continue
                    # Extract hex value from srgbClr or sysClr
                    srgb = child.find(f"{{{ns_a}}}srgbClr")
                    if srgb is not None:
                        color_map[theme_name] = f"#{srgb.get('val', '000000')}"
                    else:
                        sys_clr = child.find(f"{{{ns_a}}}sysClr")
                        if sys_clr is not None:
                            color_map[theme_name] = (
                                f"#{sys_clr.get('lastClr', '000000')}"
                            )
                    # Add alias mappings
                    if theme_name in aliases:
                        alias = aliases[theme_name]
                        if theme_name in color_map:
                            color_map[alias] = color_map[theme_name]
    except (AttributeError, TypeError, IndexError):
        # Theme elements missing or malformed; degrade gracefully
        pass
    except etree.XMLSyntaxError:
        logging.warning(
            "Malformed theme XML in slide master; skipping theme color resolution"
        )
    return color_map


MAX_THEME_REF_DEPTH = 50


def _resolve_theme_refs_in_content(
    content: dict,
    theme_colors: dict,
    *,
    max_depth: int = MAX_THEME_REF_DEPTH,
) -> dict:
    """Replace @theme_name references with resolved hex values in content.

    Raises ValueError when nesting exceeds *max_depth*.
    """

    def resolve_value(val, _depth: int = 0):
        if _depth >= max_depth:
            raise ValueError(
                f"Theme reference nesting depth {_depth} exceeds limit of {max_depth}"
            )
        if isinstance(val, str) and val.startswith("@"):
            theme_name = val[1:]
            return theme_colors.get(theme_name, val)
        if isinstance(val, dict):
            return {k: resolve_value(v, _depth + 1) for k, v in val.items()}
        if isinstance(val, list):
            return [resolve_value(item, _depth + 1) for item in val]
        return val

    return resolve_value(content)


def main():
    """CLI entry point for extracting PPTX content into YAML."""
    parser = argparse.ArgumentParser(
        description="Extract content from a PPTX into YAML"
    )
    parser.add_argument("--input", required=True, help="Input PPTX file path")
    parser.add_argument("--output-dir", required=True, help="Output content directory")
    parser.add_argument(
        "--slides", help="Comma-separated slide numbers to extract (default: all)"
    )
    parser.add_argument(
        "--resolve-themes",
        action="store_true",
        help="Resolve @theme references to actual hex RGB values from the deck's theme",
    )
    args = parser.parse_args()

    pptx_path = Path(args.input)
    output_dir = Path(args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    slide_filter = None
    if args.slides:
        slide_filter = {int(s.strip()) for s in args.slides.split(",")}

    prs = Presentation(str(pptx_path))
    print(f"Extracting from: {pptx_path}")
    print(f"Slides: {len(prs.slides)}")
    w = emu_to_inches(prs.slide_width)
    h = emu_to_inches(prs.slide_height)
    print(f'Dimensions: {w}" x {h}"')

    # Detect and save global style
    global_style = detect_global_style(prs)

    # Resolve theme colors when requested
    theme_colors = {}
    if args.resolve_themes:
        theme_colors = _resolve_theme_colors(prs)
        if theme_colors:
            global_style["theme_colors"] = theme_colors
            global_style = _resolve_theme_refs_in_content(global_style, theme_colors)
            print(f"Resolved {len(theme_colors)} theme colors")

    global_dir = output_dir / "global"
    global_dir.mkdir(parents=True, exist_ok=True)
    style_path = global_dir / "style.yaml"
    with open(style_path, "w", encoding="utf-8") as f:
        yaml.dump(
            global_style,
            f,
            default_flow_style=False,
            sort_keys=False,
            allow_unicode=True,
        )
    print(f"Global style saved to {style_path}")

    # Extract slides (filtered or all)
    slide_dims = (emu_to_inches(prs.slide_width), emu_to_inches(prs.slide_height))
    extracted = 0
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        if slide_filter and slide_num not in slide_filter:
            continue
        content, slide_dir = extract_slide(
            slide, slide_num, output_dir, slide_dims=slide_dims
        )

        # Resolve @theme references to hex values when --resolve-themes is set
        if args.resolve_themes and theme_colors:
            content = _resolve_theme_refs_in_content(content, theme_colors)

        content_path = slide_dir / "content.yaml"
        with open(content_path, "w", encoding="utf-8") as f:
            yaml.dump(
                content,
                f,
                default_flow_style=False,
                sort_keys=False,
                allow_unicode=True,
            )
        print(
            f"Slide {slide_num}: {content.get('title', 'Untitled')} -> {content_path}"
        )
        extracted += 1

    print(f"\nExtraction complete. {extracted} slide(s) extracted to {output_dir}")


if __name__ == "__main__":
    main()
